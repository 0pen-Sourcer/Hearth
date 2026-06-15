"""QA check for a freshly-built PPTX — coordinate-based, no LibreOffice needed.

A deck can't be rendered to an image without LibreOffice (often absent), but
python-pptx exposes every shape's exact position + size, so we can catch the
defects that actually wreck a deck purely from geometry + text:

  - OFF-SLIDE: a shape spills past the slide edge (cut-off content).
  - COLLISION: two text-bearing shapes overlap heavily (text on top of text).
    (A label sitting on a fill/card shape is fine — the card carries no text,
    so it's excluded; only text-on-text is flagged.)
  - PLACEHOLDER text left in: lorem/ipsum/xxxx/"your text here"/TODO.
  - TINY font: runs under 8pt (unreadable on a slide).
  - EMPTY slide: nothing meaningful on it.

Usage:
    python qa_check_pptx.py <path-to.pptx>

Prints per-slide findings + a final QA RESULT line. Always exits 0 — advisory,
meant for the model that just built the deck to decide whether to fix + rebuild.
"""
from __future__ import annotations
import re
import sys

_PLACEHOLDER = re.compile(r"\b(lorem|ipsum|xxxx+|placeholder|your text here|todo|tbd|click to add)\b", re.I)


def _rect(sh):
    """(left, top, width, height) in EMU, or None if unpositioned."""
    try:
        if None in (sh.left, sh.top, sh.width, sh.height):
            return None
        return (int(sh.left), int(sh.top), int(sh.width), int(sh.height))
    except Exception:
        return None


def _overlap_area(a, b):
    ow = max(0, min(a[0] + a[2], b[0] + b[2]) - max(a[0], b[0]))
    oh = max(0, min(a[1] + a[3], b[1] + b[3]) - max(a[1], b[1]))
    return ow * oh


def _text_of(sh):
    try:
        if sh.has_text_frame:
            return sh.text_frame.text or ""
    except Exception:
        pass
    return ""


def _min_font_pt(sh):
    smallest = None
    try:
        if not sh.has_text_frame:
            return None
        for para in sh.text_frame.paragraphs:
            for run in para.runs:
                sz = run.font.size
                if sz is not None:
                    pt = sz.pt
                    if (run.text or "").strip() and (smallest is None or pt < smallest):
                        smallest = pt
    except Exception:
        pass
    return smallest


def main() -> int:
    if len(sys.argv) < 2:
        print("usage: python qa_check_pptx.py <pptx>")
        return 0
    try:
        from pptx import Presentation
    except ImportError:
        print("QA: python-pptx not installed — skipping check.")
        return 0
    try:
        prs = Presentation(sys.argv[1])
    except Exception as e:
        print(f"QA: couldn't open PPTX ({type(e).__name__}: {e}) — does the path exist?")
        return 0

    SW, SH = int(prs.slide_width), int(prs.slide_height)
    tol = int(0.015 * SW)  # ~1.5% edge tolerance
    flags = []
    print(f"QA: {len(prs.slides)} slide(s), canvas {SW/914400:.2f}\" x {SH/914400:.2f}\"")

    for si, slide in enumerate(prs.slides, 1):
        issues = []
        text_shapes = []   # (rect, text) for collision check
        meaningful = 0
        for sh in slide.shapes:
            r = _rect(sh)
            txt = _text_of(sh).strip()
            if r:
                meaningful += 1
                # off-slide
                if (r[0] < -tol or r[1] < -tol
                        or r[0] + r[2] > SW + tol or r[1] + r[3] > SH + tol):
                    issues.append(f"off-slide shape ({txt[:24] or 'graphic'!r})")
                if txt:
                    text_shapes.append((r, txt))
            if txt and _PLACEHOLDER.search(txt):
                issues.append(f"placeholder text: {txt[:30]!r}")
            mf = _min_font_pt(sh)
            if mf is not None and mf < 8:
                issues.append(f"tiny font {mf:.0f}pt: {txt[:24]!r}")
        # text-on-text collisions
        for i in range(len(text_shapes)):
            for j in range(i + 1, len(text_shapes)):
                a, b = text_shapes[i][0], text_shapes[j][0]
                ov = _overlap_area(a, b)
                smaller = min(a[2] * a[3], b[2] * b[3])
                if smaller > 0 and ov > 0.40 * smaller:
                    issues.append(f"text collision: {text_shapes[i][1][:18]!r} ~ {text_shapes[j][1][:18]!r}")
        if meaningful == 0:
            issues.append("empty slide")

        if issues:
            flags.extend(issues)
            print(f"  slide {si}: " + "; ".join(issues))
        else:
            print(f"  slide {si}: OK")

    if flags:
        print("QA RESULT: FIX NEEDED — adjust the flagged coordinates/text, rebuild ONCE, re-run this check.")
    else:
        print("QA RESULT: OK — no off-slide shapes, text collisions, placeholders, or tiny fonts.")
    return 0


if __name__ == "__main__":
    sys.exit(main())
