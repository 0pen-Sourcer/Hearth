"""Outline -> .pptx builder used by the `make-pptx` skill.

Slides separated by `---`. First line of each slide is the title.
Bullet lines start with `-` / `*` / `+`. Lines starting with `image:`
embed a local PNG/JPG. Plain lines become subtitle text under the title.

Run:
    python build_pptx.py --outline path/to/outline.txt --out path/to/out.pptx
                         [--title "Deck Title"]
"""
from __future__ import annotations

import argparse
import sys
from pathlib import Path


def _import_pptx():
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
        return locals()
    except ImportError as e:
        print(f"python-pptx missing: {e}. pip install python-pptx",
              file=sys.stderr)
        sys.exit(2)


def _split_slides(text: str):
    blocks = []
    cur: list = []
    for raw in text.splitlines():
        line = raw.rstrip()
        if line.strip() == "---":
            if cur:
                blocks.append(cur)
                cur = []
            continue
        cur.append(line)
    if cur:
        blocks.append(cur)
    return [b for b in blocks if any(ln.strip() for ln in b)]


def _hex_to_rgb(P, hex_str: str):
    h = hex_str.lstrip("#")
    if len(h) == 3:
        h = "".join(c * 2 for c in h)
    try:
        return P["RGBColor"](int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    except Exception:
        return None


def main() -> int:
    import json
    ap = argparse.ArgumentParser()
    ap.add_argument("--outline", required=True)
    ap.add_argument("--out", required=True)
    ap.add_argument("--title", default="")
    ap.add_argument("--style", default="",
                    help='Inline JSON. PREFER --style-file (PowerShell '
                         'mangles inline JSON).')
    ap.add_argument("--style-file", default="",
                    help='Path to .json with style overrides. Write it via '
                         'write_file then pass the path here. Keys: accent, '
                         'title_color, body_color, font, title_size, body_size.')
    args = ap.parse_args()
    style = {}
    sf = getattr(args, "style_file", "") or ""
    if sf:
        try:
            style = json.loads(Path(sf).read_text(encoding="utf-8"))
        except Exception as e:
            print(f"warn: --style-file ignored, bad JSON: {e}", file=sys.stderr)
    if args.style:
        try:
            style.update(json.loads(args.style))
        except Exception as e:
            print(f"warn: --style ignored, bad JSON: {e}", file=sys.stderr)

    src = Path(args.outline)
    if not src.is_file():
        print(f"error: outline not found: {src}", file=sys.stderr)
        return 1
    out = Path(args.out)
    out.parent.mkdir(parents=True, exist_ok=True)

    P = _import_pptx()
    prs = P["Presentation"]()
    prs.slide_width = P["Inches"](13.333)   # 16:9
    prs.slide_height = P["Inches"](7.5)

    slides = _split_slides(src.read_text(encoding="utf-8"))
    if not slides:
        print("error: no slides found in outline", file=sys.stderr)
        return 1

    # Style resolution
    font_name = style.get("font", "Calibri")
    title_size = int(style.get("title_size", 40))
    body_size = int(style.get("body_size", 18))
    title_color = _hex_to_rgb(P, style.get("title_color") or style.get("accent") or "")
    body_color = _hex_to_rgb(P, style.get("body_color", ""))

    def _style_run(run, size, color):
        run.font.name = font_name
        run.font.size = P["Pt"](size)
        if color is not None:
            run.font.color.rgb = color

    title_layout = prs.slide_layouts[0]
    content_layout = prs.slide_layouts[1]

    for i, block in enumerate(slides):
        title_text = next((ln for ln in block if ln.strip()), "")
        rest = [ln for ln in block if ln.strip() and ln != title_text]
        if i == 0:
            slide = prs.slides.add_slide(title_layout)
            slide.shapes.title.text = args.title or title_text
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    _style_run(run, title_size + 4, title_color)
            if slide.placeholders and len(slide.placeholders) > 1:
                sub = slide.placeholders[1]
                sub.text = "\n".join(
                    ln.lstrip("-*+ ").strip() for ln in rest if not ln.strip().startswith("image:"))
                for para in sub.text_frame.paragraphs:
                    for run in para.runs:
                        _style_run(run, body_size + 2, body_color)
        else:
            slide = prs.slides.add_slide(content_layout)
            slide.shapes.title.text = title_text
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs:
                    _style_run(run, title_size - 8, title_color)
            body_ph = None
            for ph in slide.placeholders:
                if ph.placeholder_format.idx == 1:
                    body_ph = ph
                    break
            if body_ph is not None:
                tf = body_ph.text_frame
                tf.clear()
                first = True
                for ln in rest:
                    s = ln.strip()
                    if s.startswith("image:"):
                        path = s[6:].strip()
                        try:
                            slide.shapes.add_picture(
                                path, P["Inches"](1), P["Inches"](2),
                                width=P["Inches"](6))
                        except Exception:
                            pass
                        continue
                    bullet_text = s.lstrip("-*+ ").strip()
                    if not bullet_text:
                        continue
                    p = tf.paragraphs[0] if first else tf.add_paragraph()
                    p.text = bullet_text
                    for run in p.runs:
                        _style_run(run, body_size, body_color)
                    first = False

    prs.save(str(out))
    print(str(out.resolve()))
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
