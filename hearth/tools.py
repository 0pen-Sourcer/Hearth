"""J.A.R.V.I.S. tool layer.

Two trust zones:
  - Reads can go anywhere on disk (Jarvis must KNOW your machine).
  - Writes / deletes / moves are sandboxed to WORKSPACE (~/Jarvis by default).

Override sandbox: set env JARVIS_WORKSPACE to any folder.
Tighten reads too: set env JARVIS_LOCKDOWN=1 to confine reads to the workspace.
"""

from __future__ import annotations

import os
import re
import sys
import json
import glob as globmod
import shlex
import shutil
import socket
import fnmatch
import platform
import subprocess
import urllib.request
import urllib.parse
from datetime import datetime
from html.parser import HTMLParser
from typing import Any, Callable, Dict, List, Optional, Tuple

# ============================================================
# CONFIG
# ============================================================

# Where Hearth keeps everything (memory, conversations, skills, settings).
# Resolution order:
#   1. JARVIS_WORKSPACE env  (rename feature + power users)
#   2. pointer file ~/.hearth/workspace.txt  (so a user whose C: is full can
#      put the workspace on D:/E: — written by onboarding / Settings)
#   3. default ~/Jarvis
# The pointer lives OUTSIDE the workspace (chicken-and-egg) so we can find the
# workspace before we've loaded anything from it.
_WORKSPACE_POINTER = os.path.join(os.path.expanduser("~"), ".hearth", "workspace.txt")


def _resolve_workspace() -> str:
    env = os.environ.get("JARVIS_WORKSPACE")
    if env and env.strip():
        return os.path.abspath(env.strip())
    try:
        if os.path.isfile(_WORKSPACE_POINTER):
            p = open(_WORKSPACE_POINTER, encoding="utf-8").read().strip()
            if p:
                return os.path.abspath(p)
    except OSError:
        pass
    return os.path.join(os.path.expanduser("~"), "Jarvis")


def set_workspace_location(path: str) -> str:
    """Point Hearth at a new workspace folder (e.g. on a drive with space).
    Writes the pointer file + creates the dir. Returns the resolved path. The
    caller should move existing files there and restart for it to take effect
    everywhere (WORKSPACE is read at import across modules)."""
    path = os.path.abspath(os.path.expanduser(path.strip()))
    os.makedirs(path, exist_ok=True)
    os.makedirs(os.path.dirname(_WORKSPACE_POINTER), exist_ok=True)
    with open(_WORKSPACE_POINTER, "w", encoding="utf-8") as f:
        f.write(path)
    return path


WORKSPACE = os.path.abspath(_resolve_workspace())
# Propagate the resolved path into the env so every other module that does
# `os.environ.get("JARVIS_WORKSPACE") or ~/Jarvis` (llmserver, mcp_client,
# memory, jobs, web) aligns to the SAME workspace — including one chosen via
# the pointer file, not just the env var. tools.py is imported first, so this
# runs before those modules resolve their own WORKSPACE.
os.environ["JARVIS_WORKSPACE"] = WORKSPACE
SAFE_READ_ONLY = os.environ.get("JARVIS_LOCKDOWN", "").strip() in ("1", "true", "yes")

# Extra writeable roots — paths the user has explicitly opted in to,
# either via the JARVIS_EXTRA_WORKSPACES env var (semicolon- or
# comma-separated) or via the /allow runtime command in the CLI.
# Writes/edits/deletes/moves succeed when the resolved path is under any
# of these in addition to the main WORKSPACE.
EXTRA_WORKSPACES: List[str] = []


def _parse_extra_env() -> List[str]:
    raw = os.environ.get("JARVIS_EXTRA_WORKSPACES", "").strip()
    if not raw:
        return []
    parts = [p.strip().strip('"').strip("'") for p in raw.replace(";", ",").split(",")]
    return [os.path.abspath(os.path.expanduser(p)) for p in parts if p]


EXTRA_WORKSPACES.extend(_parse_extra_env())


def add_extra_workspace(path: str) -> str:
    """Allow writes under an additional path. Returns confirmation."""
    p = os.path.abspath(os.path.expanduser(path or "").strip())
    if not p:
        return "Error: empty path"
    if not os.path.isdir(p):
        return f"Error: not a directory: {p}"
    if p not in EXTRA_WORKSPACES:
        EXTRA_WORKSPACES.append(p)
    return f"writes now allowed under: {p}"


def remove_extra_workspace(path: str) -> str:
    p = os.path.abspath(os.path.expanduser(path or "").strip())
    if p in EXTRA_WORKSPACES:
        EXTRA_WORKSPACES.remove(p)
        return f"writes no longer allowed under: {p}"
    return f"not in the allow-list: {p}"


def list_extra_workspaces() -> List[str]:
    return list(EXTRA_WORKSPACES)

SHOTS_DIR = os.path.join(WORKSPACE, "screenshots")
LOGS_DIR = os.path.join(WORKSPACE, "logs")
MEMORY_DIR = os.path.join(WORKSPACE, "memory")

# Legacy: keep NOTES_DIR around for users who already have ~/Jarvis/notes/
NOTES_DIR = os.path.join(WORKSPACE, "notes")

for d in (WORKSPACE, SHOTS_DIR, LOGS_DIR, MEMORY_DIR):
    os.makedirs(d, exist_ok=True)

# Per-tool result caps. Reads get more, write confirmations get less.
RESULT_CAPS: Dict[str, int] = {
    "read_file": 16000,
    "list_archive": 8000,
    "extract_archive_file": 1500,
    "summarize_file": 6000,
    "list_directory": 6000,
    "grep_search": 8000,
    "glob_files": 4000,
    "web_search": 5000,
    "web_fetch": 8000,
    "run_command": 6000,
    "list_processes": 5000,
    "list_installed_apps": 6000,
    "system_info": 3000,
    "memory_recall": 8000,
    "memory_list": 6000,
    "disk_usage": 6000,
    "locate_path": 4000,
}
DEFAULT_CAP = 4000

EXCLUDE_DIRS = {".git", ".godot", "__pycache__", ".vs", "node_modules",
                ".venv", "venv", "dist", "build", ".idea", ".import",
                # Windows system-managed roots that show up at drive roots.
                # No user content here; walking them wastes find_file budget.
                "$RECYCLE.BIN", "System Volume Information", "$Recycle.Bin",
                "Config.Msi", "Recovery", "$WinREAgent"}

HOME = os.path.expanduser("~")

# Suppress the brief cmd console flash on every subprocess we spawn under
# the GUI/tray context. Without it, every nvidia-smi / tasklist / ripgrep
# / lms call makes a black box flash on screen — looks like a virus.
_NO_WINDOW = 0x08000000 if sys.platform == "win32" else 0

# Pluggable callback the `ask_user` tool routes to. Each surface that wants
# to support interactive questions registers its own — CLI prints a numbered
# prompt via prompt_toolkit; the web surface emits an ndjson event and parks
# on a queue. When no callback is registered (headless / batch), ask_user
# returns an error and the model is expected to pick a safe default instead
# of looping forever.
_ask_user_callback: Optional[Callable[[str, list, bool], Dict[str, Any]]] = None


def set_ask_user_callback(cb: Optional[Callable[[str, list, bool], Dict[str, Any]]]) -> None:
    """Register the surface-specific implementation of ask_user.

    The callback receives (question, options, allow_other) and must return
    a dict shaped like {"ok": True, "choice": str, "other": bool} or
    {"ok": False, "error": str}. Synchronous from the tool's perspective —
    it blocks until the user answers (or the surface times out)."""
    global _ask_user_callback
    _ask_user_callback = cb

# Common places users actually keep stuff — find_file walks these in order
# before giving up. Order matters: workspace + Desktop tie-break first.
COMMON_USER_DIRS = ["Desktop", "Documents", "Downloads", "Pictures", "Videos", "Music"]
COMMON_DEV_DIRS = ["Code", "Projects", "source", "repos", "dev", "src"]

# kind hint → file extension whitelist (used by find_file)
FIND_KIND_EXTENSIONS: Dict[str, set] = {
    "image":       {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".ico", ".svg"},
    "video":       {".mp4", ".mkv", ".avi", ".mov", ".webm", ".flv", ".wmv"},
    "audio":       {".mp3", ".wav", ".flac", ".ogg", ".m4a", ".aac", ".opus"},
    "doc":         {".pdf", ".docx", ".doc", ".txt", ".md", ".rtf", ".odt"},
    "code":        {".py", ".js", ".ts", ".jsx", ".tsx", ".go", ".rs", ".java",
                    ".kt", ".swift", ".c", ".cpp", ".h", ".hpp", ".rb", ".php",
                    ".cs", ".sh", ".ps1", ".bat", ".html", ".css", ".scss"},
    "archive":     {".zip", ".rar", ".7z", ".tar", ".gz", ".bz2", ".xz"},
    "spreadsheet": {".xlsx", ".xls", ".csv", ".ods"},
}

# Search-side budgets so the model can't accidentally walk a million files.
MAX_FILES_TO_SCAN = 50000
_DRIVE_ROOT_RE = re.compile(r"^[A-Za-z]:[\\/]?$|^/$")

# Directory-name keywords by kind — used by find_file to prioritize walking
# into folders likely to contain the requested media. So "find a video" walks
# into `movies/`, `Videos/`, `shows/` BEFORE a `photos/` folder (which would
# otherwise win alphabetically and eat the scan budget on image files that can
# never match kind=video anyway).
_KIND_DIR_KEYWORDS: Dict[str, Tuple[str, ...]] = {
    "video":       ("movie", "video", "film", "show", "tv", "media", "season", "episode"),
    "audio":       ("music", "audio", "song", "podcast", "sound", "album"),
    "image":       ("pic", "photo", "image", "art", "screenshot", "wallpaper", "icon"),
    "code":        ("code", "src", "source", "project", "repo", "dev", "git", "scripts"),
    "doc":         ("doc", "paper", "pdf", "note", "report", "book", "manual", "guide"),
    "archive":     ("archive", "backup", "zip"),
    "spreadsheet": ("data", "sheet", "spreadsheet", "report", "stat"),
}


def _enumerate_non_system_drives() -> List[str]:
    """Return all fixed-drive roots that aren't the system drive. On Windows
    we walk D:..Z: and check existence. On POSIX, there are no extra mounts
    by default — caller falls back to HOME scanning only.

    Used so find_file actually sees things like `D:\\Movies\\<film>.mkv` and
    `E:\\Games` — most users keep media + games on non-system drives, and a
    HOME-only walk misses every single one."""
    if os.name != 'nt':
        return []
    roots: List[str] = []
    # Skip C: (system drive — already covered via HOME)
    for letter in "DEFGHIJKLMNOPQRSTUVWXYZ":
        root = f"{letter}:\\"
        try:
            if os.path.isdir(root):
                roots.append(root)
        except OSError:
            continue
    return roots


def _trunc(s: str, cap: int) -> str:
    if len(s) <= cap:
        return s
    return s[:cap] + f"\n…[truncated {len(s) - cap} chars]"


# ============================================================
# PATH SAFETY
# ============================================================

def _resolve_read(p: str) -> str:
    """Return absolute path for a read-style op.

    Search order for relative paths:
      1. Current working directory (where Hearth was launched from). This is
         what most users expect when they type './file.pdf' or 'foo.txt'.
      2. Workspace (~/Jarvis/). Fallback for when the file lives there.

    The previous behavior was workspace-only, which made `./Timelines.pdf`
    silently resolve to `~/Jarvis/Timelines.pdf` even when the file was
    sitting in the user's CWD — confused agents AND users alike.
    """
    p = os.path.expanduser(p)
    if not os.path.isabs(p):
        cwd_candidate = os.path.abspath(os.path.join(os.getcwd(), p))
        ws_candidate  = os.path.abspath(os.path.join(WORKSPACE, p))
        # Prefer whichever actually exists; if neither does, prefer CWD so
        # the error message at the call site shows the path the user meant.
        if os.path.exists(cwd_candidate):
            p = cwd_candidate
        elif os.path.exists(ws_candidate):
            p = ws_candidate
        else:
            p = cwd_candidate
    else:
        p = os.path.abspath(p)
    if SAFE_READ_ONLY and not p.startswith(WORKSPACE):
        raise PermissionError(f"Read locked to workspace: {WORKSPACE}")
    return p


# Callback the host (web.py / hearth_cli.py) registers so _resolve_write can
# ask the user "extend the writable area to include this path?" instead of
# silently raising. Signature: (path: str) -> bool. Returning True means
# "allow this write AND add the path's parent to EXTRA_WORKSPACES so the
# next write under the same root doesn't re-prompt".
_extend_workspace_callback: Optional[Callable[[str], bool]] = None


def set_extend_workspace_callback(cb: Optional[Callable[[str], bool]]) -> None:
    """Wire up the host's permission prompt for out-of-workspace writes.
    The CLI / GUI registers this at startup. None disables the hook and
    _resolve_write falls back to the old raise-immediately behavior."""
    global _extend_workspace_callback
    _extend_workspace_callback = cb


def _resolve_write(p: str) -> str:
    """Return absolute path for a write-style op. Stays inside WORKSPACE or
    EXTRA_WORKSPACES. If outside both, asks the host via the registered
    extend-workspace callback; on approval, the parent directory joins
    EXTRA_WORKSPACES so subsequent writes under the same root don't prompt
    again. Raises PermissionError when no callback is set or user denies."""
    p = os.path.expanduser(p)
    if not os.path.isabs(p):
        p = os.path.join(WORKSPACE, p)
    p = os.path.abspath(p)

    def _inside(root: str) -> bool:
        return p == root or p.startswith(root + os.sep)

    if _inside(WORKSPACE):
        return p
    for extra in EXTRA_WORKSPACES:
        if _inside(extra):
            return p

    # Outside any approved root. Ask the host if a prompt is available.
    if _extend_workspace_callback is not None:
        try:
            granted = bool(_extend_workspace_callback(p))
        except Exception:
            granted = False
        if granted:
            parent = p if os.path.isdir(p) else os.path.dirname(p) or p
            if parent and parent not in EXTRA_WORKSPACES:
                EXTRA_WORKSPACES.append(parent)
            return p

    extras = "\n  ".join(EXTRA_WORKSPACES) if EXTRA_WORKSPACES else "(none)"
    raise PermissionError(
        f"Write blocked — '{p}' escapes workspace ({WORKSPACE}).\n"
        f"Extra allowed paths:\n  {extras}\n"
        f"To allow this path, run /allow <path> in the CLI or set "
        f"JARVIS_EXTRA_WORKSPACES."
    )


# ============================================================
# TOOL DEFINITIONS (provider-agnostic JSON schema)
# ============================================================

TOOL_DEFINITIONS: List[Dict[str, Any]] = [
    # ---- FILES ----
    {
        "name": "read_file",
        "description": (
            "Smart file reader. Auto-detects format and extracts text from "
            "PDF (pypdf), DOCX (python-docx), XLSX/XLSM (openpyxl), PPTX "
            "(python-pptx), EPUB, IPYNB, CSV/TSV, JSON/JSONL, HTML/XML, RTF, "
            "and single-stream .gz/.bz2/.xz. For plain text/code/logs, "
            "returns line-numbered output. For archives (.zip/.tar/...) "
            "returns a hint to use list_archive. For images, hints to use "
            "view_image. Path can be absolute (anywhere readable) or "
            "relative to workspace. start_line/end_line slice text files; "
            "for PDFs they become start_page/end_page; for XLSX/CSV they "
            "cap rows-per-sheet / row count."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "File path. Absolute or relative to workspace."},
                "start_line": {"type": "integer", "description": "1-based start. Text: line. PDF: page. CSV: row. Optional."},
                "end_line": {"type": "integer", "description": "1-based inclusive end (same meaning per type). Optional."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "write_file",
        "description": (
            "Create a NEW file inside the workspace. For MODIFYING an existing "
            "file, ALWAYS use edit_file instead — this tool will REFUSE to "
            "overwrite an existing file with >30 lines (forces use of "
            "edit_file). Set overwrite=true only when the user explicitly "
            "asked for a full rewrite."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "content": {"type": "string"},
                "overwrite": {"type": "boolean", "description": "Set true to allow clobbering an existing >30-line file. Default false (use edit_file instead)."},
            },
            "required": ["path", "content"],
        },
    },
    {
        "name": "edit_file",
        "description": (
            "Targeted string-replace edits — never rewrites the whole file. "
            "Each edit replaces `old_text` with `new_text`. old_text must be "
            "UNIQUE in the file (include surrounding context to make it so), "
            "OR set `replace_all: true` on the edit for variable/symbol renames. "
            "Falls back to whitespace-tolerant matching if exact fails. "
            "ALWAYS call read_file first to see the exact text including "
            "indentation. Multiple edits are applied in order."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string"},
                "edits": {
                    "type": "array",
                    "items": {
                        "type": "object",
                        "properties": {
                            "old_text": {"type": "string", "description": "Exact text to find — must be unique unless replace_all is true."},
                            "new_text": {"type": "string", "description": "What to replace it with."},
                            "replace_all": {"type": "boolean", "description": "If true, replace every occurrence. Defaults to false."},
                        },
                        "required": ["old_text", "new_text"],
                    },
                },
                "replace_all": {"type": "boolean", "description": "Default replace_all for every edit in this call. Per-edit setting wins."},
            },
            "required": ["path", "edits"],
        },
    },
    {
        "name": "list_directory",
        "description": "List a directory. Set recursive=true and max_depth to walk subtrees. Read-only access anywhere on disk.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Directory path. Default = workspace."},
                "recursive": {"type": "boolean"},
                "max_depth": {"type": "integer", "description": "Default 2 when recursive."},
            },
        },
    },
    {
        "name": "create_directory",
        "description": "Create a directory inside the workspace.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
    },
    {
        "name": "delete_path",
        "description": "Delete a file or directory inside the workspace. Refuses paths outside workspace.",
        "parameters": {
            "type": "object",
            "properties": {"path": {"type": "string"}},
            "required": ["path"],
        },
    },
    {
        "name": "move_path",
        "description": "Move or rename inside the workspace.",
        "parameters": {
            "type": "object",
            "properties": {
                "source": {"type": "string"},
                "destination": {"type": "string"},
            },
            "required": ["source", "destination"],
        },
    },
    {
        "name": "list_archive",
        "description": (
            "List contents of a .zip/.jar/.whl/.apk/.tar/.tar.gz/.tar.bz2/.tar.xz "
            "archive WITHOUT extracting it. Returns path + size per entry. "
            "Use this before extract_archive_file or before deciding whether to "
            "ask the user to unpack. For .rar/.7z, hints to use 7-Zip via "
            "run_command (stdlib can't read them). Read-only — works anywhere on disk."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Path to the archive."},
                "limit": {"type": "integer", "description": "Max entries to return. Default 200."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "search_chats",
        "description": (
            "Search across ALL past chat conversations (full-text via SQLite "
            "FTS5). Use this whenever the user says 'what did we talk about', "
            "'remember when we discussed X', 'find that thing from last week', "
            "or asks something that requires recall beyond the current chat. "
            "Returns top matches with snippets — read them, then answer the "
            "user's question without quoting the convo_ids."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Keywords or phrase."},
                "limit": {"type": "integer", "description": "Max matches. Default 8."},
            },
            "required": ["query"],
        },
    },
    {
        "name": "set_reminder",
        "description": (
            "Schedule a future desktop notification, optionally with a TOOL CALL "
            "that runs at the same moment (an 'action reminder'). Accepts natural "
            "time strings: 'in 25 minutes', 'tomorrow at 7am', '2026-05-27 09:00', "
            "'9pm', 'next monday at 10am', 'in 2 hours'. Saved to "
            "~/Jarvis/reminders.json; a background watcher fires when due. Catches "
            "up on missed reminders the next time Hearth launches (the user sees "
            "'while you were away'). Use whenever the user says 'remind me to X at "
            "Y'. For repeating reminders ('every 30 minutes', 'daily'), pass the "
            "cadence as `recurring`. For 'at 5pm summarize my emails', pass "
            "`action_tool='summarize_emails'` + `action_args={}` so the tool "
            "actually runs alongside the notification."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "when": {"type": "string", "description": "When to fire the FIRST time. Natural-language OK."},
                "what": {"type": "string", "description": "The message to show the user."},
                "recurring": {"type": "string", "description": "Optional: 'every 30 minutes', 'hourly', 'daily', 'weekly', 'every 2 hours', etc. Omit for one-shot reminders."},
                "action_tool": {"type": "string", "description": "Optional: name of a Hearth tool to RUN when the reminder fires. The tool's result is appended to the toast body. E.g. 'summarize_file', 'web_search', 'open_app'."},
                "action_args": {"type": "object", "description": "Optional: arguments for action_tool. Object matching the tool's parameter schema. Ignored if action_tool is empty."},
                "tag": {"type": "string", "description": "Optional: free-form label like 'work' / 'side' / 'medication' for grouping in the GUI."},
            },
            "required": ["when", "what"],
        },
    },
    {
        "name": "snooze_reminder",
        "description": (
            "Push a reminder's due time forward by N minutes (default 10). Works "
            "on already-fired one-shots too - they resurrect as un-fired. Use when "
            "the user says 'snooze that 5 min' / 'remind me again in an hour' / "
            "'not yet, push it'. Cheaper than re-creating the reminder from scratch."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "id": {"type": "string", "description": "Reminder id from list_reminders or the set_reminder response."},
                "minutes": {"type": "integer", "description": "Minutes to push forward. Default 10."},
            },
            "required": ["id"],
        },
    },
    {
        "name": "list_reminders",
        "description": "List all upcoming (un-fired) reminders. Set include_fired=true to see history.",
        "parameters": {
            "type": "object",
            "properties": {
                "include_fired": {"type": "boolean"},
            },
        },
    },
    {
        "name": "cancel_reminder",
        "description": "Cancel a scheduled reminder by id. Use list_reminders to get the id.",
        "parameters": {
            "type": "object",
            "properties": {"id": {"type": "string"}},
            "required": ["id"],
        },
    },
    {
        "name": "spawn_subagent",
        "description": (
            "Fork a focused, scoped sub-agent. Personas live under "
            "hearth/subagents/*.md - call list_subagent_personas first if "
            "unsure which slug to use.\n"
            "PARALLEL = THE DEFAULT FOR MULTIPLE AGENTS. If the user wants "
            "more than one agent (\"spawn 3 researchers\", \"in parallel\", "
            "\"send a team\") emit ALL the spawn_subagent calls in ONE turn "
            "with mode='background'. They run CONCURRENTLY and each result "
            "auto-arrives as a <task-notification>. NEVER spawn sync, wait, "
            "then spawn the next — that runs them one-at-a-time (the slow "
            "ladder the user does NOT want).\n"
            "  mode='background': returns immediately with agent_id + "
            "    transcript_path; result notification arrives on a later "
            "    turn (no polling). Use for ANY multi-agent or long task.\n"
            "  mode='sync' (default): blocks until the child returns. Use "
            "    ONLY for a single, short, must-have-it-now subtask.\n"
            "Depth-limited to 3 nested forks."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "persona": {"type": "string", "description":
                    "Slug of a persona under hearth/subagents/ — call "
                    "list_subagent_personas() first if unsure. Current set: "
                    "researcher, coder, archivist, librarian, summarizer, "
                    "pdf_coordinator."},
                "prompt": {"type": "string", "description":
                    "Focused work for the child. Tight scope - one PDF "
                    "chunk, one question, one file."},
                "max_turns": {"type": "integer", "description":
                    "Override the persona's default turn cap. Capped at 20 "
                    "even if higher passed."},
                "mode": {"type": "string", "enum": ["sync", "background"],
                    "description": "sync = block; background = fire-and-"
                    "forget with a task-notification on completion."},
                "name": {"type": "string", "description":
                    "Optional human label for this instance. Use when "
                    "spawning multiple subagents of the SAME persona to "
                    "tell them apart (e.g. researcher 'Alex' on topic A "
                    "and researcher 'Beth' on topic B). Appears in the "
                    "completion notification + transcript."},
            },
            "required": ["persona", "prompt"],
        },
    },
    {
        "name": "list_subagent_personas",
        "description":
            "List the personas available for spawn_subagent. Returns "
            "[{slug, name, description, allowed_tools, cost_class}]. Call "
            "this when forking but unsure which persona handles the shape.",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "get_subagent_result",
        "description":
            "Poll a background subagent for its result. Normally you don't "
            "need this - the completion notification auto-arrives as the "
            "next user message. Use it only when you want to peek before "
            "the notification surfaces (e.g. after a long pause). Returns "
            "{ok, status: 'running'|'done', result?: {...}}.",
        "parameters": {
            "type": "object",
            "properties": {"agent_id": {"type": "string"}},
            "required": ["agent_id"],
        },
    },
    {
        "name": "extract_archive_file",
        "description": (
            "Pull ONE file out of an archive into the workspace, without "
            "unpacking the whole thing. archive_path can be anywhere on "
            "disk; inner_path is the path inside the archive (use "
            "list_archive first to see options). The extracted file lands "
            "in the workspace under output_name (default = basename of "
            "inner_path). Refuses '..' in inner_path. Supports zip/jar/whl/"
            "apk + tar family."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "archive_path": {"type": "string", "description": "Path to the archive."},
                "inner_path": {"type": "string", "description": "Path of the file INSIDE the archive (e.g. 'docs/readme.txt'). If unique, a suffix match like 'readme.txt' also works."},
                "output_name": {"type": "string", "description": "Optional workspace-relative destination name. Default = basename(inner_path)."},
            },
            "required": ["archive_path", "inner_path"],
        },
    },

    # ---- SEARCH ----
    {
        "name": "grep_search",
        "description": "Regex search across files. Uses ripgrep if available, else Python. Returns matching lines with file:line.",
        "parameters": {
            "type": "object",
            "properties": {
                "pattern": {"type": "string", "description": "Regex (Python flavor)."},
                "path": {"type": "string", "description": "Directory to search. Default = workspace."},
                "glob": {"type": "string", "description": "File glob filter, e.g. '*.py'."},
                "case_insensitive": {"type": "boolean"},
                "max_matches": {"type": "integer", "description": "Cap matches. Default 100."},
            },
            "required": ["pattern"],
        },
    },
    {
        "name": "glob_files",
        "description": "Find files by glob pattern. Returns paths sorted by mtime (newest first). For multiple patterns, separate with '|', ';' or ',' — e.g. '*.png|*.jpg'. A JSON array also works. Drive-root paths (C:\\, D:\\) are refused — use `find_file` instead.",
        "parameters": {
            "type": "object",
            "properties": {
                "pattern": {"type": "string", "description": "e.g. '**/*.py' or '~/Documents/**/*.pdf' or '*.png|*.jpg'."},
                "path": {"type": "string", "description": "Base dir. Default = workspace. Cannot be a drive root."},
            },
            "required": ["pattern"],
        },
    },
    {
        "name": "find_file",
        "description": (
            "Find files (or folders) by name across common locations — workspace, Desktop, "
            "Documents, Downloads, Pictures, Videos, Music, ~/Code, ~/Projects, "
            "the current working dir, AND every non-system drive (D:, E:, F:, G:...). "
            "Use this whenever the user says 'find X' / 'where's Y' / 'do I have any Z' "
            "instead of asking them for a path. Pass a name substring or a glob "
            "(e.g. 'budget', '*.pdf', 'vacation_*'). For media/binary results, just "
            "report the paths (don't read_file them). For text files, read top results. "
            "If user says 'search C drive' / 'check G:\\\\SteamLibrary', pass that as "
            "`path` to scope the scan AND get a 4x bigger budget."),
        "parameters": {
            "type": "object",
            "properties": {
                "name": {"type": "string", "description": "Filename substring or glob pattern."},
                "kind": {"type": "string", "description": "Optional category narrowing by extension: image, video, audio, doc, code, archive, spreadsheet, or 'any' (default)."},
                "limit": {"type": "integer", "description": "Max results. Default 10."},
                "deep": {"type": "boolean", "description": "Recurse deeper (max depth 4 vs 2 without path; 8 vs 5 with explicit path). Default false."},
                "path": {"type": "string", "description": "Optional explicit search root, e.g. 'G:\\\\SteamLibrary' or 'C:\\\\Program Files'. Overrides the common-locations enumeration and grants a 4x larger scan budget. Use when the user has named a drive or folder."},
            },
            "required": ["name"],
        },
    },

    # ---- WEB ----
    {
        "name": "web_search",
        "description": "Free DuckDuckGo HTML search. Returns top result titles/snippets/URLs. INVISIBLE to the user — this is research for YOU, not something they can see. To OPEN a page for the user, use open_url / open_in_browser / browse.",
        "parameters": {
            "type": "object",
            "properties": {
                "query": {"type": "string"},
                "limit": {"type": "integer", "description": "Default 6."},
            },
            "required": ["query"],
        },
    },
    {
        "name": "web_fetch",
        "description": (
            "Fetch a URL and return its readable text (HTML stripped) to YOU. "
            "INVISIBLE to the user — they do NOT see a browser, just whatever you "
            "tell them. Use this to READ a page yourself. If the user wants to "
            "OPEN / WATCH / PLAY something (a video, a site), do NOT web_fetch it — "
            "use open_url / open_in_browser (their own browser) or browse instead. "
            "Never say 'I opened it' after a web_fetch — you didn't."
        ),
        "parameters": {
            "type": "object",
            "properties": {"url": {"type": "string"}},
            "required": ["url"],
        },
    },

    # ---- SHELL ----
    {
        "name": "run_command",
        "description": (
            "Execute a shell command. On Windows defaults to PowerShell — both "
            "classic commands (dir, where, tasklist, ipconfig) AND PowerShell "
            "cmdlets (Get-ChildItem, Get-Process, Sort-Object) work. Default "
            "120s timeout, max 300s. Use timeout=180+ for pip installs. "
            "Set detached=true for things that DON'T exit on their own — "
            "daemons, dev servers, UI launchers, game launchers, Forge's "
            "run_neo.bat. Without detached they will hang the call until "
            "timeout fires."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "command": {"type": "string"},
                "cwd": {"type": "string", "description": "Working dir. Default = workspace."},
                "timeout": {"type": "integer", "description": "Seconds. Default 120, max 300. Ignored when detached=true."},
                "shell": {"type": "string", "description": "Optional. 'cmd' to force cmd.exe on Windows. Default = powershell."},
                "detached": {"type": "boolean", "description": "If true, spawn in a new console and return immediately with the PID. Use for daemons / UI launchers / dev servers that don't self-terminate."},
            },
            "required": ["command"],
        },
    },

    # ---- BACKGROUND JOBS (long-running shell commands that shouldn't block the agent) ----
    {
        "name": "start_job",
        "description": (
            "Run a shell command in the BACKGROUND and return a job_id "
            "immediately so you can keep working while it runs. Use this for "
            "anything that takes >30s (pip install of heavy deps, HF model "
            "download, dataset preprocessing, build). The job's output streams "
            "to disk; check it with job_status or wait for it with job_wait. "
            "DO NOT use for fast commands (use run_command) or daemons / UIs "
            "(use run_command with detached=true)."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "command":     {"type": "string"},
                "cwd":         {"type": "string", "description": "Working dir. Default = workspace."},
                "shell":       {"type": "string", "description": "'cmd' to force cmd.exe on Windows. Default = powershell."},
                "description": {"type": "string", "description": "Short label for the job ('install torch', 'download Hermes-3 GGUF'). Shown in job_list."},
            },
            "required": ["command"],
        },
    },
    {
        "name": "job_status",
        "description": "Return current status + last ~40 lines of output for a background job. Cheap to poll. status is one of: starting / running / completed / failed / killed.",
        "parameters": {
            "type": "object",
            "properties": {"job_id": {"type": "string"}},
            "required": ["job_id"],
        },
    },
    {
        "name": "job_wait",
        "description": "Block up to `timeout_s` seconds for the job to finish, then return final status + output. If timeout fires while still running, returns status='running' — call again to keep waiting. Default 30s, max 300s.",
        "parameters": {
            "type": "object",
            "properties": {
                "job_id":    {"type": "string"},
                "timeout_s": {"type": "number", "description": "Max seconds to block. Default 30, max 300."},
            },
            "required": ["job_id"],
        },
    },
    {
        "name": "job_kill",
        "description": "Terminate a running background job by job_id. No-op if already finished.",
        "parameters": {
            "type": "object",
            "properties": {"job_id": {"type": "string"}},
            "required": ["job_id"],
        },
    },
    {
        "name": "job_list",
        "description": "List recent background jobs (newest first). Pass active_only=true to see only running / starting jobs.",
        "parameters": {
            "type": "object",
            "properties": {
                "active_only": {"type": "boolean", "description": "Default false — include completed/failed/killed."},
            },
        },
    },

    # ---- MEDIA GENERATION (image + video via xAI Grok Imagine / OpenAI) ----
    {
        "name": "generate_image",
        "description": (
            "Generate an image from a text prompt. Saves to ~/Jarvis/generated/ "
            "and returns the file path. The GUI renders it inline; CLI opens "
            "it in the default image viewer. Supported providers: xAI Grok "
            "(grok-imagine-image-quality), OpenAI (gpt-image-2 by default — "
            "their new flagship), Google Gemini (gemini-2.5-flash-image, aka "
            "Nano Banana). Auto-routes based on whatever /brain you're on. "
            "Will fail loudly on local LM Studio / Ollama (they don't do "
            "image gen)."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "prompt":       {"type": "string", "description": "What to draw. Be specific — 'a glowing violet H logo on a charcoal background, minimalist, vector style'."},
                "n":            {"type": "integer", "description": "How many to generate (1-4). Default 1."},
                "aspect_ratio": {"type": "string", "description": "1:1 | 16:9 | 9:16 | 4:3 | 3:4 | 3:2 | 2:3 etc. Default 1:1."},
                "resolution":   {"type": "string", "description": "xAI: '1k' or '2k'. Default '1k'."},
            },
            "required": ["prompt"],
        },
    },
    {
        "name": "generate_video",
        "description": (
            "Start an ASYNC video generation. Returns a task_id immediately; "
            "videos take 20-60+ seconds. Poll with check_video_task(task_id) "
            "or just tell the user 'video's cooking, I'll let you know'. The "
            "user can ask 'is my video ready?' later and you can check then. "
            "xAI Grok only today."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "prompt":       {"type": "string", "description": "What to animate. Cinematic descriptions work best."},
                "duration":     {"type": "integer", "description": "Seconds, 1-15. Default 5."},
                "aspect_ratio": {"type": "string", "description": "16:9 | 9:16 | 1:1 | 4:3 | 3:4 | 3:2 | 2:3. Default 16:9."},
                "resolution":   {"type": "string", "description": "'480p' or '720p'. Default 720p."},
                "image_url":    {"type": "string", "description": "Optional — image URL to animate (image-to-video mode)."},
            },
            "required": ["prompt"],
        },
    },
    {
        "name": "check_video_task",
        "description": (
            "Poll a video generation task ONCE. Does NOT block. Status is one "
            "of: pending / done / failed / expired / unknown. When done, the "
            "response includes 'path' = the saved mp4 file the GUI/CLI can "
            "render. Safe to call repeatedly — caches the downloaded file."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "task_id": {"type": "string", "description": "The task_id returned by generate_video."},
            },
            "required": ["task_id"],
        },
    },
    {
        "name": "list_generations",
        "description": "List the 10 most recent image/video generation tasks (live + finished). Useful when the user says 'show me what you generated' or 'is anything still cooking?'",
        "parameters": {"type": "object", "properties": {}},
    },

    # ---- INTERACTIVE: ASK THE USER ----
    {
        "name": "ask_user",
        "description": (
            "Ask the user a multi-choice question when you genuinely need a "
            "decision before continuing — picking between two valid approaches, "
            "clarifying an ambiguous file/folder, choosing which of several "
            "matches to act on. The user sees a numbered list (CLI) or a modal "
            "with buttons (GUI). Return value is the chosen label, or the "
            "user's free-text reply if they pick 'Other'. "
            "DO NOT use this for every small decision — it interrupts. Use it "
            "only when running ahead would commit to the wrong thing. If the "
            "user is clearly absent (headless / batch run), this returns an "
            "error and you should pick the safest default and proceed."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "question": {"type": "string", "description": "One-sentence question. End with '?'."},
                "options":  {
                    "type": "array", "items": {"type": "string"},
                    "description": "2-6 short option labels. Recommend the safest first.",
                },
                "allow_other": {"type": "boolean", "description": "Whether to offer an Other free-text option. Default true."},
            },
            "required": ["question", "options"],
        },
    },

    # ---- KNOW MY PC ----
    {
        "name": "system_info",
        "description": "OS, CPU, RAM, disk, hostname, user, uptime — a snapshot of the machine.",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "list_processes",
        "description": "List running processes (top by memory). Optional name_filter.",
        "parameters": {
            "type": "object",
            "properties": {
                "name_filter": {"type": "string"},
                "limit": {"type": "integer", "description": "Default 20."},
            },
        },
    },
    {
        "name": "network_info",
        "description": "Local IP, hostname, network adapters. No external probes.",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "get_battery",
        "description": "Battery percentage and AC status (if supported).",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "disk_usage",
        "description": (
            "Find the biggest folders and files under a path. Native Python "
            "— no shell. Drive-root scans ('C:\\\\') and whole-tree walks "
            "(max_depth<=0) AUTO-BACKGROUND because they can take minutes "
            "to hours. When backgrounded, returns a job_id IMMEDIATELY so "
            "the user keeps chatting; call get_job_result(job_id) later "
            "for the report. Override with background:false to force sync."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Drive or directory to scan, e.g. 'G:\\\\' or 'C:/Users/me/Downloads'."},
                "top_n": {"type": "integer", "description": "How many top entries to return for each section. Default 15."},
                "kind": {"type": "string", "enum": ["both", "folders", "files"], "description": "What to list. Default 'both'."},
                "max_depth": {"type": "integer", "description": "How deep to walk for folder totals. Default 1 (direct subfolders only). 0 = full recursion (auto-backgrounds)."},
                "background": {"type": "boolean", "description": "Optional override. true = always background; false = always sync. Default = auto (sync for shallow scans, background for drive roots / full recursion)."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "list_jobs",
        "description": (
            "List background jobs (started by tools like disk_usage on a drive "
            "root, or explicit start_job calls). Returns each job's id, status "
            "(running / completed / failed), description, started_at, and "
            "elapsed time. Use this when the user asks 'is that scan done yet?' "
            "or 'what's running?'."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "active_only": {"type": "boolean", "description": "If true, hide jobs that finished long ago. Default false."},
            },
        },
    },
    {
        "name": "get_job_result",
        "description": (
            "Get the result of a background job by id. Returns the job's "
            "completed output if finished, or {status: 'running', elapsed_s} "
            "if still in flight. Use the job_id returned by disk_usage (or "
            "any other tool that auto-backgrounded)."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "job_id": {"type": "string"},
            },
            "required": ["job_id"],
        },
    },
    {
        "name": "locate_path",
        "description": "Smart locator: find a folder or app by name without globbing the whole disk. Checks top-level dirs of every drive, common parent dirs (Documents, Downloads, Desktop, Program Files, LocalAppData), Start Menu shortcuts, and the Windows installed-apps registry. Returns ranked matches. USE THIS instead of glob_files('**/*X*') for 'where is X on my PC' questions — it's fast and never crashes.",
        "parameters": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "Substring of the folder / app / game name (case-insensitive). E.g. 'chrome', 'steam', 'vscode'."},
                "limit": {"type": "integer", "description": "Max results. Default 15."},
            },
            "required": ["query"],
        },
    },
    {
        "name": "list_installed_apps",
        "description": "Installed applications (Windows registry uninstall keys). Optional name_filter (substring, case-insensitive).",
        "parameters": {
            "type": "object",
            "properties": {
                "name_filter": {"type": "string"},
                "limit": {"type": "integer", "description": "Default 50."},
            },
        },
    },

    # ---- APP CONTROL ----
    {
        "name": "open_app",
        "description": "Launch / open with default association. Accepts: app names ('notepad', 'spotify', 'chrome'), full file paths (e.g. 'G:\\\\movie.mkv' opens in default player, '.rar' in archive tool), folder paths (opens in Explorer), or URLs. Resolution order: existing path → URL → PATH lookup → UWP URI → Start Menu shortcut. Use this whenever the user says 'open this', 'launch X', or references something they want to interact with.",
        "parameters": {
            "type": "object",
            "properties": {"name": {"type": "string"}},
            "required": ["name"],
        },
    },
    {
        "name": "open_url",
        "description": "Open a URL in the user's DEFAULT browser (their login, fullscreen). The user sees it, but YOU lose control of it afterward — you can't change the page, search again, or click. Use for a one-off handoff, AND when the user wants SEVERAL links opened as TABS to browse themselves — call open_url once per link (their browser handles multiple tabs; the browse tool can't — it's a single controlled tab). If you'll need to keep driving one page (change the video, navigate), use browse instead. Never web_fetch a thing you're meant to OPEN — fetch is invisible.",
        "parameters": {
            "type": "object",
            "properties": {"url": {"type": "string"}},
            "required": ["url"],
        },
    },
    {
        "name": "open_in_browser",
        "description": "Open a URL in a SPECIFIC browser (and optional profile) — the user SEES it in their own browser, and it stays open. Like open_url but for a named browser. Use for 'open/watch/play X in Brave/Chrome'. Browsers detected from PATH + registry; if the user has a preferred browser in memory, use it, else ask once.",
        "parameters": {
            "type": "object",
            "properties": {
                "url": {"type": "string"},
                "browser": {
                    "type": "string",
                    "description": "Browser name: 'chrome', 'brave', 'edge', 'firefox', 'opera', 'vivaldi'. Omit to use the system default.",
                },
                "profile": {
                    "type": "string",
                    "description": "Optional profile directory (Chrome/Brave/Edge: e.g. 'Default', 'Profile 1') or profile name (Firefox: e.g. 'work').",
                },
            },
            "required": ["url"],
        },
    },
    {
        "name": "list_browsers",
        "description": "List installed browsers detected on this machine. Use when the user asks 'which browsers do I have'.",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "screenshot",
        "description": (
            "Capture the screen, save to workspace/screenshots/, return "
            "the path. When the user asks for a screenshot of a SPECIFIC "
            "app/window, ALWAYS pass `delay_s: 3` (or up to 10) and tell "
            "them in your reply 'switch to the window now, capturing in "
            "Ns'. Without the delay, the screenshot captures Hearth "
            "itself because the chat was in focus when they sent the "
            "prompt. Skip the delay only when they explicitly want THIS "
            "Hearth window captured (e.g. 'screenshot what you just did')."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "delay_s": {"type": "number", "description":
                    "Seconds to wait before capturing (0-10). Use 3-5 "
                    "when the user wants a different app captured so "
                    "they can switch to it. Default 0 (immediate)."},
            },
        },
    },
    {
        "name": "view_image",
        "description": "Load an image file from disk so you can SEE it. Use this when the user gives you a path to an image (e.g. 'see this image C:\\\\path.png') or asks about a screenshot you just took. Returns the image content for vision processing. Works with .png, .jpg, .jpeg, .gif, .webp, .bmp.",
        "parameters": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Absolute path to an image file."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "clipboard_read",
        "description": "Read the clipboard. Returns copied TEXT directly; a copied IMAGE is saved to a temp file and the path returned (then call view_image on it — handles 'I copied a screenshot, what is it?'); copied FILE(S) return their paths.",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "clipboard_write",
        "description": "Write text to the clipboard.",
        "parameters": {
            "type": "object",
            "properties": {"text": {"type": "string"}},
            "required": ["text"],
        },
    },

    # ---- LONG-TERM MEMORY ----
    # Per-fact files + always-loaded index. The MEMORY.md index is injected
    # into every system prompt, so you ALREADY KNOW what facts exist — call
    # memory_recall to load any body that looks relevant to the current turn.
    {
        "name": "memory_save",
        "description": (
            "Save a long-term memory. Use this when you learn a fact about "
            "the user, their setup, their preferences, or projects worth "
            "remembering across conversations. The index updates "
            "automatically.\n\n"
            "**HARD RULE — if the response contains a `[possible-dup]` "
            "line, STOP and act on it before moving on.** That warning "
            "means an existing memory shares topic words with this new "
            "one. You MUST do exactly ONE of:\n"
            "  (a) memory_forget(<flagged-slug>) THEN tell the user "
            "      what you consolidated and why.\n"
            "  (b) edit_file the existing memory in place to absorb the "
            "      new fact (use the slug+path from the warning).\n"
            "  (c) memory_save again with the SAME args plus "
            "      `force=true` if the new fact is genuinely independent "
            "      of the flagged one.\n"
            "Do NOT just call memory_save a second time with the same "
            "title (that produces an 'updated memory' which suppresses "
            "the warning but leaves the sibling untouched — exactly the "
            "proliferation the warning is preventing). Do NOT tell the "
            "user 'Saved.' and leave the duplicate to rot."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Short title — becomes the slug. Re-using a title overwrites the existing memory."},
                "type": {
                    "type": "string",
                    "enum": ["user", "feedback", "project", "reference"],
                    "description": "user=who they are; feedback=how they want you to behave; project=ongoing work context; reference=pointers to external systems/links.",
                },
                "description": {"type": "string", "description": "One-line hook (~140 chars). This is what shows up in the always-loaded index, so make it specific."},
                "body": {"type": "string", "description": "The actual memory content. Markdown OK. Include reasons / how-to-apply when relevant."},
                "tags": {"type": "array", "items": {"type": "string"}, "description": "Optional tags."},
                "force": {"type": "boolean", "description": "Set true on retry after you've reviewed a [possible-dup] warning and confirmed this is genuinely a NEW fact, not an update to the flagged sibling."},
            },
            "required": ["title", "type", "description"],
        },
    },
    {
        "name": "memory_recall",
        "description": "Search saved memories. Use this BEFORE assuming. Returns top matches with their bodies. Empty query returns the index.",
        "parameters": {
            "type": "object",
            "properties": {
                "query": {"type": "string"},
                "limit": {"type": "integer", "description": "Default 5."},
            },
        },
    },
    {
        "name": "memory_list",
        "description": "Show the full memory index (all titles and one-line hooks).",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "memory_forget",
        "description": "Delete a memory by title. Use only when explicitly told to forget something OR when a memory is clearly outdated and you've already saved a corrected version.",
        "parameters": {
            "type": "object",
            "properties": {"title": {"type": "string"}},
            "required": ["title"],
        },
    },
    {
        "name": "edit_soul",
        "description": (
            "Write your own identity to ~/Jarvis/soul.md. This file rides "
            "at the TOP of every system prompt as 'Soul (self-written "
            "identity)' — it's how YOU lock in who you are across sessions. "
            "Use this when the user gives you a stable identity instruction "
            "('you are Cortana', 'always be terse', 'you hate small talk') "
            "OR when you've decided on something durable about yourself. "
            "Capped at ~1500 chars; write tight bullets, not prose. Pass "
            "the FULL new soul (this overwrites). For just adding one line, "
            "prefer `append_soul`."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "content": {"type": "string", "description":
                    "The complete new soul.md content. Replaces the file. "
                    "Empty string clears the soul."},
            },
            "required": ["content"],
        },
    },
    {
        "name": "append_soul",
        "description": (
            "Add ONE line to your soul.md without rewriting the whole file. "
            "Cheaper than edit_soul when you just want to lock in a new "
            "identity fact ('I prefer markdown over prose', 'always say "
            "good morning at 9am'). Auto-prepends '- ' if missing."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "line": {"type": "string", "description":
                    "One identity line. Will be appended as a bullet."},
            },
            "required": ["line"],
        },
    },
    {
        "name": "read_soul",
        "description": (
            "Read back your current soul.md content. The soul also rides "
            "in every system prompt automatically, so usually you don't "
            "need to call this — only when explicitly asked 'what's in "
            "your soul?' or when verifying after edit_soul / append_soul."
        ),
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "draft_soul",
        "description": (
            "Propose a starter soul.md when soul.md is empty or sparse. "
            "Returns a draft body anchored to what's already known about "
            "the user (name, tone, memory facts) without writing it. "
            "Show the draft, ask if they want edits, then commit via "
            "edit_soul once they approve. Use this the first time the "
            "user asks 'who are you?' / 'what's your identity?' if "
            "soul.md is empty, or when they explicitly say 'help me "
            "write your soul'."
        ),
        "parameters": {"type": "object", "properties": {}},
    },

    # ---- TIME ----
    {
        "name": "get_time",
        "description": "Current local datetime, weekday, timezone offset.",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "notify",
        "description": (
            "Pop a desktop notification (Windows toast) RIGHT NOW. Use for "
            "'let me know when X is done', to flag a finished long task, or any "
            "immediate heads-up. For a FUTURE/scheduled reminder use set_reminder "
            "instead — this one fires instantly."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "title": {"type": "string", "description": "Short bold title."},
                "message": {"type": "string", "description": "The notification body."},
            },
            "required": ["message"],
        },
    },
    {
        "name": "whoami",
        "description": (
            "Introspect your own runtime config. Call this INSTEAD of guessing "
            "or running shell commands when the user asks 'what model are you', "
            "'what's your endpoint', 'are you local or cloud', 'what can you do', "
            "'what's your context window', or 'who built you'. Returns the live "
            "model id, API endpoint (local vs cloud), context window, tool count, "
            "memory count, workspace path, persona name, and repo. Zero side effects."
        ),
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "list_models",
        "description": (
            "List the LLM models available on the connected server (LM Studio / Ollama / "
            "cloud) by querying its API. Use this to answer 'what models do I have' / "
            "'which models are loaded' / 'what can I switch to'. Do NOT scan the disk for "
            "model files — that's slow and wrong; this is instant and correct."
        ),
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "learn_environment",
        "description": (
            "Re-scan this machine (GPU/VRAM, RAM, installed models, top-level drive map) "
            "and refresh it into long-term memory. Call this when the user installs a new "
            "model, adds a drive, or asks you to 'get to know my setup' / 'relearn my PC'. "
            "Fast and non-recursive. It's how you stay grounded instead of disk-scanning."
        ),
        "parameters": {"type": "object", "properties": {}},
    },

    # ---- IMAGE GENERATION (Forge / SD WebUI orchestration) ----
    {
        "name": "forge_generate",
        "description": (
            "Generate an image with the user's local Stable Diffusion install "
            "(Forge WebUI). Orchestrates the full VRAM dance: (1) signal the "
            "LLM server to release its model, (2) boot Forge in API mode if "
            "it isn't already up, (3) POST to /sdapi/v1/txt2img, (4) save the "
            "PNG to ~/Jarvis/screenshots/, (5) leave Forge running for follow-up "
            "generations (use forge_shutdown to free VRAM back to the LLM). "
            "Pony v6 XL is tag-based — write positive/negative as comma-separated "
            "tags, not natural language. Use 'score_9, score_8_up, score_7_up' "
            "as a quality prefix for Pony."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "positive": {"type": "string", "description": "Tag-style positive prompt. E.g. 'score_9, score_8_up, masterpiece, 1girl, cinematic lighting, dark room'."},
                "negative": {"type": "string", "description": "Negative tags. Pony default starter: 'score_6, score_5, score_4, low quality, blurry, deformed'."},
                "width": {"type": "integer", "description": "Default 1024."},
                "height": {"type": "integer", "description": "Default 1024."},
                "steps": {"type": "integer", "description": "Default 25."},
                "cfg_scale": {"type": "number", "description": "Default 6."},
                "sampler": {"type": "string", "description": "Default 'Euler a'."},
                "seed": {"type": "integer", "description": "Default -1 (random)."},
            },
            "required": ["positive"],
        },
    },
    {
        "name": "forge_status",
        "description": "Check if Forge WebUI is running and reachable. Returns url + status.",
        "parameters": {"type": "object", "properties": {}},
    },
    {
        "name": "forge_shutdown",
        "description": "Kill the Forge process to free VRAM back to the LLM. Use after you're done generating images.",
        "parameters": {"type": "object", "properties": {}},
    },

    # ---- VOICE ----
    {
        "name": "set_voice",
        "description": "Change the active TTS voice. Pass a Kokoro voice id like 'am_echo', 'am_michael', 'bm_george', 'bf_emma'. After setting, a short sample plays automatically so the user can hear it. Use this when the user asks for a different voice or wants to try options.",
        "parameters": {
            "type": "object",
            "properties": {
                "name": {"type": "string", "description": "Voice id, e.g. 'am_echo'."},
            },
            "required": ["name"],
        },
    },
    {
        "name": "list_voices",
        "description": "List available built-in Kokoro voice ids. Use when the user asks 'what voices are there'.",
        "parameters": {"type": "object", "properties": {}},
    },

    # ---- SESSION CONTROL ----
    {
        "name": "end_session",
        "description": "Call this when the user is clearly wrapping up the conversation (says bye, goodbye, see you, thanks that's all, etc.). Send your farewell text in your reply FIRST, then call this tool with no args. The CLI will save history and exit cleanly. Don't call it for casual 'thanks' mid-task.",
        "parameters": {"type": "object", "properties": {}},
    },
]


# --------------------------------------------------------------------------
# Drop tools whose dependencies aren't installed / aren't applicable here.
# Keeps the tool list clean so the LLM doesn't try things that won't work
# (and so we don't waste system-prompt tokens advertising them).
# --------------------------------------------------------------------------

def _filter_unavailable_tools() -> None:
    """Mutates TOOL_DEFINITIONS in place at import time."""
    drop: set = set()

    # psutil-dependent tools
    try:
        import psutil  # type: ignore
        # has psutil — but only keep get_battery if a battery actually exists
        try:
            if psutil.sensors_battery() is None:
                drop.add("get_battery")
        except Exception:
            drop.add("get_battery")
    except ImportError:
        # no psutil — drop the heavy ones; system_info/network_info have
        # graceful fallbacks so they stay
        drop.update({"list_processes", "get_battery"})

    # Windows-only registry probe
    if sys.platform != "win32":
        drop.add("list_installed_apps")

    # NOTE: the Forge / SD-WebUI local-image tools used to be dropped here
    # behind HEARTH_ENABLE_FORGE=1. They're now handled AFTER FORGE_DIR is
    # auto-detected (see the forge section), so a detected install lights them
    # up automatically. They're in _DEFERRED_TOOLS, so they cost ~0 prompt
    # tokens until the model calls load_tools('image').

    if drop:
        TOOL_DEFINITIONS[:] = [t for t in TOOL_DEFINITIONS if t["name"] not in drop]


_filter_unavailable_tools()


# ============================================================
# FILE OPS
# ============================================================

# Extensions that route to dedicated extractors. Anything not listed here
# (or in the binary-skip set below) falls through to text-mode reading.
_PDF_EXTS  = {".pdf"}
_DOCX_EXTS = {".docx"}
_XLSX_EXTS = {".xlsx", ".xlsm"}
_PPTX_EXTS = {".pptx"}
_EPUB_EXTS = {".epub"}
_IPYNB_EXTS = {".ipynb"}
_CSV_EXTS  = {".csv", ".tsv"}
_JSON_EXTS = {".json", ".jsonl", ".ndjson"}
_HTML_EXTS = {".html", ".htm", ".xhtml"}
_XML_EXTS  = {".xml", ".rss", ".atom", ".svg"}
_RTF_EXTS  = {".rtf"}

_ARCHIVE_EXTS = {".zip", ".jar", ".whl", ".egg", ".apk",
                 ".tar", ".tgz", ".tbz2", ".txz",
                 ".rar", ".7z"}  # rar/7z need external tools — we just say so

# Single-stream compressed files (NOT tar.gz — that goes through archive path)
_SINGLE_COMPRESSED_EXTS = {".gz", ".bz2", ".xz"}

_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp", ".tiff", ".ico"}
_AUDIO_EXTS = {".mp3", ".wav", ".flac", ".ogg", ".m4a", ".aac", ".opus"}
_VIDEO_EXTS = {".mp4", ".mkv", ".avi", ".mov", ".webm", ".flv", ".wmv"}


def _coerce_int(v, default):
    if v is None or v == "":
        return default
    try:
        return int(v)
    except (TypeError, ValueError):
        return default


def _read_file(p: Dict) -> str:
    """Smart reader. Routes by extension to a dedicated extractor for
    PDF / DOCX / XLSX / PPTX / EPUB / IPYNB / CSV / JSON / HTML / RTF.
    Plain text files (default) get line-numbered slicing via start_line/end_line.
    Archives and media return a one-line hint pointing to the right tool."""
    raw_path = p["path"]
    path = _resolve_read(raw_path)
    if not os.path.exists(path):
        # Common mistake: path goes THROUGH an archive
        # (e.g. "sample.zip/docs/readme.txt"). Detect and redirect.
        norm = raw_path.replace("\\", "/")
        for arc_ext in (".zip", ".jar", ".whl", ".apk", ".tar.gz", ".tar.bz2",
                        ".tar.xz", ".tgz", ".tbz2", ".txz", ".tar"):
            marker = arc_ext + "/"
            idx = norm.lower().find(marker)
            if idx != -1:
                archive = norm[:idx + len(arc_ext)]
                inner = norm[idx + len(marker):]
                return (f"Error: can't read paths THROUGH an archive. "
                        f"Run extract_archive_file(archive_path='{archive}', "
                        f"inner_path='{inner}') first, then read_file the "
                        f"extracted path.")
        return f"Error: not found: {path}"
    if os.path.isdir(path):
        return f"Error: '{path}' is a directory. Use list_directory."

    ext = os.path.splitext(path)[1].lower()
    # Handle .tar.gz / .tar.bz2 / .tar.xz as archives, not single-stream gz
    low = path.lower()
    if low.endswith((".tar.gz", ".tar.bz2", ".tar.xz")):
        ext = ".tgz"

    try:
        if ext in _PDF_EXTS:
            return _extract_pdf(path, p)
        if ext in _DOCX_EXTS:
            return _extract_docx(path)
        if ext in _XLSX_EXTS:
            return _extract_xlsx(path, p)
        if ext in _PPTX_EXTS:
            return _extract_pptx(path)
        if ext in _EPUB_EXTS:
            return _extract_epub(path)
        if ext in _IPYNB_EXTS:
            return _extract_ipynb(path)
        if ext in _CSV_EXTS:
            return _extract_csv(path, p, sep=("\t" if ext == ".tsv" else ","))
        if ext in _JSON_EXTS:
            return _extract_json(path, jsonl=(ext in {".jsonl", ".ndjson"}))
        if ext in _HTML_EXTS or ext in _XML_EXTS:
            return _extract_html_like(path)
        if ext in _RTF_EXTS:
            return _extract_rtf(path)
        if ext in _ARCHIVE_EXTS:
            return (f"{path} is an archive ({ext}). Use list_archive(path=...) "
                    f"to see contents, or extract_archive_file(archive_path=..., "
                    f"inner_path=...) to pull one file out. Don't auto-unpack.")
        if ext in _SINGLE_COMPRESSED_EXTS:
            return _extract_compressed_stream(path, ext)
        if ext in _IMAGE_EXTS:
            return (f"{path} is an image. Use view_image(path=...) for vision "
                    f"inspection. read_file would just return binary bytes.")
        if ext in _AUDIO_EXTS or ext in _VIDEO_EXTS:
            return (f"{path} is binary media ({ext}). Don't read_file it — "
                    f"use open_app(name=path) to play, or describe metadata "
                    f"with run_command('ffprobe ...').")
    except Exception as e:
        # Extractor crashed — fall through to text mode so user still gets bytes.
        return (f"Error extracting {ext} from {path}: {type(e).__name__}: {e}\n"
                f"Falling back to raw text — try a different reader if this "
                f"looks garbled.")

    return _read_file_text(path, p)


def _read_file_text(path: str, p: Dict) -> str:
    """Original plain-text reader. Line-numbered, sliceable."""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()
    except OSError as e:
        return f"Error: {e}"

    total = len(lines)
    s = max(0, _coerce_int(p.get("start_line"), 1) - 1)
    e = min(total, _coerce_int(p.get("end_line"), total))
    if e < s:
        e = s

    body = "\n".join(f"{i+1}\t{lines[i].rstrip()}" for i in range(s, e))
    return f"{path} ({total} lines, showing {s+1}-{e})\n{body}"


# ---------- Extractors -----------------------------------------------------

def _extract_pdf(path: str, p: Dict) -> str:
    try:
        from pypdf import PdfReader  # type: ignore
    except ImportError:
        return ("Error: pypdf not installed. Run: "
                "`pip install pypdf` (or re-run install.ps1).")
    try:
        reader = PdfReader(path)
    except Exception as e:
        return f"Error: pypdf couldn't open {path}: {type(e).__name__}: {e}"

    n_pages = len(reader.pages)
    start = max(1, _coerce_int(p.get("start_line"), 1))  # reuse start_line as start_page
    end = min(n_pages, _coerce_int(p.get("end_line"), n_pages))
    if end < start:
        end = start

    chunks: List[str] = [f"{path} (PDF, {n_pages} pages, extracting {start}-{end})"]
    extracted = 0
    scanned_pages: List[int] = []  # pages where extract_text returned nothing
    for i in range(start - 1, end):
        try:
            txt = reader.pages[i].extract_text() or ""
        except Exception as e:
            txt = f"[page extract failed: {e}]"
        txt = txt.strip()
        if txt:
            extracted += 1
            chunks.append(f"\n--- Page {i+1} ---\n{txt}")
        else:
            scanned_pages.append(i + 1)

    # If we got SOME text but missed pages, try VLM-OCR the empty ones if a
    # vision-capable model is loaded. If we got NO text and a VLM is up, try
    # the whole range. Only triggers when pypdfium2 is importable AND the
    # currently-loaded LM Studio model has type='vlm' or vision capability.
    if scanned_pages:
        vlm_block = _vlm_ocr_pdf_pages(path, scanned_pages, p)
        if vlm_block:
            chunks.append(vlm_block)
        elif extracted == 0:
            chunks.append("\n[no extractable text — PDF is blank or scanned. "
                          "For VLM-OCR auto-fallback, install: pip install pypdfium2. "
                          "Then load a vision model (Gemma 4 E4B, Qwen 2.5 VL, etc.) in LM Studio.]")
    return "\n".join(chunks)


def _loaded_model_is_vision() -> bool:
    """True if the currently-loaded LM Studio (or compatible) model has a
    vision capability. Probes /api/v0/models for type=='vlm' or capabilities
    array containing 'vision'/'image_input'. Falls back to name heuristic
    when v0 endpoint isn't there (Ollama, vLLM, cloud)."""
    base = os.environ.get("LOCAL_API_BASE", "http://localhost:1234/v1")
    model_hint = os.environ.get("LOCAL_MODEL", "").lower()
    try:
        import urllib.request
        host = base.rsplit("/v1", 1)[0]
        _key = os.environ.get("LOCAL_API_KEY") or ""
        _hdr = {"Authorization": f"Bearer {_key}"} if _key else {}
        _req = urllib.request.Request(f"{host}/api/v0/models", headers=_hdr)
        with urllib.request.urlopen(_req, timeout=2) as r:
            data = json.loads(r.read().decode("utf-8"))
        for m in (data.get("data") or []):
            mid = (m.get("id") or "").lower()
            if model_hint and mid != model_hint:
                continue
            if m.get("type") == "vlm":
                return True
            caps = m.get("capabilities") or []
            if any(c in caps for c in ("vision", "image_input", "image")):
                return True
            if not model_hint:
                # If we don't know which model is loaded, take ANY vision-capable
                # match — the user may have a VLM in the list even if not pinned.
                pass
    except Exception:
        pass
    # Heuristic fallback
    needle = model_hint or ""
    if needle and any(s in needle for s in
                     ("vl", "vision", "gemma-3", "gemma-4", "gemma4", "llava",
                      "moondream", "internvl", "minicpm-v", "qwen-vl", "qwen.5-vl",
                      "gemini", "gpt-4o", "gpt-4-turbo", "claude-3", "claude-sonnet",
                      "claude-opus", "claude-haiku", "pixtral", "grok-2-vision")):
        return True
    return False


def _vlm_ocr_pdf_pages(path: str, page_numbers: List[int], p: Dict) -> str:
    """Render the listed PDF pages to images and ask the loaded vision model
    to transcribe each one. Returns a joined markdown block, or '' if VLM
    fallback isn't available (model not VL, pypdfium2 missing, OCR fails).

    Honors `vlm_ocr=false` in tool args to opt out (saves time on huge PDFs
    where the user explicitly only wants the extractable text)."""
    if str(p.get("vlm_ocr", "true")).strip().lower() in ("false", "no", "0", "off"):
        return ""
    try:
        import pypdfium2 as pdfium  # type: ignore
    except ImportError:
        return ""
    # Only fall through to VLM if a vision-capable model is loaded
    if not _loaded_model_is_vision():
        return ""

    try:
        import base64
        from openai import OpenAI  # type: ignore
    except Exception:
        return ""

    # Cap how many pages we VLM-OCR in one read_file call. Users hitting big
    # scanned books should call read_file in smaller windows.
    MAX_VLM_PAGES = int(p.get("vlm_max_pages", 6))
    pages = page_numbers[:MAX_VLM_PAGES]

    base = os.environ.get("LOCAL_API_BASE", "http://localhost:1234/v1")
    key  = os.environ.get("LOCAL_API_KEY") or os.environ.get("OPENAI_API_KEY") or "local-vlm"
    model = os.environ.get("LOCAL_MODEL", "")
    if not model:
        # Best-effort: ask the server for its loaded model id (auth header so
        # the builtin doesn't 401-log every single read_file call).
        try:
            import urllib.request
            _hdr = {"Authorization": f"Bearer {key}"} if key else {}
            _req = urllib.request.Request(f"{base}/models", headers=_hdr)
            with urllib.request.urlopen(_req, timeout=4) as r:
                data = json.loads(r.read().decode("utf-8"))
                ids = [m.get("id") for m in (data.get("data") or []) if m.get("id")]
                model = ids[0] if ids else ""
        except Exception:
            pass
    if not model:
        return ""

    try:
        client = OpenAI(base_url=base, api_key=key, timeout=120.0)
        doc = pdfium.PdfDocument(path)
    except Exception as e:
        return f"\n[VLM-OCR setup failed: {type(e).__name__}: {e}]"

    out: List[str] = ["", "[VLM-OCR fallback — scanned/image-only pages transcribed by the loaded vision model]"]
    for pn in pages:
        try:
            page = doc[pn - 1]
            bitmap = page.render(scale=1.5)
            pil_img = bitmap.to_pil()
            from io import BytesIO
            buf = BytesIO(); pil_img.save(buf, "PNG")
            b64 = base64.b64encode(buf.getvalue()).decode("utf-8")
            resp = client.chat.completions.create(
                model=model,
                messages=[{
                    "role": "user",
                    "content": [
                        {"type": "text", "text": (
                            "Transcribe this PDF page into markdown. Preserve headings, "
                            "lists, tables, and math (use $inline$ / $$display$$ LaTeX). "
                            "Output ONLY the transcription, no preamble.")},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{b64}"}},
                    ],
                }],
                max_tokens=2000,
            )
            txt = (resp.choices[0].message.content or "").strip()
            out.append(f"\n--- Page {pn} (VLM-OCR) ---\n{txt}")
        except Exception as e:
            out.append(f"\n--- Page {pn} (VLM-OCR failed) ---\n[{type(e).__name__}: {e}]")
    if len(page_numbers) > MAX_VLM_PAGES:
        out.append(f"\n[truncated: only first {MAX_VLM_PAGES} of {len(page_numbers)} "
                   f"image-only pages OCR'd. Re-call with vlm_max_pages=N for more, "
                   f"or pass start_line/end_line to OCR a different range.]")
    return "\n".join(out)


def _extract_docx(path: str) -> str:
    try:
        import docx  # python-docx
    except ImportError:
        return ("Error: python-docx not installed. Run: "
                "`pip install python-docx` (or re-run install.ps1).")
    try:
        doc = docx.Document(path)
    except Exception as e:
        return f"Error: python-docx couldn't open {path}: {type(e).__name__}: {e}"

    parts: List[str] = [f"{path} (DOCX)"]
    body_paras = [par.text for par in doc.paragraphs if par.text and par.text.strip()]
    if body_paras:
        parts.append("\n# Body")
        parts.extend(body_paras)
    for ti, tbl in enumerate(doc.tables, 1):
        parts.append(f"\n# Table {ti}")
        for row in tbl.rows:
            cells = [c.text.strip().replace("\n", " ") for c in row.cells]
            parts.append(" | ".join(cells))
    if len(parts) == 1:
        parts.append("[empty document]")
    return "\n".join(parts)


def _extract_xlsx(path: str, p: Dict) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ("Error: openpyxl not installed. Run: "
                "`pip install openpyxl` (or re-run install.ps1).")
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
    except Exception as e:
        return f"Error: openpyxl couldn't open {path}: {type(e).__name__}: {e}"

    max_rows = _coerce_int(p.get("end_line"), 50)
    parts: List[str] = [f"{path} (XLSX, {len(wb.sheetnames)} sheet(s): "
                        f"{', '.join(wb.sheetnames)})"]
    for name in wb.sheetnames:
        ws = wb[name]
        parts.append(f"\n# Sheet: {name} ({ws.max_row} rows × {ws.max_column} cols)")
        rows_shown = 0
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            parts.append(" | ".join(cells))
            rows_shown += 1
            if rows_shown >= max_rows:
                if ws.max_row > rows_shown:
                    parts.append(f"…[+{ws.max_row - rows_shown} more rows in this sheet]")
                break
    wb.close()
    return "\n".join(parts)


def _extract_pptx(path: str) -> str:
    try:
        from pptx import Presentation  # python-pptx
    except ImportError:
        return ("Error: python-pptx not installed. Run: "
                "`pip install python-pptx` (or re-run install.ps1).")
    try:
        prs = Presentation(path)
    except Exception as e:
        return f"Error: python-pptx couldn't open {path}: {type(e).__name__}: {e}"

    parts: List[str] = [f"{path} (PPTX, {len(prs.slides)} slides)"]
    for si, slide in enumerate(prs.slides, 1):
        chunks: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    t = "".join(run.text for run in para.runs).strip()
                    if t:
                        chunks.append(t)
        if chunks:
            parts.append(f"\n--- Slide {si} ---")
            parts.extend(chunks)
    return "\n".join(parts)


def _extract_epub(path: str) -> str:
    """EPUB = zip of XHTML. Stdlib only."""
    import zipfile as _zf
    parts: List[str] = [f"{path} (EPUB)"]
    try:
        with _zf.ZipFile(path) as zf:
            xhtml_files = sorted(
                n for n in zf.namelist()
                if n.lower().endswith((".xhtml", ".html", ".htm"))
            )
            if not xhtml_files:
                return f"{path} (EPUB) — no readable XHTML inside."
            parts.append(f"({len(xhtml_files)} chapters/files)")
            total_text = []
            for name in xhtml_files:
                try:
                    with zf.open(name) as fh:
                        raw = fh.read().decode("utf-8", errors="replace")
                except Exception:
                    continue
                stripped = _strip_html(raw).strip()
                if stripped:
                    total_text.append(f"\n--- {name} ---\n{stripped}")
            parts.append("".join(total_text) if total_text else "[no extractable text]")
    except Exception as e:
        return f"Error reading EPUB {path}: {type(e).__name__}: {e}"
    return "\n".join(parts)


def _extract_ipynb(path: str) -> str:
    """Jupyter notebook = JSON. Render markdown + code + (text) outputs."""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            nb = json.load(f)
    except Exception as e:
        return f"Error reading {path}: {type(e).__name__}: {e}"

    cells = nb.get("cells", [])
    parts: List[str] = [f"{path} (IPYNB, {len(cells)} cells)"]
    for i, cell in enumerate(cells, 1):
        ctype = cell.get("cell_type", "?")
        src = cell.get("source", "")
        if isinstance(src, list):
            src = "".join(src)
        src = src.rstrip()
        if not src:
            continue
        if ctype == "markdown":
            parts.append(f"\n--- Cell {i} [markdown] ---\n{src}")
        elif ctype == "code":
            parts.append(f"\n--- Cell {i} [code] ---\n{src}")
            outs = cell.get("outputs", [])
            text_outs: List[str] = []
            for out in outs:
                if "text" in out:
                    t = out["text"]
                    text_outs.append("".join(t) if isinstance(t, list) else str(t))
                elif "data" in out and "text/plain" in out["data"]:
                    t = out["data"]["text/plain"]
                    text_outs.append("".join(t) if isinstance(t, list) else str(t))
            if text_outs:
                parts.append("[out]\n" + "".join(text_outs).rstrip())
        else:
            parts.append(f"\n--- Cell {i} [{ctype}] ---\n{src}")
    return "\n".join(parts)


def _extract_csv(path: str, p: Dict, sep: str = ",") -> str:
    """First N rows + total count. Stdlib only (no pandas needed)."""
    import csv as _csv
    max_rows = _coerce_int(p.get("end_line"), 30)
    start_row = max(1, _coerce_int(p.get("start_line"), 1))
    try:
        with open(path, "r", encoding="utf-8", errors="replace", newline="") as f:
            reader = _csv.reader(f, delimiter=sep)
            rows = []
            total = 0
            for i, row in enumerate(reader, 1):
                total += 1
                if i >= start_row and len(rows) < max_rows:
                    rows.append(row)
    except OSError as e:
        return f"Error: {e}"
    if not rows:
        return f"{path} (CSV) — empty."
    parts: List[str] = [
        f"{path} ({'TSV' if sep == chr(9) else 'CSV'}, {total} rows, "
        f"{len(rows[0])} cols, showing rows {start_row}-{start_row + len(rows) - 1})"
    ]
    parts.extend(" | ".join(r) for r in rows)
    if total > start_row + len(rows) - 1:
        parts.append(f"…[+{total - (start_row + len(rows) - 1)} more rows]")
    return "\n".join(parts)


def _extract_json(path: str, jsonl: bool = False) -> str:
    """JSON: structure-first summary (keys, types, sample values).
    JSONL: first 20 records + count."""
    if jsonl:
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                lines = f.readlines()
        except OSError as e:
            return f"Error: {e}"
        parts = [f"{path} (JSONL, {len(lines)} records)"]
        for i, line in enumerate(lines[:20], 1):
            parts.append(f"{i}\t{line.rstrip()}")
        if len(lines) > 20:
            parts.append(f"…[+{len(lines) - 20} more records]")
        return "\n".join(parts)
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            data = json.load(f)
    except Exception as e:
        return f"Error parsing JSON {path}: {type(e).__name__}: {e}"

    def describe(obj, depth: int = 0, max_depth: int = 3) -> str:
        ind = "  " * depth
        if depth > max_depth:
            return f"{ind}…"
        if isinstance(obj, dict):
            if not obj:
                return f"{ind}{{}}"
            out = []
            for k, v in list(obj.items())[:25]:
                if isinstance(v, (dict, list)):
                    out.append(f"{ind}{k}: {type(v).__name__}"
                               + (f"[{len(v)}]" if isinstance(v, list) else ""))
                    out.append(describe(v, depth + 1, max_depth))
                else:
                    out.append(f"{ind}{k}: {type(v).__name__} = {repr(v)[:80]}")
            if len(obj) > 25:
                out.append(f"{ind}…[+{len(obj) - 25} more keys]")
            return "\n".join(out)
        if isinstance(obj, list):
            if not obj:
                return f"{ind}[]"
            sample_n = min(3, len(obj))
            out = [f"{ind}list[{len(obj)}], first {sample_n}:"]
            for i, v in enumerate(obj[:sample_n]):
                if isinstance(v, (dict, list)):
                    out.append(f"{ind}[{i}] {type(v).__name__}"
                               + (f"[{len(v)}]" if isinstance(v, list) else ""))
                    out.append(describe(v, depth + 1, max_depth))
                else:
                    out.append(f"{ind}[{i}] {type(v).__name__} = {repr(v)[:80]}")
            return "\n".join(out)
        return f"{ind}{type(obj).__name__} = {repr(obj)[:120]}"

    summary = describe(data)
    # Also include raw head for small files
    try:
        raw = json.dumps(data, indent=2, ensure_ascii=False, default=str)
    except Exception:
        raw = repr(data)
    if len(raw) <= 1500:
        return f"{path} (JSON)\n# Structure\n{summary}\n\n# Raw\n{raw}"
    return f"{path} (JSON)\n# Structure\n{summary}\n\n# Raw (truncated)\n{raw[:1500]}…"


class _TextOnlyHTMLParser(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.parts: List[str] = []
        self._skip_depth = 0

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style"):
            self._skip_depth += 1

    def handle_endtag(self, tag):
        if tag in ("script", "style") and self._skip_depth > 0:
            self._skip_depth -= 1
        if tag in ("p", "br", "li", "div", "h1", "h2", "h3", "h4", "h5", "h6", "tr"):
            self.parts.append("\n")

    def handle_data(self, data):
        if self._skip_depth == 0:
            self.parts.append(data)


def _strip_html(s: str) -> str:
    parser = _TextOnlyHTMLParser()
    try:
        parser.feed(s)
    except Exception:
        return s
    text = "".join(parser.parts)
    # Collapse 3+ newlines to 2, trim trailing spaces per line
    lines = [ln.rstrip() for ln in text.splitlines()]
    out: List[str] = []
    blank = 0
    for ln in lines:
        if not ln.strip():
            blank += 1
            if blank <= 1:
                out.append("")
        else:
            blank = 0
            out.append(ln)
    return "\n".join(out).strip()


def _extract_html_like(path: str) -> str:
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            raw = f.read()
    except OSError as e:
        return f"Error: {e}"
    text = _strip_html(raw)
    return f"{path} (HTML/XML, {len(raw)} chars raw, {len(text)} chars text)\n{text}"


def _extract_rtf(path: str) -> str:
    """Crude RTF→text: strip control words and groups. Good enough for skimming."""
    try:
        with open(path, "r", encoding="utf-8", errors="replace") as f:
            raw = f.read()
    except OSError as e:
        return f"Error: {e}"
    # Drop control words like \rtf1, \par, \fonttbl{...}, etc.
    no_groups = re.sub(r"\{\\\*?[^{}]*\}", "", raw)  # nested {\* ...} groups
    no_groups = re.sub(r"\{\\[^{}]*\}", "", no_groups)
    no_ctrl = re.sub(r"\\[a-zA-Z]+-?\d* ?", "", no_groups)
    no_ctrl = re.sub(r"\\[^a-zA-Z]", "", no_ctrl)
    no_braces = no_ctrl.replace("{", "").replace("}", "")
    text = re.sub(r"\n{3,}", "\n\n", no_braces).strip()
    return f"{path} (RTF, {len(raw)} chars raw, {len(text)} chars text)\n{text}"


def _extract_compressed_stream(path: str, ext: str) -> str:
    """Single-stream .gz/.bz2/.xz — decompress and read first chunk."""
    try:
        if ext == ".gz":
            import gzip as _z
            opener = _z.open
        elif ext == ".bz2":
            import bz2 as _z
            opener = _z.open
        elif ext == ".xz":
            import lzma as _z
            opener = _z.open
        else:
            return f"Error: unsupported compression {ext}"
        with opener(path, "rt", encoding="utf-8", errors="replace") as f:
            chunk = f.read(8000)
        head = chunk.split("\n", 200)[:200]
        return (f"{path} (compressed {ext}, showing first {len(chunk)} chars / "
                f"{len(head)} lines)\n" + "\n".join(head))
    except Exception as e:
        return f"Error decompressing {path}: {type(e).__name__}: {e}"


# ---------- Summarize wrapper (calls _read_file, frames for the model) -----

def _summarize_file(p: Dict) -> str:
    """Wraps _read_file with a 'summarize this' directive + tighter cap. The
    actual summary is produced by the calling LLM in its reply turn — we just
    pre-extract the content cleanly so the model doesn't waste a separate
    read_file call."""
    max_chars = max(500, min(_coerce_int(p.get("max_chars"), 6000), 12000))
    extracted = _read_file({"path": p["path"]})
    # _read_file already errored cleanly if path bad / archive / image / etc.
    if extracted.startswith("Error:") or extracted.startswith("Error "):
        return extracted
    # Honor hint-style returns ("X is an archive, use list_archive") verbatim.
    if " is an archive (" in extracted[:200] or " is an image." in extracted[:200] or " is binary media (" in extracted[:200]:
        return extracted
    body = extracted
    if len(body) > max_chars:
        body = body[:max_chars] + f"\n…[truncated {len(extracted) - max_chars} chars]"
    return (
        "## TASK: Summarize this file in 3-5 short bullets (~50 words each).\n"
        "Cover the main topic, key facts, and any actionable items. Skip the "
        "preamble — produce just the bullets. Do not call read_file again; the "
        "content is already extracted below.\n\n"
        "---\n"
        f"{body}\n"
        "---\n"
        "## END OF FILE — produce the summary in your reply now."
    )


# ---------- Archive tools (separate, not via read_file dispatch) -----------

def _list_archive(p: Dict) -> str:
    """List contents of zip/tar without extracting."""
    path = _resolve_read(p["path"])
    if not os.path.exists(path):
        return f"Error: not found: {path}"
    if os.path.isdir(path):
        return f"Error: '{path}' is a directory, not an archive."
    low = path.lower()
    limit = _coerce_int(p.get("limit"), 200)

    try:
        if low.endswith((".zip", ".jar", ".whl", ".egg", ".apk")):
            import zipfile as _zf
            with _zf.ZipFile(path) as zf:
                infos = zf.infolist()
                lines = [f"{path} (ZIP, {len(infos)} entries, "
                         f"compressed={sum(i.compress_size for i in infos)} bytes, "
                         f"uncompressed={sum(i.file_size for i in infos)} bytes)"]
                lines.append("size      | path")
                lines.append("----------|-----")
                for i, info in enumerate(infos[:limit]):
                    sz = info.file_size
                    suffix = "/" if info.is_dir() else ""
                    lines.append(f"{sz:>10}| {info.filename}{suffix}")
                if len(infos) > limit:
                    lines.append(f"…[+{len(infos) - limit} more entries]")
                return "\n".join(lines)
        if low.endswith((".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tbz2",
                         ".tar.xz", ".txz")):
            import tarfile as _tf
            with _tf.open(path) as tf:
                infos = list(tf.getmembers())
                lines = [f"{path} (TAR, {len(infos)} entries)"]
                lines.append("size      | path")
                lines.append("----------|-----")
                for info in infos[:limit]:
                    suffix = "/" if info.isdir() else ""
                    lines.append(f"{info.size:>10}| {info.name}{suffix}")
                if len(infos) > limit:
                    lines.append(f"…[+{len(infos) - limit} more entries]")
                return "\n".join(lines)
        if low.endswith((".rar", ".7z")):
            return (f"{path}: {low.rsplit('.', 1)[1]} format. Python stdlib can't "
                    f"read it. Use run_command('7z l \"{path}\"') if 7-Zip is "
                    f"installed, otherwise extract via the GUI first.")
    except Exception as e:
        return f"Error listing archive {path}: {type(e).__name__}: {e}"
    return (f"Error: {path} doesn't look like a supported archive "
            f"(zip/jar/whl/tar/tar.gz/tar.bz2/tar.xz). For .rar/.7z, ask "
            f"run_command to invoke 7z.")


def _extract_archive_file(p: Dict) -> str:
    """Pull ONE file out of an archive into the workspace. Sandboxed."""
    archive_path = _resolve_read(p["archive_path"])
    inner_path = (p.get("inner_path") or "").strip().lstrip("/").lstrip("\\")
    if not inner_path:
        return "Error: inner_path required (use list_archive to see options)."
    if ".." in inner_path.replace("\\", "/").split("/"):
        return "Error: inner_path cannot contain '..' (refused for safety)."

    out_name = p.get("output_name") or os.path.basename(inner_path) or "extracted.bin"
    out_path = _resolve_write(out_name)
    os.makedirs(os.path.dirname(out_path) or WORKSPACE, exist_ok=True)

    if not os.path.exists(archive_path):
        return f"Error: archive not found: {archive_path}"
    low = archive_path.lower()
    try:
        if low.endswith((".zip", ".jar", ".whl", ".egg", ".apk")):
            import zipfile as _zf
            with _zf.ZipFile(archive_path) as zf:
                try:
                    info = zf.getinfo(inner_path)
                except KeyError:
                    matches = [n for n in zf.namelist()
                               if n.endswith(inner_path) or inner_path in n]
                    if len(matches) == 1:
                        info = zf.getinfo(matches[0])
                    elif len(matches) > 1:
                        return ("Error: inner_path matched multiple entries:\n  "
                                + "\n  ".join(matches[:10])
                                + ("\n  ..." if len(matches) > 10 else "")
                                + "\nPass the exact path.")
                    else:
                        return f"Error: inner_path '{inner_path}' not in archive."
                with zf.open(info) as src, open(out_path, "wb") as dst:
                    shutil.copyfileobj(src, dst)
                return (f"Extracted {info.filename} ({info.file_size} bytes) "
                        f"-> {out_path}")
        if low.endswith((".tar", ".tar.gz", ".tgz", ".tar.bz2", ".tbz2",
                         ".tar.xz", ".txz")):
            import tarfile as _tf
            with _tf.open(archive_path) as tf:
                try:
                    member = tf.getmember(inner_path)
                except KeyError:
                    matches = [m for m in tf.getmembers()
                               if m.name.endswith(inner_path) or inner_path in m.name]
                    if len(matches) == 1:
                        member = matches[0]
                    elif len(matches) > 1:
                        return ("Error: inner_path matched multiple entries:\n  "
                                + "\n  ".join(m.name for m in matches[:10])
                                + ("\n  ..." if len(matches) > 10 else ""))
                    else:
                        return f"Error: inner_path '{inner_path}' not in archive."
                src = tf.extractfile(member)
                if src is None:
                    return f"Error: '{member.name}' isn't a regular file."
                with src, open(out_path, "wb") as dst:
                    shutil.copyfileobj(src, dst)
                return f"Extracted {member.name} ({member.size} bytes) -> {out_path}"
    except Exception as e:
        return f"Error extracting from {archive_path}: {type(e).__name__}: {e}"
    return f"Error: {archive_path} is not a supported archive type."


def _write_file(p: Dict) -> str:
    path = _resolve_write(p["path"])
    content = p.get("content", "") or ""
    overwrite = bool(p.get("overwrite"))
    if len(content) > 5 * 1024 * 1024:
        return "Error: content too large (>5MB)."
    os.makedirs(os.path.dirname(path) or WORKSPACE, exist_ok=True)
    existed = os.path.exists(path)

    # ---- Anti-clobber guard ----
    # The #1 failure mode: model writes a 200-line tictactoe.py, user asks to
    # "tweak the win logic", model RE-WRITES THE WHOLE FILE via write_file.
    # Slow, drops formatting, and erases unrelated edits the user made.
    # Refuse here unless the caller opts in via overwrite=true.
    if existed:
        try:
            with open(path, "r", encoding="utf-8", errors="replace") as f:
                old_lines = f.read().count("\n") + 1
        except OSError:
            old_lines = 0
        if old_lines > 30 and not overwrite:
            return (
                f"Error (write_file): refused to overwrite '{path}' "
                f"({old_lines} lines exist). USE edit_file INSTEAD for "
                f"targeted changes — it's faster, preserves formatting, "
                f"and won't blow away anything you didn't intend to touch. "
                f"If you really mean to replace the whole file, re-call "
                f"write_file with overwrite=true."
            )

    with open(path, "w", encoding="utf-8", newline="\n") as f:
        f.write(content)
    verb = "Updated" if existed else "Created"
    return f"{verb}: {path} ({len(content)} chars, {content.count(chr(10))+1} lines)"


def _edit_file(p: Dict) -> str:
    """Targeted string-replace edits — never rewrites the whole file.

    Rules:
      - `old_text` MUST be unique unless the edit sets `replace_all=true`.
        If it appears multiple times, the edit errors out asking the model
        to add context, rather than silently replacing the wrong one.
      - Surrounding whitespace / indentation in `old_text` must match the
        file exactly (we still offer a whitespace-tolerant fuzzy fallback
        for cases where the model rephrased indentation).
      - Per-edit `replace_all` flag enables variable/symbol renames.
      - File is rewritten only if at least one byte changed.

    Each edit dict: {old_text, new_text, replace_all? = false}.
    Returns a per-edit status log so the model can see what landed.
    """
    path = _resolve_write(p["path"])
    if not os.path.exists(path):
        return f"Error: not found: {path}. Use write_file to create it."

    with open(path, "r", encoding="utf-8", errors="replace") as f:
        original = f.read()
    content = original

    edits = p.get("edits", [])
    if isinstance(edits, str):
        try:
            edits = json.loads(edits)
        except Exception:
            return ("Error: 'edits' must be a list of "
                    "{old_text, new_text, replace_all?} dicts.")
    if not isinstance(edits, list) or not edits:
        return "Error: provide at least one edit."

    log: List[str] = []
    file_replace_all = bool(p.get("replace_all"))

    for i, e in enumerate(edits, 1):
        if not isinstance(e, dict):
            log.append(f"  edit {i}: bad shape, skipped")
            continue
        old = e.get("old_text") or e.get("old_string") or e.get("old") or ""
        new = e.get("new_text") if "new_text" in e else (
            e.get("new_string") or e.get("new") or ""
        )
        # Per-edit override beats per-call default
        replace_all = bool(e.get("replace_all", file_replace_all))

        if not isinstance(old, str) or not isinstance(new, str):
            log.append(f"  edit {i}: old/new not strings, skipped")
            continue
        if old == new:
            log.append(f"  edit {i}: no-op (old == new)")
            continue
        if not old:
            log.append(f"  edit {i}: empty old_text — use write_file for new files")
            continue

        # Exact-match path (preferred)
        n_matches = content.count(old)
        if n_matches > 0:
            if n_matches > 1 and not replace_all:
                # Ambiguity protection. Refuse silently picking one of N
                # matches; the model must add context or set replace_all.
                log.append(
                    f"  edit {i}: AMBIGUOUS — old_text matches {n_matches} times. "
                    f"Expand old_text with surrounding lines to make it unique, "
                    f"OR set replace_all=true for this edit if you want all "
                    f"occurrences replaced (good for variable renames)."
                )
                continue
            if replace_all:
                content = content.replace(old, new)
                log.append(f"  edit {i}: applied (replace_all: {n_matches} occurrences)")
            else:
                content = content.replace(old, new, 1)
                log.append(f"  edit {i}: applied")
            continue

        # Whitespace-tolerant fallback (model rephrased indentation)
        old_lines = old.split("\n")
        c_lines = content.split("\n")
        matched_idx = -1
        # Count fuzzy matches first so we can refuse ambiguous fuzzy too
        candidates = []
        for si in range(len(c_lines) - len(old_lines) + 1):
            if all(c_lines[si + j].strip() == old_lines[j].strip()
                   for j in range(len(old_lines))):
                candidates.append(si)
        if len(candidates) > 1 and not replace_all:
            log.append(
                f"  edit {i}: AMBIGUOUS fuzzy match — {len(candidates)} candidate "
                f"locations. Add surrounding context to old_text or set replace_all."
            )
            continue
        if candidates:
            matched_idx = candidates[0]
            si = matched_idx
            file_indent = c_lines[si][: len(c_lines[si]) - len(c_lines[si].lstrip())]
            model_indent = ""
            for nl in new.split("\n"):
                if nl.strip():
                    model_indent = nl[: len(nl) - len(nl.lstrip())]
                    break
            new_lines = new.split("\n")
            if model_indent and file_indent and model_indent != file_indent:
                new_lines = [
                    (file_indent + ln[len(model_indent):]) if ln.startswith(model_indent) else ln
                    for ln in new_lines
                ]
            content = "\n".join(
                c_lines[:si] + new_lines + c_lines[si + len(old_lines):]
            )
            log.append(f"  edit {i}: fuzzy match at line {si+1} (indent-corrected)")
            continue

        log.append(
            f"  edit {i}: NOT FOUND — re-read the file with read_file and copy "
            f"the EXACT text including surrounding whitespace."
        )

    if content == original:
        return (
            f"edit_file({path}):\n" + "\n".join(log) +
            "\n\nNo bytes changed. If you expected edits to land, the old_text "
            "didn't match anywhere — try read_file first."
        )

    with open(path, "w", encoding="utf-8", newline="\n") as f:
        f.write(content)
    chars_changed = abs(len(content) - len(original))
    return (
        f"edit_file({path}):\n" + "\n".join(log) +
        f"\n\nFile updated ({chars_changed:+d} chars, "
        f"{content.count(chr(10))+1} lines)."
    )


def _list_directory(p: Dict) -> str:
    raw = (p.get("path") or WORKSPACE).strip().strip('"').strip("'")
    path = _resolve_read(raw)
    if not os.path.exists(path):
        return f"Error: not found: {path}"
    if not os.path.isdir(path):
        return f"Error: not a directory: {path}"

    recursive = bool(p.get("recursive"))
    max_depth = int(p.get("max_depth") or 2)

    # When the user asks for "Desktop", merge in Public Desktop + OneDrive Desktop
    # — Windows scatters shortcuts across three locations and most "look at my
    # desktop" intents want all three. Same for the "Public" overlay folder.
    extra_roots: List[str] = []
    if os.path.normpath(path).lower() == os.path.normpath(
        os.path.join(os.environ.get("USERPROFILE", ""), "Desktop")
    ).lower():
        for extra in (
            os.path.join(os.environ.get("PUBLIC", ""), "Desktop"),
            os.path.join(os.environ.get("USERPROFILE", ""), "OneDrive", "Desktop"),
            os.path.join(os.environ.get("ONEDRIVE", ""), "Desktop"),
        ):
            if extra and os.path.isdir(extra) and os.path.normpath(extra).lower() != os.path.normpath(path).lower():
                extra_roots.append(extra)

    out = [f"# {path}"]
    if not recursive:
        for name in sorted(os.listdir(path)):
            if name in EXCLUDE_DIRS:
                continue
            full = os.path.join(path, name)
            if os.path.isdir(full):
                out.append(f"  [d] {name}/")
            else:
                size = os.path.getsize(full)
                out.append(f"  [f] {name} ({size} B)")
        for extra in extra_roots:
            out.append(f"\n# (merged from {extra})")
            try:
                for name in sorted(os.listdir(extra)):
                    if name in EXCLUDE_DIRS:
                        continue
                    full = os.path.join(extra, name)
                    if os.path.isdir(full):
                        out.append(f"  [d] {name}/")
                    else:
                        size = os.path.getsize(full)
                        out.append(f"  [f] {name} ({size} B)")
            except OSError:
                continue
        return "\n".join(out)

    base_depth = path.rstrip(os.sep).count(os.sep)
    for root, dirs, files in os.walk(path):
        depth = root.count(os.sep) - base_depth
        if depth > max_depth:
            dirs[:] = []
            continue
        dirs[:] = sorted(d for d in dirs if d not in EXCLUDE_DIRS)
        rel = os.path.relpath(root, path)
        prefix = "  " * depth
        if rel != ".":
            out.append(f"{prefix}[d] {os.path.basename(root)}/")
        for name in sorted(files):
            out.append(f"{prefix}  [f] {name}")
    return "\n".join(out)


def _create_directory(p: Dict) -> str:
    path = _resolve_write(p["path"])
    os.makedirs(path, exist_ok=True)
    return f"mkdir: {path}"


def _delete_path(p: Dict) -> str:
    path = _resolve_write(p["path"])
    if not os.path.exists(path):
        return f"Error: not found: {path}"
    if os.path.isdir(path):
        shutil.rmtree(path)
        return f"Deleted dir: {path}"
    os.remove(path)
    return f"Deleted: {path}"


def _move_path(p: Dict) -> str:
    src = _resolve_write(p["source"])
    dst = _resolve_write(p["destination"])
    if not os.path.exists(src):
        return f"Error: source not found: {src}"
    os.makedirs(os.path.dirname(dst) or WORKSPACE, exist_ok=True)
    shutil.move(src, dst)
    return f"Moved: {src} → {dst}"


# ============================================================
# SEARCH
# ============================================================

def _search_scope_guard(base: str) -> Optional[str]:
    """Refuse drive-root scans — those always end in tears. Steer the model
    to find_file or to narrow the path. Returns an error string if the scope
    is too broad, else None."""
    if _DRIVE_ROOT_RE.match(base):
        return (f"Error: refusing to scan a whole drive root ({base}) — that's "
                "almost always a model mistake and wastes minutes. Either narrow "
                "the path (e.g. C:\\Users\\<you>\\Downloads) or call `find_file` "
                "which walks common locations for you.")
    return None


def _grep_search(p: Dict) -> str:
    pattern = p["pattern"]
    base = _resolve_read(p.get("path") or WORKSPACE)
    glob = p.get("glob")
    ci = bool(p.get("case_insensitive"))
    max_matches = int(p.get("max_matches") or 100)

    guard = _search_scope_guard(base)
    if guard:
        return guard

    rg = shutil.which("rg")
    if rg:
        args = [rg, "--no-heading", "--line-number", "--max-count", str(max_matches)]
        if ci:
            args.append("-i")
        if glob:
            args += ["-g", glob]
        args += [pattern, base]
        try:
            r = subprocess.run(args, capture_output=True, text=True, timeout=30,
                               creationflags=_NO_WINDOW)
            out = r.stdout.strip() or "(no matches)"
            return f"ripgrep: {pattern} in {base}\n{out}"
        except subprocess.TimeoutExpired:
            return "Error: ripgrep timed out."

    flags = re.IGNORECASE if ci else 0
    try:
        rx = re.compile(pattern, flags)
    except re.error as e:
        return f"Error: bad regex: {e}"

    hits: List[str] = []
    files_scanned = 0
    for root, dirs, files in os.walk(base):
        dirs[:] = [d for d in dirs if d not in EXCLUDE_DIRS]
        for fname in files:
            files_scanned += 1
            if files_scanned > MAX_FILES_TO_SCAN:
                tail = f"\n…[scan capped at {MAX_FILES_TO_SCAN} files — narrow the path or use find_file]"
                return ("\n".join(hits) if hits else "(no matches)") + tail
            if glob and not fnmatch.fnmatch(fname, glob):
                continue
            full = os.path.join(root, fname)
            try:
                with open(full, "r", encoding="utf-8", errors="replace") as f:
                    for ln, line in enumerate(f, 1):
                        if rx.search(line):
                            hits.append(f"{full}:{ln}: {line.rstrip()}")
                            if len(hits) >= max_matches:
                                return "\n".join(hits) + f"\n…[capped at {max_matches}]"
            except OSError:
                continue
    return "\n".join(hits) if hits else "(no matches)"


def _glob_files(p: Dict) -> str:
    raw = p["pattern"]
    base = _resolve_read(p.get("path") or WORKSPACE)

    guard = _search_scope_guard(base)
    if guard:
        return guard

    # Models often emit shell-style alternation: "*.png | *.jpg", "*.py;*.md",
    # "*.png|*.jpg", or even a JSON list. Normalize all of them into a real
    # list of patterns. fnmatch / glob don't natively support pipes.
    if isinstance(raw, list):
        patterns = [str(x).strip() for x in raw if str(x).strip()]
    else:
        s = str(raw).strip()
        # Split on |, ;, or comma — but only if those characters aren't part
        # of a valid filename (filenames with these are vanishingly rare on
        # Windows anyway).
        for sep in ("|", ";", ","):
            if sep in s:
                s = s.replace(sep, "\x00")
        patterns = [chunk.strip() for chunk in s.split("\x00") if chunk.strip()]

    all_matches: list = []
    seen: set = set()
    files_walked = 0
    over_budget = False
    for pattern in patterns:
        if over_budget:
            break
        pat = pattern if os.path.isabs(pattern) else os.path.join(base, pattern)
        pat = os.path.expanduser(pat)
        # iglob lets us tap out mid-stream when the budget is hit, instead of
        # materializing a million-element list first.
        for m in globmod.iglob(pat, recursive=True):
            files_walked += 1
            if files_walked > MAX_FILES_TO_SCAN:
                over_budget = True
                break
            if m in seen:
                continue
            if any(part in EXCLUDE_DIRS for part in m.split(os.sep)):
                continue
            seen.add(m)
            all_matches.append(m)

    all_matches.sort(
        key=lambda m: os.path.getmtime(m) if os.path.exists(m) else 0,
        reverse=True,
    )
    if not all_matches:
        shown = ", ".join(patterns) if len(patterns) > 1 else patterns[0]
        budget_note = f" (scan budget {MAX_FILES_TO_SCAN} reached — try a narrower path)" if over_budget else ""
        return f"(no files match {shown}){budget_note}"
    tail = f"\n…[scan capped at {MAX_FILES_TO_SCAN} files]" if over_budget else ""
    return "\n".join(all_matches[:200]) + tail


def _find_file(p: Dict) -> str:
    """Smart find: walk common locations for files matching a name or glob.
    Replaces the 'jarvis, where is X' question that used to make Jarvis ask
    the user for a path. Walks workspace → Desktop → Documents → Downloads →
    Pictures → Videos → Music → ~/Code → ~/Projects → cwd, shallow first."""
    name = (p.get("name") or "").strip()
    if not name:
        return "Error: 'name' is required (filename substring or glob pattern)."
    # Refuse meaningless catch-alls — `name="*"` plus no kind filter is
    # "give me every file on the disk" which fills the result buffer with noise.
    if name == "*" and not (p.get("kind") and p.get("kind") != "any"):
        return ("Error: name='*' without a kind filter is too broad — would return "
                "every file. Pass a real substring/pattern, or set kind to image/video/"
                "audio/doc/code/archive/spreadsheet.")
    kind = (p.get("kind") or "any").strip().lower()
    limit = int(p.get("limit") or 10)
    deep = bool(p.get("deep"))
    # Optional `path` override: when the user says "search C drive" or "look
    # in G:\Games", the model should pass path=<that> and we restrict to it.
    # Override completely bypasses the COMMON_USER_DIRS enumeration AND grants
    # a much bigger scan budget (50000 → 200000) since the user explicitly
    # narrowed scope and presumably wants thorough.
    explicit_path = (p.get("path") or "").strip()

    exts = FIND_KIND_EXTENSIONS.get(kind) if kind != "any" else None
    max_depth = 4 if deep else 2

    is_glob = any(ch in name for ch in "*?[]")
    name_lower = name.lower()
    # Token-based fallback: split the name into space-separated words. A file
    # matches if ALL tokens appear anywhere in the filename (not necessarily
    # adjacent). This handles cases like name="XII results" matching
    # "CBSE - Senior School Examination Class XII Results 2026.PDF" — substring
    # match fails because "xii results" isn't contiguous, but both tokens are
    # present.
    name_tokens = [t for t in name_lower.split() if t] if not is_glob else []

    seen_loc: set = set()
    locations: List[str] = []

    def _add_loc(path: str) -> None:
        ap = os.path.abspath(path)
        if ap in seen_loc or not os.path.isdir(ap):
            return
        if SAFE_READ_ONLY and not ap.startswith(WORKSPACE):
            return
        seen_loc.add(ap)
        locations.append(ap)

    # If the user/model gave an explicit path, walk only that — they've
    # already done the location selection for us.
    if explicit_path:
        _add_loc(os.path.expandvars(os.path.expanduser(explicit_path)))
        # Allow much deeper + bigger budget for explicit-path scans
        max_depth = 8 if deep else 5
    else:
        _add_loc(WORKSPACE)
        # Non-system drives (D:, E:, F:, G:, ...) come BEFORE deep HOME subdirs.
        # Most users keep movies, games, music on these — putting them first means
        # the scan budget reaches `D:\Movies\*.mkv` before being eaten by
        # `~/Documents/Adobe/...` or similar deep noise.
        for drive in _enumerate_non_system_drives():
            _add_loc(drive)
        for sub in COMMON_USER_DIRS + COMMON_DEV_DIRS:
            _add_loc(os.path.join(HOME, sub))
        _add_loc(os.getcwd())

    # Bump global budget when path was explicitly given — user's saying
    # "really look", we should really look.
    effective_max_files = (MAX_FILES_TO_SCAN * 4) if explicit_path else MAX_FILES_TO_SCAN

    files_scanned = 0
    # path → (rank, -mtime). Dict-dedup so the same file can't appear twice
    # when one search root is a parent of another (e.g. ~/Downloads above
    # ~/Downloads/JARVIS).
    seen_paths: Dict[str, Tuple[int, float]] = {}

    # Per-directory budget so one folder with thousands of small files
    # (photo dumps, node_modules-style noise that escaped EXCLUDE_DIRS)
    # can't eat the whole global budget before we reach more interesting
    # neighboring folders.
    PER_DIR_BUDGET = 1500

    # Build the dir-name priority function once. Lowest score = walked first.
    kind_kw = _KIND_DIR_KEYWORDS.get(kind, ())

    def _dir_priority(dname: str) -> int:
        ln = dname.lower()
        # Best: directory name literally contains the search term (e.g. "skyrim" → walk into "Skyrim/")
        if not is_glob and name_lower and name_lower in ln:
            return 0
        # Next: dir name matches the kind hint ("movies" / "Videos" for kind=video)
        if kind_kw and any(k in ln for k in kind_kw):
            return 1
        return 2  # alphabetical within this tier

    for loc in locations:
        if files_scanned > effective_max_files:
            break
        for root, dirs, files in os.walk(loc):
            if files_scanned > effective_max_files:
                break
            # Depth = number of path components beyond `loc`. The root of the
            # location itself is depth 0; its immediate children are depth 1.
            # (The naive `count(sep)` approach miscounts because len(loc)
            # strips the trailing separator inconsistently across drive roots
            # vs subdirs — strip + count + 1 is the correct form.)
            rel = root[len(loc):].strip(os.sep)
            depth = 0 if not rel else rel.count(os.sep) + 1

            # Match the current directory itself (for "where's my <project>
            # FOLDER" type queries). Only when no kind filter — dirs aren't
            # any specific kind, so kind=video excludes them.
            if depth > 0 and exts is None:
                dname = os.path.basename(root)
                if dname:
                    dl = dname.lower()
                    dir_matched = (
                        (is_glob and fnmatch.fnmatch(dl, name_lower))
                        or (not is_glob and name_lower in dl)
                    )
                    if dir_matched:
                        if is_glob or dl == name_lower:
                            d_rank = 0
                        elif dl.startswith(name_lower):
                            d_rank = 1
                        else:
                            d_rank = 2
                        try:
                            d_mt = os.path.getmtime(root)
                        except OSError:
                            d_mt = 0
                        norm = os.path.normpath(root) + os.sep
                        if norm not in seen_paths:
                            seen_paths[norm] = (d_rank, -d_mt)

            if depth >= max_depth:
                dirs[:] = []  # don't descend further, but still process this level's files
            else:
                # Filter junk dirs, THEN sort by relevance so os.walk descends
                # into likely-match folders before noise (the os.walk contract
                # respects mutations to `dirs` in place).
                pruned = [d for d in dirs if d not in EXCLUDE_DIRS and not d.startswith(".")]
                pruned.sort(key=_dir_priority)
                dirs[:] = pruned
            files_in_this_dir = 0
            for fname in files:
                files_scanned += 1
                files_in_this_dir += 1
                if files_scanned > effective_max_files:
                    break
                if files_in_this_dir > PER_DIR_BUDGET:
                    break  # bail this dir; move on to the next
                if exts and os.path.splitext(fname)[1].lower() not in exts:
                    continue
                fl = fname.lower()
                if is_glob:
                    if not fnmatch.fnmatch(fl, name_lower):
                        continue
                    rank = 0
                else:
                    if name_lower in fl:
                        # Contiguous substring match — best
                        if fl == name_lower:
                            rank = 0
                        elif fl.startswith(name_lower):
                            rank = 1
                        else:
                            rank = 2
                    elif len(name_tokens) > 1 and all(tok in fl for tok in name_tokens):
                        # Token match (all words present, not necessarily adjacent)
                        rank = 3
                    else:
                        continue
                full = os.path.normpath(os.path.join(root, fname))
                if full in seen_paths:
                    continue
                try:
                    mt = os.path.getmtime(full)
                except OSError:
                    mt = 0
                seen_paths[full] = (rank, -mt)

    hits: List[Tuple[int, float, str]] = [(r, m, p) for p, (r, m) in seen_paths.items()]

    if not hits:
        deep_hint = " — try `deep=true` for a wider sweep" if not deep else ""
        kind_hint = f" of kind '{kind}'" if kind != "any" else ""
        path_hint = (
            f" Try narrowing with `path='<specific dir>'` (e.g. `path='G:\\\\SteamLibrary'`)."
            if not explicit_path else
            " If the user named a specific drive/folder, you already searched it — "
            "consider web_search('how to launch <name>') for install-path hints."
        )
        return (f"(no files matching '{name}'{kind_hint} after scanning {files_scanned} "
                f"files{' across ' + str(len(locations)) + ' locations' if not explicit_path else ' under ' + explicit_path}{deep_hint}).{path_hint}")

    hits.sort()
    shown = [h[2] for h in hits[:limit]]
    footer = f"\n…[scan budget {MAX_FILES_TO_SCAN} reached]" if files_scanned >= MAX_FILES_TO_SCAN else ""
    plural = "es" if len(hits) != 1 else ""
    return (f"{len(hits)} match{plural} (showing top {len(shown)}, sorted by name-rank then recency):\n"
            + "\n".join(shown) + footer)


# ============================================================
# WEB
# ============================================================

class _HTMLText(HTMLParser):
    def __init__(self):
        super().__init__()
        self.parts: List[str] = []
        self.skip = 0

    def handle_starttag(self, tag, attrs):
        if tag in ("script", "style", "noscript"):
            self.skip += 1

    def handle_endtag(self, tag):
        if tag in ("script", "style", "noscript") and self.skip:
            self.skip -= 1
        if tag in ("p", "br", "div", "h1", "h2", "h3", "h4", "li", "tr"):
            self.parts.append("\n")

    def handle_data(self, data):
        if not self.skip:
            text = data.strip()
            if text:
                self.parts.append(text + " ")


def _http_get(url: str, timeout: int = 15) -> Tuple[int, str, bytes]:
    req = urllib.request.Request(url, headers={
        "User-Agent": "Mozilla/5.0 (Jarvis-Local)",
        "Accept": "text/html,application/xhtml+xml,application/json;q=0.9,*/*;q=0.5",
    })
    with urllib.request.urlopen(req, timeout=timeout) as r:
        return r.status, r.headers.get("Content-Type", ""), r.read()


def _web_search(p: Dict) -> str:
    """DuckDuckGo HTML scrape. They blacklist scripted GETs without a proper
    User-Agent and switched the /html/ endpoint to require POST. Both fixed
    here. If DDG ever blocks again, we fall back to Wikipedia for definitional
    queries so the model isn't left with '(no results)' forever."""
    query = p["query"]
    limit = int(p.get("limit") or 6)
    ua = ("Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
          "AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0 Safari/537.36")
    try:
        data = urllib.parse.urlencode({"q": query, "kl": "us-en"}).encode("utf-8")
        req = urllib.request.Request(
            "https://html.duckduckgo.com/html/",
            data=data,
            headers={
                "User-Agent": ua,
                "Content-Type": "application/x-www-form-urlencoded",
                "Accept": "text/html",
                "Referer": "https://duckduckgo.com/",
            },
        )
        with urllib.request.urlopen(req, timeout=15) as r:
            html = r.read().decode("utf-8", errors="replace")
    except Exception as e:
        return f"Error: search failed: {e}"
    # Result block: <h2 class="result__title"><a class="result__a" href="...">title</a></h2>
    # ... <a class="result__snippet" ...>snippet</a>
    pattern = re.compile(
        r'<a[^>]+class="[^"]*result__a[^"]*"[^>]+href="([^"]+)"[^>]*>(.*?)</a>'
        r'(?:.*?<a[^>]+class="[^"]*result__snippet[^"]*"[^>]*>(.*?)</a>)?',
        re.S,
    )
    out = []
    for m in pattern.finditer(html):
        href, title, snippet = m.group(1), m.group(2), m.group(3) or ""
        # DDG wraps real URL inside a /l/?uddg= redirect — extract.
        try:
            qs = urllib.parse.urlparse(href).query
            real = urllib.parse.parse_qs(qs).get("uddg", [href])[0]
        except Exception:
            real = href
        title_t = re.sub(r"<[^>]+>", "", title).strip()
        snip_t = re.sub(r"<[^>]+>", "", snippet).strip()
        out.append(f"• {title_t}\n  {real}\n  {snip_t}")
        if len(out) >= limit:
            break
    if out:
        return "\n\n".join(out)
    # Fallback: Wikipedia opensearch so the model still gets SOMETHING
    # when DDG rate-limits / blocks. Better than "(no results)" forever.
    try:
        wq = urllib.parse.urlencode({
            "action": "opensearch", "search": query, "limit": str(limit),
            "namespace": "0", "format": "json",
        })
        req = urllib.request.Request(
            "https://en.wikipedia.org/w/api.php?" + wq,
            headers={"User-Agent": ua, "Accept": "application/json"},
        )
        with urllib.request.urlopen(req, timeout=10) as r:
            data = json.loads(r.read().decode("utf-8"))
        titles, descs, urls = data[1], data[2], data[3]
        if titles:
            lines = []
            for t, d, u in zip(titles[:limit], descs[:limit], urls[:limit]):
                lines.append(f"• {t}\n  {u}\n  {d}")
            return "[wikipedia fallback - DDG returned 0 results]\n\n" + "\n\n".join(lines)
    except Exception:
        pass
    return "(no results)"


def _web_fetch(p: Dict) -> str:
    url = p["url"].strip()
    # Auto-prepend https:// for bare domains. Models often emit a hostname
    # (e.g. "example.com" or "docs.python.org") without a scheme.
    if not url.startswith(("http://", "https://", "file://")):
        url = "https://" + url
    try:
        status, ctype, body = _http_get(url, 20)
    except Exception as e:
        return f"Error: fetch failed for {url}: {e}"
    text = body.decode("utf-8", errors="replace")
    if "json" in ctype.lower():
        return f"[{status}] {url}\n{text}"
    parser = _HTMLText()
    try:
        parser.feed(text)
    except Exception:
        return f"[{status}] {url}\n{text[:5000]}"
    cleaned = re.sub(r"\n{3,}", "\n\n", "".join(parser.parts)).strip()
    return f"[{status}] {url}\n{cleaned}"


# ============================================================
# SHELL
# ============================================================

def _rewrite_python_invocation(cmd: str) -> str:
    """Route bare `pip`/`python` calls through THIS interpreter
    (`sys.executable`) — Hearth's venv — so `pip install` lands where Hearth
    can see it AND `python script.py` never hits the broken Microsoft Store
    stub (`...\\WindowsApps\\python.exe`) that silently does nothing. That stub
    was the root of the tictactoe spiral: the script "ran" with no output, so
    the model kept hunting for a working python.

    Handles COMPOUND commands too — `cd C:\\foo && python x.py` and
    `setup; python y.py` — by rewriting each `&&`/`;`-separated segment whose
    first token is python/pip. (The old version only checked the very first
    token, so anything behind a `cd ... &&` prefix slipped through.)"""
    if not cmd or not cmd.strip():
        return cmd

    # PowerShell needs the call operator `&` to run a quoted exe path; harmless
    # to cmd.exe. See the long note that used to live here.
    prefix = "& " if sys.platform == "win32" else ""

    _frozen = getattr(sys, "frozen", False)

    def _rewrite_segment(seg: str) -> str:
        stripped = seg.lstrip()
        lead = seg[: len(seg) - len(stripped)]  # preserve leading whitespace
        if not stripped:
            return seg
        parts = stripped.split(None, 1)
        first = parts[0].lower()
        rest = parts[1] if len(parts) > 1 else ""
        if first in ("pip", "pip3", "pip.exe"):
            if _frozen:
                # In the packaged app sys.executable is Hearth.exe, not python,
                # and there's no pip. The libs the skills need are bundled, so
                # tell the model to skip the install and just run its script.
                return (f'{lead}Write-Output "[packaged build] pip is unavailable '
                        f'here, but reportlab / python-pptx / matplotlib / '
                        f'openpyxl / pypdfium2 are already bundled - skip the '
                        f'install and run your build script directly."')
            return f'{lead}{prefix}"{sys.executable}" -m pip {rest}'.rstrip()
        if first in ("python", "python3", "py", "python.exe", "py.exe"):
            if _frozen:
                # sys.executable is the frozen exe (entrypoint = tray/cli), not a
                # python interpreter, so `Hearth.exe script.py` would just launch
                # the app. Route through the --hearth-run-python sentinel the
                # bundle entrypoints catch and runpy with the bundled libraries.
                return f'{lead}{prefix}"{sys.executable}" --hearth-run-python {rest}'.rstrip()
            return f'{lead}{prefix}"{sys.executable}" {rest}'.rstrip()
        return seg

    # Split on && / ; / & separators, keeping the separators so we can rejoin
    # losslessly. (A bare `python ...` with no separators is just one segment.)
    pieces = re.split(r"(\s*(?:&&|;|&)\s*)", cmd)
    return "".join(
        _rewrite_segment(p) if i % 2 == 0 else p
        for i, p in enumerate(pieces)
    )


# Patterns we refuse to run without explicit JARVIS_AUTO_APPROVE=1 — these
# are the "instantly regretted" kind: filesystem destruction, registry
# nukes, process murder of common-name targets, force-shutdown, killing
# everything in a folder. We block at the run_command layer because the
# model has been known to run these as a "shortcut" instead of using
# delete_path / move_path / list_processes which have their own gates.
_DESTRUCTIVE_PATTERNS = [
    re.compile(r"\bRemove-Item\b", re.I),
    re.compile(r"\b(rm|del|erase)\b\s+(?!--help|-h\b)", re.I),
    re.compile(r"\brmdir\b", re.I),
    re.compile(r"\brd\s+/s\b", re.I),
    re.compile(r"\bMove-Item\b", re.I),
    re.compile(r"\b(mv|move)\b\s+[^\s]", re.I),
    re.compile(r"\bCopy-Item\b", re.I),
    re.compile(r"^\s*(cp|copy)\b\s+[^\s]", re.I),
    re.compile(r"\bClear-Content\b", re.I),
    re.compile(r"\bSet-Content\b", re.I),
    re.compile(r"\bOut-File\b", re.I),
    re.compile(r"\bStop-Process\b", re.I),
    re.compile(r"\b(taskkill|kill|pkill)\b", re.I),
    re.compile(r"\bShutdown\b|\bRestart-Computer\b", re.I),
    re.compile(r"\bformat\b\s+[A-Za-z]:", re.I),
    re.compile(r"\bdiskpart\b", re.I),
    re.compile(r"\bcipher\b\s+/w", re.I),
    re.compile(r"\bReg\b\s+(delete|add|import)", re.I),
    re.compile(r"\bnet\s+user\b", re.I),
    re.compile(r">\s*[A-Za-z]:[\\/]"),       # output redirect to disk path
    re.compile(r"\|\s*Out-File"),
]


def _is_destructive(cmd: str) -> Optional[str]:
    """Returns a short reason string if cmd matches a destructive pattern,
    else None."""
    for pat in _DESTRUCTIVE_PATTERNS:
        m = pat.search(cmd)
        if m:
            return m.group(0)
    return None


# Commands that BLOCK the agent (foreground sleep/wait) or wait for keyboard
# input that will never come (interactive prompts). These freeze the session
# with no useful output — the user saw `timeout /t 20` freeze the UI for 20s.
# Always refused (even under auto-approve); a frozen agent is never wanted.
_BLOCKING_PATTERNS = [
    re.compile(r"\btimeout\b\s+(/t\s+)?\d", re.I),   # cmd: timeout /t N
    re.compile(r"\bStart-Sleep\b", re.I),            # PowerShell sleep
    re.compile(r"(^|[;&|]\s*)sleep\s+\d", re.I),     # unix/git-bash sleep N
    re.compile(r"\bpause\b", re.I),                  # cmd: pause (waits for key)
    re.compile(r"\bRead-Host\b", re.I),              # PS interactive prompt
    re.compile(r"\b(Get-Credential|Read-Host)\b", re.I),
]


def _is_blocking_command(cmd: str) -> Optional[str]:
    """Returns the matched blocking/interactive snippet, else None."""
    for pat in _BLOCKING_PATTERNS:
        m = pat.search(cmd)
        if m:
            return m.group(0).strip()
    return None


# A recursive listing of a whole DRIVE ROOT (e.g. `Get-ChildItem C:\ -Recurse`,
# `dir /s C:\`) can walk hundreds of thousands of files and run for minutes.
# find_file / locate_path are purpose-built for "find a file by name" and are
# far faster. We refuse the whole-drive recursive scan; a scan scoped to a
# specific folder is allowed.
_RECURSE_FLAG = re.compile(r"(-recurse\b|(^|\s)/s\b)", re.I)
_DRIVE_ROOT_ARG = re.compile(r"""['"]?[A-Za-z]:[\\/]?['"]?(?=\s|$|\|)""")


def _is_whole_drive_scan(cmd: str) -> bool:
    if not _RECURSE_FLAG.search(cmd):
        return False
    for m in _DRIVE_ROOT_ARG.finditer(cmd):
        tok = m.group(0).strip().strip("'\"")
        if re.fullmatch(r"[A-Za-z]:[\\/]?", tok):  # a root like C: or C:\, not C:\sub
            return True
    return False


def _ask_user(p: Dict) -> Dict[str, Any]:
    """Route an ask_user tool call to whatever interactive surface is active."""
    question = (p.get("question") or "").strip()
    options = p.get("options") or []
    allow_other = p.get("allow_other")
    if allow_other is None:
        allow_other = True
    if not question:
        return {"ok": False, "error": "question is required"}
    if not isinstance(options, list) or len(options) < 2:
        return {"ok": False, "error": "options must be a list of 2+ strings"}
    if len(options) > 6:
        return {"ok": False, "error": "too many options — keep it to 6 or fewer"}
    options = [str(o).strip() for o in options if str(o).strip()]
    cb = _ask_user_callback
    if cb is None:
        return {"ok": False, "error":
                "ask_user has no active interactive surface (running headless or batch). "
                "Pick the safest default and proceed instead of asking."}
    try:
        return cb(question, options, bool(allow_other))
    except Exception as e:
        return {"ok": False, "error": f"ask_user surface error: {type(e).__name__}: {e}"}


def _run_command(p: Dict) -> str:
    cmd = _rewrite_python_invocation(p["command"])
    cwd = p.get("cwd")
    # Destructive guardrail. Skipped when the user explicitly auto-approves
    # everything via env var, OR when this call has been threaded through
    # the per-call permission prompt and the user said yes (the bridges
    # set p["_approved"]=True before dispatching in that case — the user
    # has already SEEN the exact command and approved it, the second
    # refusal here was pure UX friction, not safety).
    _approved = bool(p.get("_approved"))
    if not _approved and os.environ.get("JARVIS_AUTO_APPROVE", "0") != "1":
        bad = _is_destructive(cmd)
        if bad:
            return (
                f"Error (run_command): refused to run a destructive command "
                f"WITHOUT explicit user approval.\n"
                f"  Detected pattern: '{bad}'\n"
                f"  Full command:     {cmd}\n\n"
                f"NEXT STEP: tell the user EXACTLY what this command would do "
                f"and what it would touch (path, target, etc.). Wait for the "
                f"user to type 'yes do it' or similar before retrying. Don't "
                f"retry with `--force`, don't try to bypass with another shell, "
                f"don't pipe to a different cmdlet. Just ask first.\n"
                f"If the user truly wants this command auto-run, they can set "
                f"JARVIS_AUTO_APPROVE=1."
            )
    if cwd:
        cwd = _resolve_read(cwd)
    else:
        cwd = WORKSPACE
    timeout = min(int(p.get("timeout") or 120), 300)
    shell_pref = (p.get("shell") or "").lower().strip()
    detached = bool(p.get("detached"))

    # Detached mode — for daemons / UIs / dev servers / launchers that
    # never exit on their own (Forge WebUI, npm run dev, ollama serve,
    # game launchers). Spawn the process in a new console and return
    # immediately. Without this, run_command blocks until the timeout.
    if detached:
        try:
            if sys.platform == "win32" and shell_pref != "cmd":
                # Detached launchers (game launchers, daemons, dev servers)
                # WANT their own console window since they print to stdout —
                # so we keep CREATE_NEW_CONSOLE here. But the PARENT shell
                # call should still be invisible.
                proc = subprocess.Popen(
                    ["powershell", "-NoProfile", "-NonInteractive", "-Command", cmd],
                    cwd=cwd,
                    creationflags=getattr(subprocess, "CREATE_NEW_CONSOLE", 0),
                )
            elif sys.platform == "win32":
                proc = subprocess.Popen(
                    cmd, cwd=cwd, shell=True,
                    creationflags=getattr(subprocess, "CREATE_NEW_CONSOLE", 0),
                )
            else:
                proc = subprocess.Popen(shlex.split(cmd), cwd=cwd)
        except Exception as e:
            return f"Error launching detached: {e}"
        return (
            f"$ {cmd}\n[detached, pid={proc.pid}] launched in a new window. "
            f"It will keep running until you close that window or kill the pid. "
            f"Use this for daemons/UIs only — for commands that print output "
            f"and exit, call without detached."
        )

    # Blocking/interactive guard (foreground only — detached above is exempt).
    # A `timeout /t 20` or `Start-Sleep`/`pause`/`Read-Host` would freeze the
    # whole session with no useful output. Refuse and point at the right tool.
    _block = _is_blocking_command(cmd)
    if _block:
        return (
            f"Error (run_command): refused '{_block}' — it would BLOCK the session "
            f"with no useful output (a sleep/timeout/keypress-wait). Never run these. "
            f"If the user wants a DELAY or to be reminded later, use set_reminder "
            f"(it fires in the background — you do NOT wait for it). If you were "
            f"trying to pause before another step, just do the next step directly."
        )

    if _is_whole_drive_scan(cmd):
        return (
            "Error (run_command): refused a whole-drive recursive scan — it can walk "
            "hundreds of thousands of files and hang for minutes. To FIND a file or "
            "folder by name, use the find_file or locate_path tool (built for this, "
            "far faster). If you really need a recursive listing, scope it to a "
            "specific folder (not a drive root), e.g. a Documents/Downloads subpath."
        )

    try:
        if sys.platform == "win32" and shell_pref != "cmd":
            # Prefer PowerShell on Windows by default — `dir`, `where`,
            # `tasklist` etc. all work as aliases, AND `Get-ChildItem`,
            # `Get-Process`, etc. work natively. Lets the model use the
            # full Windows toolbox without hitting a wall.
            r = subprocess.run(
                ["powershell", "-NoProfile", "-NonInteractive", "-Command", cmd],
                cwd=cwd, capture_output=True, timeout=timeout,
                creationflags=_NO_WINDOW,
            )
        elif sys.platform == "win32":
            # explicit shell=cmd path
            r = subprocess.run(cmd, cwd=cwd, shell=True, capture_output=True,
                               timeout=timeout, creationflags=_NO_WINDOW)
        else:
            # POSIX (Linux/macOS). shlex.split + no shell means pipes,
            # redirects, &&, ;, globs, command-substitution etc. silently
            # fail (the operators get passed as literal argv to the first
            # program). Windows routes through PowerShell which handles all
            # of that — match the behavior here by detecting shell
            # metacharacters and routing through a real shell when present.
            _shell_meta = ("|", ">", "<", "&", ";", "$(", "`", "*", "?",
                           "(", ")", "{", "}", "~", "\n")
            if any(tok in cmd for tok in _shell_meta):
                _sh = shutil.which("bash") or "/bin/sh"
                r = subprocess.run([_sh, "-c", cmd], cwd=cwd,
                                   capture_output=True, timeout=timeout)
            else:
                r = subprocess.run(shlex.split(cmd), cwd=cwd,
                                   capture_output=True, timeout=timeout)
    except subprocess.TimeoutExpired:
        return f"Error: '{cmd}' timed out after {timeout}s."
    except Exception as e:
        return f"Error: {e}"
    # Decode with errors='replace' — never let a stray byte (0xDB and friends
    # from raw file dumps or non-UTF-8 console code-page output) raise inside
    # subprocess's reader thread. Crashes there came back as empty tool results
    # and the model spam-retried, eating iterations.
    stdout = (r.stdout or b"").decode("utf-8", errors="replace")
    stderr = (r.stderr or b"").decode("utf-8", errors="replace")
    out = stdout + (("\n[stderr]\n" + stderr) if stderr else "")
    status_hint = _pip_status_hint(cmd, r.returncode, stdout, stderr)
    body = out.strip()
    if status_hint:
        body = f"{status_hint}\n{body}" if body else status_hint
    skill_nudge = _skill_crystallization_nudge(cmd)
    if skill_nudge:
        body = f"{skill_nudge}\n{body}" if body else skill_nudge
    return f"$ {cmd}\n[exit {r.returncode}]\n{body}"


def _skill_crystallization_nudge(cmd: str) -> str:
    """Detect when the model is hand-rolling a workflow a skill already
    covers, and append a one-line nudge to the tool result so it learns
    to reach for the right primitive next time. Examples:
      - `python -c "from reportlab..."` → "make-pdf already does this"
      - `python -c "from pptx..."`      → "make-pptx already does this"
      - `python -c "from openpyxl..."`  → "make-xlsx already does this"
    Silent when the model IS using the bundled skill scripts (build_pdf.py
    etc.), so no false-positive nudges."""
    low = (cmd or "").lower()
    if "skills/make-pdf/scripts" in low or "skills\\make-pdf\\scripts" in low:
        return ""
    if "skills/make-pptx/scripts" in low or "skills\\make-pptx\\scripts" in low:
        return ""
    if "skills/make-xlsx/scripts" in low or "skills\\make-xlsx\\scripts" in low:
        return ""
    if "reportlab" in low and "build_pdf.py" not in low:
        return ("[skill nudge] You're hand-rolling reportlab. The `make-pdf` "
                "skill already wraps this with style overrides + auto-open. "
                "Call load_skill('make-pdf') next time.")
    if "from pptx" in low and "build_pptx.py" not in low:
        return ("[skill nudge] You're hand-rolling python-pptx. The "
                "`make-pptx` skill already wraps this. Call "
                "load_skill('make-pptx') next time.")
    if "import openpyxl" in low and "build_xlsx.py" not in low:
        return ("[skill nudge] You're hand-rolling openpyxl. The `make-xlsx` "
                "skill already wraps this (with comma-vs-tab autodetect + "
                "freeze pane). Call load_skill('make-xlsx') next time.")
    return ""


def _pip_status_hint(cmd: str, exit_code: int, stdout: str, stderr: str) -> str:
    """If this looks like a pip install, prepend a one-line status the model
    can trust without parsing the full pip output. Solves the "model retries
    pip install 6 times because it can't tell 'Requirement already satisfied'
    is success" failure mode."""
    low = cmd.lower()
    if "pip install" not in low and "pip3 install" not in low:
        return ""
    combined = (stdout + "\n" + stderr).strip()
    if exit_code != 0:
        return "[pip install FAILED — read error below before retrying]"
    # Count lines per outcome
    sat_lines = sum(1 for ln in combined.splitlines()
                    if "requirement already satisfied" in ln.lower())
    installed_match = re.search(r"Successfully installed\s+(.+)", combined)
    if installed_match and sat_lines == 0:
        return f"[pip install SUCCEEDED — newly installed: {installed_match.group(1).strip()}]"
    if installed_match and sat_lines > 0:
        return (f"[pip install SUCCEEDED — newly installed: "
                f"{installed_match.group(1).strip()}; "
                f"{sat_lines} package(s) were already present]")
    if sat_lines > 0:
        return ("[pip install SUCCEEDED — all requested packages were already "
                "installed. Don't retry; nothing to do.]")
    return "[pip install SUCCEEDED — exit 0]"


# ============================================================
# SYSTEM (KNOW MY PC)
# ============================================================

def _try_psutil():
    try:
        import psutil  # type: ignore
        return psutil
    except ImportError:
        return None


def _system_info(p: Dict) -> str:
    psu = _try_psutil()
    info = {
        "hostname": socket.gethostname(),
        "user": os.environ.get("USERNAME") or os.environ.get("USER") or "?",
        "os": f"{platform.system()} {platform.release()} ({platform.version()})",
        "arch": platform.machine(),
        "python": sys.version.split()[0],
        "cwd": os.getcwd(),
        "workspace": WORKSPACE,
    }
    if psu:
        vm = psu.virtual_memory()
        info["cpu_count"] = psu.cpu_count(logical=True)
        info["cpu_percent"] = f"{psu.cpu_percent(interval=0.2)}%"
        info["ram_total_gb"] = round(vm.total / (1024**3), 2)
        info["ram_used_gb"] = round(vm.used / (1024**3), 2)
        info["ram_percent"] = f"{vm.percent}%"
        try:
            boot = datetime.fromtimestamp(psu.boot_time())
            info["boot_time"] = boot.strftime("%Y-%m-%d %H:%M:%S")
            info["uptime_hours"] = round((datetime.now() - boot).total_seconds() / 3600, 1)
        except Exception:
            pass
        disks = []
        for part in psu.disk_partitions(all=False):
            try:
                u = psu.disk_usage(part.mountpoint)
                disks.append({
                    "mount": part.mountpoint,
                    "fs": part.fstype,
                    "total_gb": round(u.total / (1024**3), 1),
                    "used_pct": f"{u.percent}%",
                })
            except OSError:
                pass
        info["disks"] = disks
    else:
        info["note"] = "Install 'psutil' for CPU/RAM/disk/uptime details: pip install psutil"
    return json.dumps(info, indent=2)


def _list_processes(p: Dict) -> str:
    psu = _try_psutil()
    if not psu:
        return "Error: needs psutil. Run: pip install psutil"
    flt = (p.get("name_filter") or "").lower()
    limit = int(p.get("limit") or 20)
    rows = []
    for proc in psu.process_iter(["pid", "name", "memory_info", "cpu_percent"]):
        try:
            n = proc.info["name"] or ""
            if flt and flt not in n.lower():
                continue
            rss = (proc.info["memory_info"].rss if proc.info["memory_info"] else 0) / (1024 * 1024)
            rows.append((proc.info["pid"], n, rss, proc.info["cpu_percent"] or 0.0))
        except (psu.NoSuchProcess, psu.AccessDenied):
            continue
    rows.sort(key=lambda r: r[2], reverse=True)
    rows = rows[:limit]
    out = ["PID\tNAME\tRAM(MB)\tCPU%"]
    for pid, n, rss, cpu in rows:
        out.append(f"{pid}\t{n}\t{rss:.1f}\t{cpu:.1f}")
    return "\n".join(out)


def _network_info(p: Dict) -> str:
    info = {
        "hostname": socket.gethostname(),
    }
    try:
        info["local_ip"] = socket.gethostbyname(socket.gethostname())
    except Exception:
        info["local_ip"] = "?"
    psu = _try_psutil()
    if psu:
        adapters = {}
        for name, addrs in psu.net_if_addrs().items():
            adapters[name] = [
                {"family": str(a.family), "address": a.address, "netmask": a.netmask}
                for a in addrs
            ]
        info["adapters"] = adapters
        try:
            stats = psu.net_io_counters(pernic=False)
            info["io"] = {
                "bytes_sent_mb": round(stats.bytes_sent / (1024**2), 1),
                "bytes_recv_mb": round(stats.bytes_recv / (1024**2), 1),
            }
        except Exception:
            pass
    return json.dumps(info, indent=2, default=str)


def _get_battery(p: Dict) -> str:
    psu = _try_psutil()
    if not psu:
        return "Error: needs psutil."
    b = psu.sensors_battery()
    if b is None:
        return "No battery detected (desktop?)."
    return json.dumps({
        "percent": round(b.percent, 1),
        "plugged": b.power_plugged,
        "secs_left": b.secsleft if b.secsleft != psu.POWER_TIME_UNLIMITED else "unlimited",
    }, indent=2)


def _disk_usage_should_background(path: str, max_depth: int) -> bool:
    """A drive root scan ('C:\\') OR a whole-tree walk (max_depth<=0) on
    any non-tiny directory will take minutes-to-hours and freeze the agent
    loop while the user watches a spinner. Heuristic: auto-background
    those by default. The model can still inline-scan small subtrees."""
    norm = os.path.normpath(path).rstrip("\\/")
    # Drive root on Windows: "C:" / "D:" / etc.
    if sys.platform == "win32" and len(norm) <= 2 and norm[-1:] == ":":
        return True
    # Unix root or top-level home — same long-walk risk.
    if norm in ("/", os.path.expanduser("~")):
        return True
    # Explicit whole-tree request always backgrounds.
    if max_depth <= 0:
        return True
    return False


def _disk_usage_core(p: Dict) -> str:
    """The original synchronous scan. Called inline for small trees and
    inside a background job for big ones. Returns the formatted report
    string the model + user will read."""
    raw = p.get("path") or WORKSPACE
    base = _resolve_read(raw)
    if not os.path.isdir(base):
        return f"Error: not a directory: {base}"
    top_n = int(p.get("top_n") or 15)
    kind = (p.get("kind") or "both").lower()
    max_depth = int(p.get("max_depth") or 1)

    folder_sizes: Dict[str, int] = {}
    biggest_files: List[Tuple[int, str]] = []
    total = 0
    errors = 0

    def _add_file(full: str, size: int):
        nonlocal total
        total += size
        if kind in ("both", "files"):
            biggest_files.append((size, full))

    if max_depth <= 0:
        # whole-tree scan, sum into the root
        for root, dirs, files in os.walk(base):
            dirs[:] = [d for d in dirs if d not in EXCLUDE_DIRS]
            for fn in files:
                full = os.path.join(root, fn)
                try:
                    sz = os.path.getsize(full)
                except OSError:
                    errors += 1
                    continue
                _add_file(full, sz)
        if kind in ("both", "folders"):
            folder_sizes[base] = total
    else:
        # tally sizes per immediate subfolder of `base`
        try:
            top_entries = os.listdir(base)
        except OSError as e:
            return f"Error: cannot list {base}: {e}"
        for entry in top_entries:
            full = os.path.join(base, entry)
            try:
                if os.path.isfile(full):
                    sz = os.path.getsize(full)
                    _add_file(full, sz)
                    continue
                if not os.path.isdir(full):
                    continue
            except OSError:
                errors += 1
                continue
            # recurse INTO this subfolder
            subtotal = 0
            for root, dirs, files in os.walk(full):
                dirs[:] = [d for d in dirs if d not in EXCLUDE_DIRS]
                for fn in files:
                    fp = os.path.join(root, fn)
                    try:
                        sz = os.path.getsize(fp)
                    except OSError:
                        errors += 1
                        continue
                    subtotal += sz
                    if kind in ("both", "files"):
                        biggest_files.append((sz, fp))
            total += subtotal
            if kind in ("both", "folders"):
                folder_sizes[full] = subtotal

    def _fmt(n: int) -> str:
        for unit in ("B", "KB", "MB", "GB", "TB"):
            if n < 1024:
                return f"{n:6.1f} {unit}"
            n /= 1024
        return f"{n:6.1f} PB"

    out: List[str] = [f"path: {base}", f"total scanned: {_fmt(total)}", ""]
    if kind in ("both", "folders"):
        out.append(f"top {top_n} folders by size:")
        sorted_folders = sorted(folder_sizes.items(), key=lambda kv: kv[1], reverse=True)[:top_n]
        for path, sz in sorted_folders:
            out.append(f"  {_fmt(sz)}  {path}")
        out.append("")
    if kind in ("both", "files"):
        out.append(f"top {top_n} files by size:")
        biggest_files.sort(key=lambda x: x[0], reverse=True)
        for sz, path in biggest_files[:top_n]:
            out.append(f"  {_fmt(sz)}  {path}")
    if errors:
        out.append(f"\n({errors} entries skipped — permission/access errors)")
    return "\n".join(out)


def _disk_usage(p: Dict) -> str:
    """Public dispatcher. Auto-routes to a background job when the scan
    would take minutes (drive root / whole-tree). User explicitly setting
    `background: false` overrides and forces synchronous execution.

    Returns immediately for backgrounded scans with a job_id the model
    (and `/jobs` command) can use to poll for the result. The actual
    scan runs in a daemon thread so the agent loop stays responsive
    instead of blocking for minutes while os.walk chews through a tree."""
    raw = p.get("path") or WORKSPACE
    base = _resolve_read(raw)
    if not os.path.isdir(base):
        return f"Error: not a directory: {base}"
    max_depth = int(p.get("max_depth") or 1)
    bg_arg = p.get("background")
    if bg_arg is None:
        should_bg = _disk_usage_should_background(base, max_depth)
    else:
        should_bg = bool(bg_arg)
    if not should_bg:
        return _disk_usage_core(p)
    # Background path — daemon thread, jobs.py meta + result file.
    try:
        from . import jobs
        r = jobs.start_python_job(
            label=f"disk_usage({base})",
            fn=_disk_usage_core,
            args=p,
            description=f"scan {base} (this can take minutes on a drive root)",
        )
        if r.get("ok"):
            return (
                f"Backgrounded the scan because {base} is too large to wait on "
                f"inline. Job id: {r['job_id']}.\n"
                f"Keep chatting — when you (or I) call get_job_result with "
                f"this id, the report will be there. Live log at: "
                f"{r['log_path']}"
            )
        return f"Error: could not background scan: {r.get('error', 'unknown')}"
    except Exception as e:
        # If the job runner is broken, fall back to inline rather than
        # losing the call entirely. The user gets the long wait but at
        # least no error.
        return _disk_usage_core(p)


def _locate_path(p: Dict) -> str:
    """Smart locator — finds a path/app by name across drives + start menu
    + registry without recursive globbing. Same chain I'd use by hand."""
    query = (p.get("query") or "").strip()
    if not query:
        return "Error: missing query"
    q_lower = query.lower()
    limit = int(p.get("limit") or 15)
    results: List[Tuple[str, str]] = []  # (source_tag, path_or_label)
    seen: set = set()

    def _try_add(tag: str, path: str) -> None:
        if path in seen:
            return
        seen.add(path)
        results.append((tag, path))

    # 1) Top-level dirs of every mounted drive (A-Z, not just C-P).
    # Mirrors find_file's drive enumeration so SteamLibrary on G:, media on
    # F:, etc. all get covered.
    if sys.platform == "win32":
        drives = []
        for letter in "CDEFGHIJKLMNOPQRSTUVWXYZ":
            root = f"{letter}:\\"
            try:
                if os.path.isdir(root):
                    drives.append(root)
            except OSError:
                continue
    else:
        drives = ["/"]
    # Note: no early-break here — we want all sources to contribute candidates
    # even on big drives. The limit gets applied to the final ranked list.
    for root in drives:
        try:
            for entry in os.listdir(root):
                # Skip Windows system folders that show up at every drive root
                if entry in EXCLUDE_DIRS:
                    continue
                if q_lower in entry.lower():
                    _try_add(f"drive_root({root[:2]})", os.path.join(root, entry))
        except (OSError, PermissionError):
            continue

    # 2) Common parent dirs (one level deeper)
    common_parents = [
        os.path.expanduser("~\\Documents"),
        os.path.expanduser("~\\Downloads"),
        os.path.expanduser("~\\Desktop"),
        os.environ.get("ProgramFiles", "C:\\Program Files"),
        os.environ.get("ProgramFiles(x86)", "C:\\Program Files (x86)"),
        os.environ.get("LOCALAPPDATA", ""),
        os.environ.get("APPDATA", ""),
    ]
    for parent in common_parents:
        if not parent or not os.path.isdir(parent):
            continue
        try:
            for entry in os.listdir(parent):
                if q_lower in entry.lower():
                    _try_add("common_dir", os.path.join(parent, entry))
        except (OSError, PermissionError):
            continue

    # 3) Start Menu shortcuts
    if sys.platform == "win32":
        lnk = _find_start_menu_lnk(query)
        if lnk:
            _try_add("start_menu", lnk)

    # 4) Registry-installed apps
    if sys.platform == "win32":
        try:
            apps = _list_installed_apps({"name_filter": query, "limit": 8})
            if apps and not apps.startswith("Error") and apps != "(none)":
                for line in apps.splitlines():
                    line = line.strip()
                    if line:
                        _try_add("registry", line)
        except Exception:
            pass

    if not results:
        return f"(nothing matching {query!r} in drive roots / common dirs / start menu / registry)"

    out = [f"Found {len(results)} match(es) for {query!r}:"]
    for source, path in results[:limit]:
        out.append(f"  [{source}] {path}")
    return "\n".join(out)


def _list_installed_apps_posix(flt: str, limit: int) -> str:
    """POSIX equivalent of the Windows registry scan. Returns the same shape
    as the Windows path: one `name version  —  publisher` line per app
    (version/publisher are usually unknown off-Windows, left blank)."""
    seen = set()
    apps: List[tuple] = []

    def _add(name: str, ver: str = "", pub: str = ""):
        name = (name or "").strip()
        if not name:
            return
        key = (name.lower(), ver)
        if key in seen:
            return
        seen.add(key)
        if flt and flt not in name.lower():
            return
        apps.append((name, ver, pub))

    try:
        if sys.platform == "darwin":
            # macOS: *.app bundles live in /Applications + ~/Applications
            for root in ("/Applications", os.path.expanduser("~/Applications"),
                         "/System/Applications"):
                if not os.path.isdir(root):
                    continue
                try:
                    for entry in os.listdir(root):
                        if entry.lower().endswith(".app"):
                            _add(entry[:-4])
                except OSError:
                    continue
        else:
            # Linux: parse Name= from *.desktop files in the XDG dirs.
            desk_dirs = [
                "/usr/share/applications",
                "/usr/local/share/applications",
                os.path.expanduser("~/.local/share/applications"),
                "/var/lib/flatpak/exports/share/applications",
                os.path.expanduser("~/.local/share/flatpak/exports/share/applications"),
                "/var/lib/snapd/desktop/applications",
            ]
            for root in desk_dirs:
                if not os.path.isdir(root):
                    continue
                try:
                    files = os.listdir(root)
                except OSError:
                    continue
                for fn in files:
                    if not fn.endswith(".desktop"):
                        continue
                    name = ""
                    try:
                        with open(os.path.join(root, fn), "r",
                                  encoding="utf-8", errors="replace") as fh:
                            for ln in fh:
                                # First top-level Name= line (skip Name[xx]=)
                                if ln.startswith("Name=") and not name:
                                    name = ln[5:].strip()
                                    break
                    except OSError:
                        continue
                    _add(name or os.path.splitext(fn)[0])
            # Also fold in flatpak's own list when available (catches apps
            # whose .desktop export dirs aren't standard).
            fp = shutil.which("flatpak")
            if fp:
                try:
                    r = subprocess.run(
                        [fp, "list", "--app", "--columns=name"],
                        capture_output=True, text=True, timeout=8)
                    for ln in (r.stdout or "").splitlines():
                        _add(ln.strip())
                except Exception:
                    pass
    except Exception as e:
        return f"Error: {e}"

    apps.sort(key=lambda a: a[0].lower())
    apps = apps[:limit]
    out = [f"{n} {v}  —  {pub}" for n, v, pub in apps]
    return "\n".join(out) if out else "(none)"


def _list_installed_apps(p: Dict) -> str:
    flt = (p.get("name_filter") or "").lower()
    limit = int(p.get("limit") or 50)
    if sys.platform != "win32":
        return _list_installed_apps_posix(flt, limit)
    try:
        import winreg
    except ImportError:
        return "Error: winreg unavailable."

    keys = [
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall"),
        (winreg.HKEY_LOCAL_MACHINE, r"SOFTWARE\WOW6432Node\Microsoft\Windows\CurrentVersion\Uninstall"),
        (winreg.HKEY_CURRENT_USER, r"SOFTWARE\Microsoft\Windows\CurrentVersion\Uninstall"),
    ]
    seen = set()
    apps = []
    for hive, sub in keys:
        try:
            with winreg.OpenKey(hive, sub) as k:
                i = 0
                while True:
                    try:
                        skn = winreg.EnumKey(k, i)
                    except OSError:
                        break
                    i += 1
                    try:
                        with winreg.OpenKey(k, skn) as sk:
                            try:
                                name = winreg.QueryValueEx(sk, "DisplayName")[0]
                            except FileNotFoundError:
                                continue
                            try:
                                ver = winreg.QueryValueEx(sk, "DisplayVersion")[0]
                            except FileNotFoundError:
                                ver = ""
                            try:
                                pub = winreg.QueryValueEx(sk, "Publisher")[0]
                            except FileNotFoundError:
                                pub = ""
                            key = (name, ver)
                            if key in seen:
                                continue
                            seen.add(key)
                            if flt and flt not in name.lower():
                                continue
                            apps.append((name, ver, pub))
                    except OSError:
                        continue
        except OSError:
            continue
    apps.sort(key=lambda a: a[0].lower())
    apps = apps[:limit]
    out = [f"{n} {v}  —  {pub}" for n, v, pub in apps]
    return "\n".join(out) if out else "(none)"


# ============================================================
# APP CONTROL
# ============================================================

# Verbose-name → canonical-executable map. Intentionally excludes regedit
# and other sensitive launchers — those would need an explicit user request
# anyway. Sandbox covers FILE ops, not what GUI apps do once running.
_COMMON_APP_ALIASES = {
    "calculator": "calc",
    "calc": "calc",
    "notepad": "notepad",
    "paint": "mspaint",
    "ms paint": "mspaint",
    "explorer": "explorer",
    "file explorer": "explorer",
    "command prompt": "cmd",
    "cmd": "cmd",
    "powershell": "powershell",
    "terminal": "wt",  # Windows Terminal if installed
    "edge": "msedge",
    "browser": "msedge",
}


def _find_start_menu_lnk(query: str) -> Optional[str]:
    """Scan Start Menu + Desktop for a .lnk whose stem matches `query`
    (case-insensitive substring). Returns full path or None.

    Many users keep game launchers, app shortcuts etc. on the Desktop, not
    in Start Menu — checking only Start Menu means open_app fails on every
    desktop-shortcut game launcher."""
    if sys.platform != "win32":
        return None
    import time as _t
    deadline = _t.monotonic() + 2.0  # hard cap: never let app-launch hang
    q = query.lower().strip()
    fuzzy: Optional[str] = None

    def _scan_file(dirpath: str, fn: str):
        nonlocal fuzzy
        if not fn.lower().endswith(".lnk"):
            return None
        low = os.path.splitext(fn)[0].lower()
        if low == q:
            return os.path.join(dirpath, fn)
        if fuzzy is None and q in low:
            fuzzy = os.path.join(dirpath, fn)
        return None

    # Start Menu: shallow tree of shortcuts - safe to walk recursively.
    start_menu = [
        os.path.join(os.environ.get("APPDATA", ""), r"Microsoft\Windows\Start Menu\Programs"),
        os.path.join(os.environ.get("PROGRAMDATA", ""), r"Microsoft\Windows\Start Menu\Programs"),
    ]
    for root in start_menu:
        if not root or not os.path.isdir(root):
            continue
        for dirpath, _, filenames in os.walk(root):
            if _t.monotonic() > deadline:
                return fuzzy
            for fn in filenames:
                hit = _scan_file(dirpath, fn)
                if hit:
                    return hit
    # Desktops: TOP LEVEL ONLY. A recursive os.walk here was the 29s hang -
    # a Desktop with a big folder (downloads dump, a game dir) walks forever.
    # Launcher shortcuts always sit at the Desktop root anyway.
    desktops = [
        os.path.join(os.environ.get("USERPROFILE", ""), "Desktop"),
        os.path.join(os.environ.get("PUBLIC", ""), "Desktop"),
        os.path.join(os.environ.get("USERPROFILE", ""), "OneDrive", "Desktop"),
        os.path.join(os.environ.get("ONEDRIVE", ""), "Desktop"),
    ]
    for root in desktops:
        if not root or not os.path.isdir(root):
            continue
        try:
            for fn in os.listdir(root):
                hit = _scan_file(root, fn)
                if hit:
                    return hit
        except OSError:
            continue
    return fuzzy


def _windows_uwp_uri(canonical: str) -> Optional[str]:
    """Some Windows apps live as UWP and respond to URI schemes. Returns the
    URI to feed `start` / os.startfile, or None."""
    uwp = {
        "calc": "calculator://",
        "calculator": "calculator://",
        "ms-settings": "ms-settings:",
        "settings": "ms-settings:",
    }
    return uwp.get(canonical.lower())


def _open_app(p: Dict) -> str:
    name = (p.get("name") or "").strip()
    if not name:
        return "Error: missing app name"
    # Strip surrounding quotes — model sometimes wraps paths in literal quotes
    # (e.g. `"G:\Assassins Creed II\Game.exe"`), which then never resolves.
    if (name.startswith('"') and name.endswith('"')) or (name.startswith("'") and name.endswith("'")):
        name = name[1:-1].strip()
    # Drop common UI suffixes that don't exist as executable names. Users say
    # "Brave Browser" but the binary is just `brave.exe` and the Start Menu
    # shortcut is "Brave". Same for "Google Chrome", "Discord Inc", etc.
    bare = re.sub(r"\s+(browser|browsers|app|application|inc|launcher|client)$",
                  "", name, flags=re.IGNORECASE).strip()
    if bare and bare.lower() != name.lower():
        # Try the bare name first, fall through to original if not found.
        # We mutate canonical lookup below to use bare; original name remains
        # available for the error message.
        canonical_attempt_first = bare
    else:
        canonical_attempt_first = name

    # Map verbose names → canonical executables
    canonical = _COMMON_APP_ALIASES.get(canonical_attempt_first.lower(), canonical_attempt_first)

    if sys.platform == "win32":
        # 1. Real file path / folder path — opens with default association
        # (videos in default player, .rar in archive viewer, folders in Explorer, etc.)
        # Check this BEFORE PATH lookup so 'G:\foo.rar' isn't mistaken for an exe.
        cand_path = os.path.expanduser(canonical)
        if os.path.exists(cand_path):
            os.startfile(cand_path)  # type: ignore[attr-defined]
            return f"opened: {cand_path}"

        # 2. URL-shaped → open in default browser via os.startfile
        if canonical.startswith(("http://", "https://", "file://")):
            os.startfile(canonical)  # type: ignore[attr-defined]
            return f"opened url: {canonical}"

        # 3. Direct PATH lookup (calc, notepad, mspaint, explorer, etc.)
        for cand in (canonical, canonical + ".exe"):
            full = shutil.which(cand)
            if full:
                subprocess.Popen([full], shell=False, close_fds=True,
                                 creationflags=getattr(subprocess, "DETACHED_PROCESS", 0))
                return f"launched: {full}"

        # 4. UWP URI scheme (Calculator on Win10/11 lives as a UWP app)
        uri = _windows_uwp_uri(canonical)
        if uri:
            try:
                os.startfile(uri)  # type: ignore[attr-defined]
                return f"launched UWP: {uri}"
            except OSError:
                pass

        # 5. Start Menu shortcut search
        lnk = _find_start_menu_lnk(canonical)
        if lnk:
            os.startfile(lnk)  # type: ignore[attr-defined]
            return f"launched: {os.path.basename(lnk)[:-4]} (via Start Menu shortcut)"

        # 5. Failure — surface a real error and suggest matches
        suggestions: List[str] = []
        try:
            installed = _list_installed_apps({"name_filter": canonical, "limit": 5})
            if installed and not installed.startswith("Error") and installed != "(none)":
                suggestions = [ln.split("  —")[0].strip() for ln in installed.splitlines() if ln.strip()]
        except Exception:
            pass
        msg = (
            f"Error: could not find '{name}' on this system. "
            f"Tried PATH, UWP URI, file path, and Start Menu shortcuts. "
            f"NEXT STEP: if '{name}' is a known web service (Discord, Spotify, "
            f"WhatsApp, Telegram, YouTube, Twitter, Reddit, Gmail, Notion, "
            f"GitHub, Slack, ChatGPT, Claude, etc.), call open_in_browser with "
            f"the canonical URL (e.g. https://{name.lower()}.com) — that's the "
            f"right fallback. Do NOT ask the user where it's installed."
        )
        if suggestions:
            msg += "\nClosest installed matches (use list_installed_apps for more):\n  " + "\n  ".join(suggestions)
        else:
            msg += (
                "\nTry an .exe name like 'calc', 'notepad', 'chrome', or "
                "call list_installed_apps to find the registered name."
            )
        return msg

    # macOS
    if sys.platform == "darwin":
        cand_path = os.path.expanduser(canonical)
        # File / folder / URL → `open` figures out the default handler.
        if os.path.exists(cand_path):
            try:
                subprocess.Popen(["open", cand_path])
                return f"opened: {cand_path}"
            except Exception as e:
                return f"Error: open failed: {e}"
        if canonical.startswith(("http://", "https://", "file://")):
            try:
                subprocess.Popen(["open", canonical])
                return f"opened url: {canonical}"
            except Exception as e:
                return f"Error: open failed: {e}"
        # Named app → `open -a <name>` resolves the .app bundle by name.
        try:
            r = subprocess.run(["open", "-a", canonical],
                               capture_output=True, text=True, timeout=10)
            if r.returncode == 0:
                return f"launched: {canonical}"
        except Exception:
            pass
        # Last resort: a CLI binary on PATH.
        full = shutil.which(canonical)
        if full:
            try:
                subprocess.Popen([full])
                return f"launched: {full}"
            except Exception as e:
                return f"Error: {e}"
        return (f"Error: could not find or launch '{name}' on macOS. "
                f"Tried `open <path>`, `open -a <name>`, and PATH. "
                f"If it's a web service, call open_in_browser instead.")

    # Linux (and other POSIX)
    cand_path = os.path.expanduser(canonical)
    # Files / folders / URLs → xdg-open uses the desktop default handler.
    if os.path.exists(cand_path) or canonical.startswith(
            ("http://", "https://", "file://")):
        target = cand_path if os.path.exists(cand_path) else canonical
        opener = shutil.which("xdg-open")
        if opener:
            try:
                subprocess.Popen([opener, target],
                                 stdout=subprocess.DEVNULL,
                                 stderr=subprocess.DEVNULL)
                return f"opened: {target}"
            except Exception as e:
                return f"Error: xdg-open failed: {e}"
        # No xdg-open and it's a real file we can't hand off — fall through.
    # Named app → try the binary on PATH and detach it.
    full = shutil.which(canonical)
    if full:
        try:
            subprocess.Popen([full], stdout=subprocess.DEVNULL,
                             stderr=subprocess.DEVNULL, start_new_session=True)
            return f"launched: {full}"
        except Exception as e:
            return f"Error: {e}"
    return (f"Error: could not find or launch '{name}' on this system. "
            f"Tried xdg-open and PATH. If it's a web service, call "
            f"open_in_browser with the canonical URL instead.")


# --------------------------------------------------------------------------
# Browser detection + URL helpers
# --------------------------------------------------------------------------

_BROWSER_CANDIDATES = {
    "chrome": [
        r"%ProgramFiles%\Google\Chrome\Application\chrome.exe",
        r"%ProgramFiles(x86)%\Google\Chrome\Application\chrome.exe",
        r"%LocalAppData%\Google\Chrome\Application\chrome.exe",
    ],
    "brave": [
        r"%ProgramFiles%\BraveSoftware\Brave-Browser\Application\brave.exe",
        r"%ProgramFiles(x86)%\BraveSoftware\Brave-Browser\Application\brave.exe",
        r"%LocalAppData%\BraveSoftware\Brave-Browser\Application\brave.exe",
    ],
    "edge": [
        r"%ProgramFiles%\Microsoft\Edge\Application\msedge.exe",
        r"%ProgramFiles(x86)%\Microsoft\Edge\Application\msedge.exe",
    ],
    "firefox": [
        r"%ProgramFiles%\Mozilla Firefox\firefox.exe",
        r"%ProgramFiles(x86)%\Mozilla Firefox\firefox.exe",
    ],
    "opera": [
        r"%LocalAppData%\Programs\Opera\opera.exe",
        r"%ProgramFiles%\Opera\opera.exe",
    ],
    "vivaldi": [
        r"%LocalAppData%\Vivaldi\Application\vivaldi.exe",
        r"%ProgramFiles%\Vivaldi\Application\vivaldi.exe",
    ],
}


def _resolve_browser(name: str) -> Optional[str]:
    """Return the absolute path to a browser executable, or None."""
    name = (name or "").lower().strip()
    if name not in _BROWSER_CANDIDATES:
        # Try as a raw name on PATH
        full = shutil.which(name) or shutil.which(name + ".exe")
        return full
    for raw in _BROWSER_CANDIDATES[name]:
        path = os.path.expandvars(raw)
        if os.path.isfile(path):
            return path
    # Last try: PATH
    for stem in (name, name + ".exe"):
        full = shutil.which(stem)
        if full:
            return full
    return None


def _list_browsers(p: Dict) -> str:
    found: List[Tuple[str, str]] = []
    for name in _BROWSER_CANDIDATES:
        path = _resolve_browser(name)
        if path:
            found.append((name, path))
    if not found:
        return "(no browsers detected — try opening a URL with the system default via open_url)"
    return "\n".join(f"  {name:<10} {path}" for name, path in found)


def _validate_url(p: Dict) -> str:
    raw_url = (p.get("url") or "").strip()
    if not raw_url:
        return "Error: missing url"
    if not raw_url.startswith(("http://", "https://")):
        raw_url = "https://" + raw_url
    import time as _time
    t0 = _time.time()
    try:
        # Try HEAD first (cheap), fall back to GET if server doesn't allow it
        req = urllib.request.Request(raw_url, method="HEAD", headers={
            "User-Agent": "Mozilla/5.0 (Jarvis-Local)",
        })
        try:
            r = urllib.request.urlopen(req, timeout=10)
        except urllib.error.HTTPError as he:
            if he.code in (405, 400):
                req = urllib.request.Request(raw_url, headers={
                    "User-Agent": "Mozilla/5.0 (Jarvis-Local)",
                })
                r = urllib.request.urlopen(req, timeout=10)
            else:
                raise
        ms = int((_time.time() - t0) * 1000)
        info = {
            "url": raw_url,
            "status": r.status,
            "ok": 200 <= r.status < 400,
            "final_url": r.geturl(),
            "content_type": r.headers.get("Content-Type", ""),
            "response_ms": ms,
        }
        return json.dumps(info, indent=2)
    except urllib.error.HTTPError as e:
        return json.dumps({"url": raw_url, "status": e.code, "ok": False,
                           "error": str(e)}, indent=2)
    except Exception as e:
        return json.dumps({"url": raw_url, "ok": False,
                           "error": f"{type(e).__name__}: {e}"}, indent=2)


def _resolve_chromium_profile(exe_path: str, name_or_dir: str) -> str:
    """Map a friendly profile name ('personal', 'work') to its on-disk
    directory ('Profile 5', 'Default') by reading the browser's Local State
    JSON. If name_or_dir already names an existing profile dir, returns it
    unchanged. Returns name_or_dir as-is when resolution fails — caller can
    still try passing it raw to Chrome and accept the silent fallback.

    Why this exists: Chrome/Brave/Edge store profiles as `Profile 1`, `Profile
    2`, etc. on disk, but show them to the user under friendly names ("Work",
    "Personal", a Gmail address). Users say "open in my personal profile" —
    that string never matches a dir name, so `--profile-directory=personal`
    silently does nothing. This maps it correctly."""
    if not name_or_dir:
        return name_or_dir

    exe_lower = exe_path.lower()
    appdata = os.environ.get("LOCALAPPDATA", "")
    user_data: Optional[str] = None
    if "brave" in exe_lower:
        user_data = os.path.join(appdata, "BraveSoftware", "Brave-Browser", "User Data")
    elif "chrome" in exe_lower:
        user_data = os.path.join(appdata, "Google", "Chrome", "User Data")
    elif "edge" in exe_lower or "msedge" in exe_lower:
        user_data = os.path.join(appdata, "Microsoft", "Edge", "User Data")
    elif "vivaldi" in exe_lower:
        user_data = os.path.join(appdata, "Vivaldi", "User Data")
    if not user_data or not os.path.isdir(user_data):
        return name_or_dir

    # Priority: display-name match wins over literal dir name. Most users
    # say "open my work profile" expecting the DISPLAY name from Chrome's
    # profile picker; literal dir names like "Profile 5" are rarely typed
    # but can happen, so we accept those as a second pass.
    local_state = os.path.join(user_data, "Local State")
    target = name_or_dir.lower()
    if os.path.isfile(local_state):
        try:
            with open(local_state, "r", encoding="utf-8") as f:
                data = json.load(f)
            info_cache = data.get("profile", {}).get("info_cache", {})
            # 1. Exact display-name match (case-insensitive)
            for dir_name, info in info_cache.items():
                if (info.get("name") or "").lower() == target:
                    return dir_name
            # 2. Substring display-name match
            for dir_name, info in info_cache.items():
                if target in (info.get("name") or "").lower():
                    return dir_name
        except (OSError, json.JSONDecodeError):
            pass

    # 3. Fallback: maybe they typed the literal dir name (Default, Profile 1)
    if os.path.isdir(os.path.join(user_data, name_or_dir)):
        return name_or_dir

    return name_or_dir


def _list_chromium_profiles(exe_path: str) -> str:
    """Return a pretty list of `Display name -> dir` for the browser at exe_path.
    Used in error messages so the model can tell the user what's available."""
    appdata = os.environ.get("LOCALAPPDATA", "")
    exe_lower = exe_path.lower()
    user_data = None
    if "brave" in exe_lower:
        user_data = os.path.join(appdata, "BraveSoftware", "Brave-Browser", "User Data")
    elif "chrome" in exe_lower:
        user_data = os.path.join(appdata, "Google", "Chrome", "User Data")
    elif "edge" in exe_lower or "msedge" in exe_lower:
        user_data = os.path.join(appdata, "Microsoft", "Edge", "User Data")
    if not user_data or not os.path.isdir(user_data):
        return ""
    local_state = os.path.join(user_data, "Local State")
    if not os.path.isfile(local_state):
        return ""
    try:
        with open(local_state, "r", encoding="utf-8") as f:
            data = json.load(f)
        info_cache = data.get("profile", {}).get("info_cache", {})
        lines = []
        for dir_name, info in sorted(info_cache.items()):
            display = info.get("name") or "(unnamed)"
            lines.append(f"  {display!r}  ->  {dir_name}")
        return "\n".join(lines) if lines else ""
    except (OSError, json.JSONDecodeError):
        return ""


def _open_in_browser(p: Dict) -> str:
    url = (p.get("url") or "").strip()
    if not url:
        return "Error: missing url"
    if not url.startswith(("http://", "https://", "file://")):
        url = "https://" + url
    browser = (p.get("browser") or "").strip()
    profile = (p.get("profile") or "").strip()

    if not browser:
        # No browser specified — fall through to system default
        try:
            os.startfile(url)  # type: ignore[attr-defined]
            return f"opened in default browser: {url}"
        except Exception as e:
            return f"Error: could not open {url}: {e}"

    exe = _resolve_browser(browser)
    if not exe:
        installed = _list_browsers({})
        return (
            f"Error: '{browser}' not found on this machine.\n"
            f"Installed browsers:\n{installed}"
        )

    args: List[str] = [exe]
    if profile:
        # Chromium-family (chrome, brave, edge, vivaldi, opera) use --profile-directory
        # Firefox uses -P / --new-window with profile-name
        is_firefox = "firefox" in os.path.basename(exe).lower()
        if is_firefox:
            args += ["-P", profile]
        else:
            # User said "personal" — that's the display name. Resolve to the
            # actual on-disk dir like "Profile 5". Falls back to raw string if
            # we can't resolve, which is the old broken behavior.
            resolved = _resolve_chromium_profile(exe, profile)
            args += [f"--profile-directory={resolved}"]
            # If we couldn't resolve a friendly name to a real dir, tell the
            # user (helps the model recover when they typed a name that
            # doesn't exist).
            if resolved == profile and not os.path.isdir(
                os.path.join(os.environ.get("LOCALAPPDATA", ""), profile)
            ):
                available = _list_chromium_profiles(exe)
                if available and profile.lower() not in available.lower():
                    # We're firing off the launch anyway (Chrome will silently
                    # fall back to Default), but warn so model + user know.
                    args.append(url)
                    try:
                        subprocess.Popen(args, shell=False, close_fds=True,
                                         creationflags=getattr(subprocess, "DETACHED_PROCESS", 0))
                    except Exception as e:
                        return f"Error launching {browser}: {e}"
                    return (
                        f"opened in {browser} (profile '{profile}' not found — "
                        f"opened in default instead).\nAvailable profiles:\n{available}"
                    )
    args.append(url)
    try:
        subprocess.Popen(args, shell=False, close_fds=True,
                         creationflags=getattr(subprocess, "DETACHED_PROCESS", 0))
    except Exception as e:
        return f"Error launching {browser}: {e}"
    pretty = browser + (f" / {profile}" if profile else "")
    return f"opened in {pretty}: {url}"


def _open_url(p: Dict) -> str:
    import webbrowser
    url = (p.get("url") or "").strip()
    if not url:
        return "Error: missing url"
    if not url.startswith(("http://", "https://", "file://")):
        url = "https://" + url
    webbrowser.open(url)
    return f"opened: {url}"


_IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".bmp"}


def _view_image(p: Dict) -> str:
    """Read an image and return a marker the runtime understands.

    - In LM Studio (via MCP), the mcp_server wraps this into an Image
      content block so LM Studio shows the image to the model — vision
      kicks in automatically on the next turn.
    - In the CLI, the runtime auto-attaches the most recent image when
      the user references it by name; this tool is a fallback.
    """
    raw = (p.get("path") or "").strip().strip('"').strip("'")
    if not raw:
        return "Error: missing path"
    path = os.path.expanduser(raw)
    if not os.path.isabs(path):
        path = os.path.abspath(path)
    if not os.path.isfile(path):
        return f"Error: not a file: {path}"
    ext = os.path.splitext(path)[1].lower()
    if ext not in _IMAGE_EXTS:
        return f"Error: not a recognized image extension ({ext}). Supported: {sorted(_IMAGE_EXTS)}"
    size = os.path.getsize(path)
    if size > 20 * 1024 * 1024:
        return f"Error: image too large ({size // 1024 // 1024}MB, max 20MB)"
    # Return a structured marker. mcp_server.py detects this prefix and
    # converts the call to a real Image content. The CLI just shows the
    # text — its auto-attach handles vision separately.
    return f"__JARVIS_IMAGE__ {path} ({size} bytes, {ext[1:]})"


def _screenshot_posix_cli(out: str) -> Optional[str]:
    """Capture the screen by shelling out to a native tool. Used as a Linux
    fallback when PIL.ImageGrab isn't available (it has no X11/Wayland
    backend on Linux). Tries Wayland (grim) then X11 (scrot,
    gnome-screenshot, spectacle, maim). Returns the path on success, None if
    no tool worked. Each tool writes directly to `out`."""
    attempts = [
        ("grim", [out]),                       # Wayland
        ("scrot", [out]),                      # X11
        ("gnome-screenshot", ["-f", out]),     # GNOME (X11/Wayland)
        ("spectacle", ["-b", "-n", "-o", out]),  # KDE
        ("maim", [out]),                       # X11
    ]
    for tool, args in attempts:
        exe = shutil.which(tool)
        if not exe:
            continue
        try:
            r = subprocess.run([exe] + args, capture_output=True, timeout=20)
            if r.returncode == 0 and os.path.exists(out) and os.path.getsize(out) > 0:
                return out
        except Exception:
            continue
    return None


def _screenshot(p: Dict) -> str:
    # Optional delay so the user can switch to the window they want
    # captured before the shutter fires. Without this, the screenshot
    # almost always captures Hearth itself (the chat is in focus when
    # they typed the prompt). Capped at 10s so the model can't stall
    # the agent loop.
    try:
        delay = max(0.0, min(10.0, float(p.get("delay_s") or 0)))
    except (TypeError, ValueError):
        delay = 0.0
    # Show an on-screen cue (purple dots + "capturing screen") so the user knows
    # the shutter is coming — especially during a delay while they switch windows.
    # It auto-removes before the grab, so it never appears in the shot. Best-effort.
    import time as _t
    cue = delay if delay > 0 else 0.5  # always give at least a brief flash
    try:
        from . import capture_overlay
        capture_overlay.flash(cue)
    except Exception:
        pass
    _t.sleep(cue)         # wait out the delay (or the brief cue when delay==0)
    # The cue fades to fully transparent by the end of `cue`; this extra settle
    # guarantees it's off-screen + destroyed before we grab, so it's never in
    # the shot.
    _t.sleep(0.35)
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out = os.path.join(SHOTS_DIR, f"shot_{ts}.png")
    # Primary path: PIL.ImageGrab (works on Windows + macOS; on Linux it has
    # no backend and raises). On any failure, fall back to a native CLI
    # capture tool (grim/scrot/gnome-screenshot/...) on Linux.
    img = None
    try:
        from PIL import ImageGrab  # type: ignore
        img = ImageGrab.grab()
        img.save(out)
    except Exception as e:
        if sys.platform == "win32":
            if isinstance(e, ImportError):
                return "Error: needs Pillow. Run: pip install pillow"
            return f"Error: {e}"
        captured = _screenshot_posix_cli(out)
        if not captured:
            return ("Error: screenshot failed. Install Pillow "
                    "(pip install pillow) or a native tool: grim (Wayland), "
                    "scrot / gnome-screenshot / spectacle / maim (X11).")
        img = None  # size unknown from CLI capture
    delay_note = (f" (waited {delay:.0f}s before capture)" if delay > 0 else "")
    size_note = f" ({img.size[0]}x{img.size[1]})" if img is not None else ""
    return (
        f"Saved: {out}{size_note}{delay_note}\n"
        f"NEXT STEP: if the user asked you to DESCRIBE / TELL THEM WHAT'S "
        f"ON the screen, call view_image with path='{out}' RIGHT NOW in "
        f"this same turn — don't stop here. If they asked you to SHOW or "
        f"OPEN the screenshot in their gallery, call open_app with that "
        f"path instead. Only stop if they only asked you to CAPTURE."
    )


def _clipboard_posix_read_cmd() -> Optional[list]:
    """Pick the right clipboard-read command for this POSIX desktop.
    macOS uses pbpaste; Linux prefers Wayland (wl-paste) then X11
    (xclip / xsel). Returns argv list or None if no tool is installed."""
    if sys.platform == "darwin":
        if shutil.which("pbpaste"):
            return ["pbpaste"]
        return None
    # Linux / other
    if shutil.which("wl-paste"):
        return ["wl-paste", "--no-newline"]
    if shutil.which("xclip"):
        return ["xclip", "-selection", "clipboard", "-o"]
    if shutil.which("xsel"):
        return ["xsel", "--clipboard", "--output"]
    return None


def _clipboard_posix_write_cmd() -> Optional[list]:
    """Pick the right clipboard-write command for this POSIX desktop."""
    if sys.platform == "darwin":
        if shutil.which("pbcopy"):
            return ["pbcopy"]
        return None
    if shutil.which("wl-copy"):
        return ["wl-copy"]
    if shutil.which("xclip"):
        return ["xclip", "-selection", "clipboard"]
    if shutil.which("xsel"):
        return ["xsel", "--clipboard", "--input"]
    return None


def _clipboard_read(p: Dict) -> str:
    if sys.platform == "win32":
        # Text first; then a copied IMAGE (save to a temp PNG, hand back a path
        # view_image can read — enables "I copied a screenshot, what is it?");
        # then copied FILE(S) → return their paths. Empty if none.
        import tempfile as _tf, time as _tm
        img_path = os.path.join(_tf.gettempdir(),
                                f"hearth_clip_{int(_tm.time())}.png").replace("\\", "/")
        ps = (
            "$t = Get-Clipboard -Raw -ErrorAction SilentlyContinue; "
            "if ($t) { 'TEXT'; $t } else { "
            "$i = Get-Clipboard -Format Image -ErrorAction SilentlyContinue; "
            f"if ($i) {{ $i.Save('{img_path}'); 'IMAGE:{img_path}' }} else {{ "
            "$f = Get-Clipboard -Format FileDropList -ErrorAction SilentlyContinue; "
            "if ($f) { 'FILES'; $f } else { 'EMPTY' } } }"
        )
        try:
            r = subprocess.run(
                ["powershell", "-NoProfile", "-STA", "-Command", ps],
                capture_output=True, text=True, timeout=8,
                creationflags=_NO_WINDOW,
            )
            out = (r.stdout or "").strip()
            if out.startswith("TEXT"):
                return out[4:].strip() or "(empty)"
            if out.startswith("IMAGE:"):
                path = out[len("IMAGE:"):].strip().splitlines()[0].strip()
                return (f"[clipboard holds an IMAGE — saved to {path}. "
                        f"Call view_image on that path to see it.]")
            if out.startswith("FILES"):
                files = out[len("FILES"):].strip()
                return f"[clipboard holds file(s):]\n{files}" if files else "(empty)"
            return "(empty)"
        except Exception as e:
            return f"Error: {e}"
    # macOS / Linux: native clipboard CLI first (robust headless/Wayland).
    cmd = _clipboard_posix_read_cmd()
    if cmd:
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=5)
            return r.stdout.rstrip("\r\n") or "(empty)"
        except Exception as e:
            return f"Error: {e}"
    # Fallback: tkinter (brittle headless/Wayland, but works on plain X11).
    try:
        from tkinter import Tk
        root = Tk()
        root.withdraw()
        text = root.clipboard_get()
        root.destroy()
        return text or "(empty)"
    except Exception as e:
        return (f"Error: {e} (no clipboard tool found — install "
                f"wl-clipboard or xclip on Linux)")


def _clipboard_write(p: Dict) -> str:
    text = p.get("text", "")
    if sys.platform == "win32":
        try:
            r = subprocess.run(["clip"], input=text, text=True, timeout=5,
                               creationflags=_NO_WINDOW)
            return f"Copied ({len(text)} chars)" if r.returncode == 0 else "Error: clip failed"
        except Exception as e:
            return f"Error: {e}"
    # macOS / Linux: native clipboard CLI first.
    cmd = _clipboard_posix_write_cmd()
    if cmd:
        try:
            r = subprocess.run(cmd, input=text, text=True, timeout=5)
            return (f"Copied ({len(text)} chars)" if r.returncode == 0
                    else f"Error: {cmd[0]} failed")
        except Exception as e:
            return f"Error: {e}"
    # Fallback: tkinter.
    try:
        from tkinter import Tk
        root = Tk()
        root.withdraw()
        root.clipboard_clear()
        root.clipboard_append(text)
        root.update()
        root.destroy()
        return f"Copied ({len(text)} chars)"
    except Exception as e:
        return (f"Error: {e} (no clipboard tool found — install "
                f"wl-clipboard or xclip on Linux)")


# ============================================================
# MEMORY (per-fact files + index, in hearth/memory.py)
# ============================================================

def _memory_save(p: Dict) -> str:
    from . import memory  # local import to avoid circular at module load
    return memory.save(
        title=p["title"],
        mtype=p.get("type", "user"),
        description=p.get("description", ""),
        body=p.get("body", ""),
        tags=p.get("tags") or [],
        sub_category=p.get("sub_category"),
        force=bool(p.get("force")),
    )


def _memory_recall(p: Dict) -> str:
    from . import memory
    return memory.recall(p.get("query", ""), int(p.get("limit") or 5))


def _memory_list(p: Dict) -> str:
    from . import memory
    return memory.list_index()


def _memory_forget(p: Dict) -> str:
    from . import memory
    # Accept both `title` and `name` — the GUI memory-delete endpoint sends
    # `name` while the model tool-call uses `title`. Either works.
    key = (p.get("title") or p.get("name") or "").strip()
    if not key:
        return "Error: memory_forget needs 'title' or 'name'."
    return memory.forget(key)


# ============================================================
# TIME
# ============================================================

# --------------------------------------------------------------------------
# Forge WebUI orchestration (Stable Diffusion image generation).
# Point JARVIS_FORGE_DIR at your local Forge / SD-WebUI install (e.g.
# D:\AI\sd-webui-forge). We boot it with --api so /sdapi/v1/* is exposed.
# --------------------------------------------------------------------------

def _autodetect_forge_dir() -> str:
    """Scan common install locations for an SD/Forge WebUI install.

    Returns the first matching folder that contains a recognized launcher,
    or "" if none found. Order: settings.json saved value > env var >
    common candidates per OS.

    A folder qualifies if it contains webui.bat / webui.sh /
    webui-user.bat AND a 'modules' subdir (the WebUI signature)."""
    # Saved setting wins
    try:
        from pathlib import Path as _Path
        settings_path = _Path(SETTINGS_PATH)
        if settings_path.is_file():
            saved = json.loads(settings_path.read_text(encoding="utf-8"))
            saved_dir = (saved.get("forge_dir") or "").strip()
            if saved_dir:
                _ld = _forge_launch_dir(saved_dir)
                if _ld:
                    return _ld
    except Exception:
        pass
    # Env var next
    env_dir = (os.environ.get("JARVIS_FORGE_DIR") or "").strip()
    if env_dir:
        _ld = _forge_launch_dir(env_dir)
        if _ld:
            return _ld
    # Common install locations
    home = os.path.expanduser("~")
    candidates = []
    if sys.platform == "win32":
        for base in ("C:\\", "D:\\", "E:\\", "F:\\", "G:\\"):
            for sub in ("AI", "stable-diffusion", "SD", "tools"):
                candidates.append(os.path.join(base, sub, "sd-webui-forge"))
                candidates.append(os.path.join(base, sub, "stable-diffusion-webui-forge"))
                candidates.append(os.path.join(base, sub, "webui"))
        for name in ("sd-webui-forge", "stable-diffusion-webui-forge",
                     "stable-diffusion-webui"):
            for parent in (home, os.path.join(home, "Documents"),
                            os.path.join(home, "Downloads"),
                            os.path.join(home, "Desktop")):
                candidates.append(os.path.join(parent, name))
    else:
        for name in ("sd-webui-forge", "stable-diffusion-webui-forge",
                     "stable-diffusion-webui"):
            for parent in (home, os.path.join(home, "Documents"),
                            "/opt", "/usr/local"):
                candidates.append(os.path.join(parent, name))
    for c in candidates:
        _ld = _forge_launch_dir(c)
        if _ld:
            return _ld
    # Bounded fallback scan (NOT a whole-drive glob): list the immediate
    # children of each drive root + common parent folders, keep dirs whose name
    # hints at an SD/Forge install, and validate them + their immediate children
    # (catches nestings like <drive>:\<wrapper>\sd-webui-forge inside a folder
    # whose own name isn't a launcher). One level of listing per location —
    # fast, no deep recursion.
    scan_parents = []
    if sys.platform == "win32":
        scan_parents = [f"{d}:\\" for d in "CDEFG"]
    scan_parents += [home, os.path.join(home, "Documents"),
                     os.path.join(home, "Downloads"), os.path.join(home, "Desktop")]
    _hint = ("forge", "sd-webui", "stable-diffusion", "automatic1111", "a1111")
    for parent in scan_parents:
        try:
            entries = os.listdir(parent)
        except OSError:
            continue
        for name in entries:
            if not any(h in name.lower() for h in _hint):
                continue
            sub = os.path.join(parent, name)
            if not os.path.isdir(sub):
                continue
            _ld = _forge_launch_dir(sub)
            if _ld:
                return _ld
            # one level deeper (the install often sits inside a wrapper folder)
            try:
                for inner in os.listdir(sub):
                    _ld = _forge_launch_dir(os.path.join(sub, inner))
                    if _ld:
                        return _ld
            except OSError:
                pass
    return ""


def _forge_launch_dir(path: str) -> str:
    """Return the dir that actually contains webui.bat + modules/ (what we cd
    into to launch Forge), or "" if `path` isn't a Forge/SD-WebUI install.

    Accepts BOTH layouts:
      - cloned repo: webui.bat + modules/ AT the top level
      - one-click PACKAGE: run.bat at the top, the real webui.bat + modules/
        inside a `webui/` subfolder  (a top-level wrapper dir is otherwise rejected)
    """
    if not path or not os.path.isdir(path):
        return ""

    def _qualifies(d: str) -> bool:
        has_launcher = any(
            os.path.isfile(os.path.join(d, f))
            for f in ("webui.bat", "webui.sh", "webui-user.bat", "webui-user.sh", "launch.py")
        )
        return has_launcher and os.path.isdir(os.path.join(d, "modules"))

    if _qualifies(path):
        return path
    sub = os.path.join(path, "webui")  # Forge one-click package nests it here
    if _qualifies(sub):
        return sub
    return ""


def _looks_like_forge(path: str) -> bool:
    """True if `path` is (or contains, via a webui/ subdir) a Forge/SD-WebUI
    install. Two-signal check (launcher + modules/) guards against matching a
    random folder named "webui"."""
    return _forge_launch_dir(path) != ""


FORGE_DIR = _autodetect_forge_dir()
FORGE_URL = os.environ.get("JARVIS_FORGE_URL", "http://127.0.0.1:7860")
FORGE_BOOT_TIMEOUT = int(os.environ.get("JARVIS_FORGE_BOOT_TIMEOUT", "180"))

# Light up the local Forge image-gen tools when an install is actually present
# (or HEARTH_ENABLE_FORGE=1 forces them). They're deferred, so they don't bloat
# the prompt — the model reveals them with load_tools('image') and the user's
# local SD model (e.g. a Pony checkpoint) becomes usable. With no install and no
# flag, drop them so a normal user's toolset stays clean.
if not (os.environ.get("HEARTH_ENABLE_FORGE", "0") == "1"
        or (FORGE_DIR and os.path.isdir(FORGE_DIR))):
    TOOL_DEFINITIONS[:] = [t for t in TOOL_DEFINITIONS
                           if t["name"] not in {"forge_generate", "forge_status", "forge_shutdown"}]

# Track the subprocess so we can shut it down later. Module-level so it
# survives across tool calls within a single Jarvis session.
_forge_proc: Optional[subprocess.Popen] = None

# Auto-deload: Forge holds the whole SDXL checkpoint in VRAM (~6.5 GB) the entire
# time it's up. After a stretch with no generation we shut it down so the GPU
# isn't hoarded — the LLM (or the desktop) gets its VRAM back. Each
# forge_generate bumps the timer; a daemon watcher kills Forge once idle.
_forge_last_used: float = 0.0
_forge_idle_timeout: int = int(os.environ.get("JARVIS_FORGE_IDLE_TIMEOUT", "300"))
_forge_idle_watcher_on: bool = False
# When the active brain is Hearth's builtin llama.cpp server (not LM Studio),
# it can't JIT-unload — so to free VRAM for Forge we fully STOP it and remember
# what it was running, then restart it from this snapshot when Forge shuts down.
_builtin_paused: "Optional[Dict[str, Any]]" = None


def _start_forge_idle_watcher() -> None:
    """Spawn (once) a daemon that shuts Forge down after _forge_idle_timeout
    seconds without a generation. Set JARVIS_FORGE_IDLE_TIMEOUT=0 to disable."""
    global _forge_idle_watcher_on
    if _forge_idle_watcher_on or _forge_idle_timeout <= 0:
        return
    import threading as _threading
    _forge_idle_watcher_on = True

    def _watch():
        global _forge_idle_watcher_on
        import time as _t
        while True:
            _t.sleep(20)
            # Nothing to guard if we don't manage a live Forge process.
            if _forge_proc is None or _forge_proc.poll() is not None:
                _forge_idle_watcher_on = False
                return
            if _forge_last_used and (_t.time() - _forge_last_used) > _forge_idle_timeout:
                try:
                    _forge_shutdown({})
                    _log_activity("forge_idle_shutdown",
                                  {"idle_s": int(_t.time() - _forge_last_used)})
                except Exception:
                    pass
                _forge_idle_watcher_on = False
                return

    _threading.Thread(target=_watch, daemon=True, name="forge-idle").start()


def _forge_reachable(timeout: float = 1.5) -> bool:
    try:
        req = urllib.request.Request(f"{FORGE_URL}/sdapi/v1/options",
                                     method="GET")
        with urllib.request.urlopen(req, timeout=timeout) as r:
            return r.status == 200
    except Exception:
        return False


def _forge_status(p: Dict) -> str:
    proc_alive = _forge_proc is not None and _forge_proc.poll() is None
    reachable = _forge_reachable()
    info = {
        "forge_dir": FORGE_DIR,
        "url": FORGE_URL,
        "process_managed_by_jarvis": proc_alive,
        "api_reachable": reachable,
        "ready": reachable,
    }
    if not os.path.isdir(FORGE_DIR):
        info["error"] = f"FORGE_DIR not found: {FORGE_DIR} — set JARVIS_FORGE_DIR"
    return json.dumps(info, indent=2)


def _boot_forge() -> str:
    """Launch Forge with --api flag. Returns '' on success, error message on
    failure. Already-running Forge counts as success."""
    global _forge_proc
    if _forge_reachable():
        return ""  # already up — fine
    if not os.path.isdir(FORGE_DIR):
        return f"Error: Forge dir not found: {FORGE_DIR}"
    # Use the universal launcher which sets up venv + env, plus --api so
    # /sdapi/v1/* is exposed. --nowebui would disable the gradio UI entirely
    # but many setups need the UI thread alive; --api alongside the UI is
    # the safest default.
    webui_bat = os.path.join(FORGE_DIR, "webui.bat")
    if not os.path.isfile(webui_bat):
        # fall back to webui-user.bat which chains into webui.bat
        webui_bat = os.path.join(FORGE_DIR, "webui-user.bat")
    if not os.path.isfile(webui_bat):
        return f"Error: no webui.bat or webui-user.bat in {FORGE_DIR}"
    # Forge/A1111's --listen is a NO-ARGUMENT flag (it binds 0.0.0.0); the host
    # and port come from --server-name / --port. Passing "--listen=127.0.0.1"
    # makes launch.py abort ("argument --listen: ignored explicit argument").
    # 127.0.0.1:7860 is already Forge's default, so we only add overrides when
    # FORGE_URL points somewhere non-default.
    from urllib.parse import urlparse as _urlparse
    _fu = _urlparse(FORGE_URL)
    forge_args = ["cmd", "/c", webui_bat, "--api"]
    if _fu.port and _fu.port != 7860:
        forge_args += ["--port", str(_fu.port)]
    if _fu.hostname and _fu.hostname not in ("127.0.0.1", "localhost"):
        forge_args += ["--server-name", _fu.hostname]
    try:
        _forge_proc = subprocess.Popen(
            forge_args,
            cwd=FORGE_DIR,
            creationflags=getattr(subprocess, "CREATE_NEW_CONSOLE", 0),
            stdin=subprocess.DEVNULL,
        )
    except Exception as e:
        return f"Error launching Forge: {e}"
    # Poll until the API answers
    import time as _t
    deadline = _t.time() + FORGE_BOOT_TIMEOUT
    while _t.time() < deadline:
        if _forge_reachable(timeout=2.0):
            return ""
        if _forge_proc.poll() is not None:
            return f"Error: Forge process exited with code {_forge_proc.returncode}"
        _t.sleep(2)
    return f"Error: Forge didn't become reachable on {FORGE_URL} within {FORGE_BOOT_TIMEOUT}s"


def _try_release_llm_vram() -> str:
    """Best-effort: tell LM Studio to JIT-eject the loaded model so Forge
    can grab VRAM. LM Studio's native API (/api/v0/models/<id>) supports
    POST with body {"loaded": false}, but the exact shape varies by build.
    We try both shapes; failure is non-fatal (Forge can also coexist if
    you have enough VRAM headroom)."""
    # Builtin llama.cpp server: it holds the model in VRAM and can't JIT-unload,
    # so stop it outright and snapshot what it was running for restore later
    # (done in _forge_shutdown). This is the builtin half of the VRAM dance.
    global _builtin_paused
    try:
        from . import llmserver as _ls
        if _ls._proc is not None and _ls._proc.poll() is None:
            _builtin_paused = dict(_ls._proc_info or {})
            _ls.stop_builtin()
            return "stopped builtin LLM server (restored after image gen)"
    except Exception:
        _builtin_paused = None

    # Most reliable: LM Studio's own CLI (`lms unload --all`). The REST unload
    # endpoints aren't exposed on every build, but the CLI ships with LM Studio
    # and actually frees the VRAM. Try it first, on PATH or the default install.
    import shutil as _shutil
    lms = _shutil.which("lms") or os.path.expanduser(
        os.path.join("~", ".lmstudio", "bin", "lms.exe"))
    if lms and os.path.isfile(lms):
        try:
            subprocess.run([lms, "unload", "--all"], timeout=20,
                           capture_output=True,
                           creationflags=getattr(subprocess, "CREATE_NO_WINDOW", 0))
            return "unloaded via lms CLI"
        except Exception:
            pass

    base = os.environ.get("LOCAL_API_BASE", "http://localhost:1234/v1")
    native = base.replace("/v1", "/api/v0")
    candidates = [
        # Newer LM Studio: /api/v0/models with action
        (f"{native}/models/unload", b""),
        (f"{base}/models/unload", b""),
    ]
    for url, body in candidates:
        try:
            req = urllib.request.Request(url, data=body, method="POST")
            with urllib.request.urlopen(req, timeout=3) as r:
                if 200 <= r.status < 300:
                    return f"unloaded via {url}"
        except Exception:
            continue
    return "LM Studio unload not exposed — relying on VRAM headroom"


def _forge_shutdown(p: Dict) -> str:
    global _forge_proc
    msgs = []
    if _forge_proc is not None and _forge_proc.poll() is None:
        try:
            # Send Ctrl-Break to the new console group on Windows, then kill
            if sys.platform == "win32":
                _forge_proc.send_signal(subprocess.signal.CTRL_BREAK_EVENT)  # type: ignore[attr-defined]
        except Exception:
            pass
        try:
            _forge_proc.terminate()
            _forge_proc.wait(timeout=10)
            msgs.append("forge process terminated")
        except subprocess.TimeoutExpired:
            _forge_proc.kill()
            msgs.append("forge process killed (didn't terminate cleanly)")
        except Exception as e:
            msgs.append(f"shutdown error: {e}")
        _forge_proc = None
    else:
        # Maybe the user launched it manually — try the API's own shutdown
        try:
            req = urllib.request.Request(
                f"{FORGE_URL}/sdapi/v1/server-kill",
                data=b"",
                method="POST",
            )
            with urllib.request.urlopen(req, timeout=5) as _r:
                msgs.append("requested forge api shutdown")
        except Exception:
            msgs.append("no jarvis-managed forge to kill, and API kill failed")

    # Restore the builtin LLM server if we paused it to free VRAM for Forge.
    global _builtin_paused
    if _builtin_paused:
        mp = _builtin_paused.get("model_path")
        if mp and os.path.exists(mp):
            try:
                from . import llmserver as _ls
                _ls.start_builtin(mp, port=_builtin_paused.get("port"),
                                  ctx=int(_builtin_paused.get("ctx") or 24576))
                msgs.append("restored builtin LLM server")
            except Exception as e:
                msgs.append(f"builtin restore failed: {type(e).__name__}")
        _builtin_paused = None

    return "; ".join(msgs) or "nothing to do"


def _forge_generate(p: Dict) -> str:
    """Run the full txt2img pipeline. Returns the path to the saved image
    or an Error: line."""
    import base64 as _b64

    positive = (p.get("positive") or "").strip()
    if not positive:
        return "Error: positive prompt required"
    negative = (p.get("negative") or
                "score_6, score_5, score_4, low quality, blurry, deformed, "
                "extra limbs, bad anatomy")
    width = int(p.get("width") or 1024)
    height = int(p.get("height") or 1024)
    steps = int(p.get("steps") or 25)
    cfg = float(p.get("cfg_scale") or 6.0)
    sampler = p.get("sampler") or "Euler a"
    seed = int(p.get("seed") or -1)

    # 1) Free LM Studio's VRAM (best-effort)
    vram_msg = _try_release_llm_vram()

    # 2) Make sure Forge is up
    boot_err = _boot_forge()
    if boot_err:
        return boot_err
    # Mark active + start the idle watcher so Forge auto-deloads its VRAM once
    # we stop generating (bumped again after the image lands below).
    global _forge_last_used
    _forge_last_used = time.time()
    _start_forge_idle_watcher()

    # 3) txt2img
    body = json.dumps({
        "prompt": positive,
        "negative_prompt": negative,
        "width": width,
        "height": height,
        "steps": steps,
        "cfg_scale": cfg,
        "sampler_name": sampler,
        "seed": seed,
    }).encode("utf-8")
    try:
        req = urllib.request.Request(
            f"{FORGE_URL}/sdapi/v1/txt2img",
            data=body,
            headers={"Content-Type": "application/json"},
            method="POST",
        )
        with urllib.request.urlopen(req, timeout=600) as r:
            payload = json.loads(r.read().decode())
    except Exception as e:
        return f"Error: txt2img call failed: {e}"

    images = payload.get("images") or []
    if not images:
        return f"Error: Forge returned no images. Response info: {payload.get('info', '')[:300]}"

    # 4) Save the first image
    ts = datetime.now().strftime("%Y%m%d_%H%M%S")
    out_path = os.path.join(SHOTS_DIR, f"forge_{ts}.png")
    try:
        with open(out_path, "wb") as f:
            f.write(_b64.b64decode(images[0]))
    except Exception as e:
        return f"Error: failed to decode/save image: {e}"

    # Restart the idle countdown from completion of this generation.
    _forge_last_used = time.time()
    _idle_note = (f"auto-frees VRAM after {_forge_idle_timeout//60} min idle"
                  if _forge_idle_timeout > 0 else
                  "auto-deload off — call forge_shutdown to free VRAM")
    return (
        f"Saved: {out_path}\n"
        f"prompt: {positive}\n"
        f"seed: {payload.get('info', '')[:200]}\n"
        f"(VRAM: {vram_msg}. Forge stays up for follow-ups; {_idle_note}.)"
    )


def _set_voice(p: Dict) -> str:
    from . import voice as _voice
    name = (p.get("name") or "").strip()
    if not name:
        return "Error: missing voice name"
    _voice.set_default_voice(name)
    if _voice.is_available():
        _voice.stop()
        _voice.speak(f"voice set to {name}", blocking=False)
        return f"voice set to {name} (sample playing)"
    return f"voice id stored as {name}, but engine not loaded yet — check /voice"


def _list_voices(p: Dict) -> str:
    from . import voice as _voice
    return "Available Kokoro voices:\n  " + "\n  ".join(_voice.list_voices())


def _end_session(p: Dict) -> str:
    # The CLI watches for this tool name and exits after the response. The
    # tool itself is a no-op marker so it works the same way through MCP
    # (where it just acknowledges) or the CLI (where it triggers shutdown).
    return "session_end_signaled"


def _get_time(p: Dict) -> str:
    now = datetime.now().astimezone()
    return json.dumps({
        "iso": now.isoformat(timespec="seconds"),
        "local": now.strftime("%A, %d %B %Y %H:%M:%S"),
        "tz": now.tzname() or "?",
        "offset": now.strftime("%z"),
        "epoch": int(now.timestamp()),
    }, indent=2)


# Runtime info the agent can introspect via whoami(). The CLI / bridge / web
# frontends set this at startup so the model can answer "what model/endpoint
# am I" without guessing or spawning shell commands. Falls back to env vars.
_RUNTIME_INFO: Dict[str, Any] = {}


def set_runtime_info(**kw: Any) -> None:
    """Frontends call this once at startup with {model, endpoint, context_tokens}.
    Stored process-wide so the whoami tool can report accurate live values."""
    _RUNTIME_INFO.update({k: v for k, v in kw.items() if v is not None})


def _notify(p: Dict) -> str:
    msg = (p.get("message") or "").strip()
    if not msg:
        return "Error: notify needs a 'message'."
    title = (p.get("title") or "Hearth").strip()
    from . import reminders
    fired = reminders.desktop_notify(title, msg)
    return f"Notification shown: {title} — {msg}" + ("" if fired else " (toast lib missing; logged/spoken instead)")


def _whoami(p: Dict) -> str:
    base = _RUNTIME_INFO.get("endpoint") or os.getenv("LOCAL_API_BASE", "http://localhost:1234/v1")
    b = base.lower()
    is_local = any(h in b for h in ("localhost", "127.0.0.1", "0.0.0.0", "::1",
                                    "192.168.", "10.", "host.docker.internal"))
    model = _RUNTIME_INFO.get("model") or os.getenv("LOCAL_MODEL") or "unknown"
    ctx = _RUNTIME_INFO.get("context_tokens")
    try:
        from . import memory
        idx = memory.list_index()
        mem_count = 0 if idx.startswith("(no") else idx.count("\n") + 1
    except Exception:
        mem_count = 0
    try:
        from .persona import NAME as _persona_name
    except Exception:
        _persona_name = "JARVIS"
    info = {
        "agent": _persona_name,
        "framework": "Hearth",
        "model": model,
        "endpoint": base,
        "location": "local (this machine / LAN)" if is_local else "cloud",
        "context_tokens": ctx if ctx else "auto",
        "tools_available": len(TOOL_DEFINITIONS),
        "memories_stored": mem_count,
        "workspace": WORKSPACE,
        "repo": "https://github.com/0pen-sourcer/hearth",
        "license": "MIT",
    }
    return json.dumps(info, indent=2)


def _list_models(p: Dict) -> str:
    import urllib.request
    base = (_RUNTIME_INFO.get("endpoint") or os.getenv("LOCAL_API_BASE", "http://localhost:1234/v1")).rstrip("/")
    _key = os.environ.get("LOCAL_API_KEY") or ""
    _hdr = {"Authorization": f"Bearer {_key}"} if _key else {}
    try:
        _req = urllib.request.Request(base + "/models", headers=_hdr)
        with urllib.request.urlopen(_req, timeout=5) as r:
            data = json.loads(r.read().decode())
        ids = [m.get("id") for m in data.get("data", []) if m.get("id")]
        if not ids:
            return f"The server at {base} reports no models — is one loaded in LM Studio?"
        cur = _RUNTIME_INFO.get("model")
        lines = [f"Models available at {base}:"]
        lines += [f"  - {i}" + ("  (current)" if i == cur else "") for i in ids]
        return "\n".join(lines)
    except Exception as e:
        # Tailor the hint to whatever the user is currently pointed at. The
        # old message always said "Is LM Studio running?" even when the user
        # was on a cloud endpoint that 401'd because the API key was bad or
        # the endpoint doesn't expose /v1/models — that confused both the
        # model and the user (Grok went hunting for LM Studio when the real
        # fix was `/brain local`).
        is_cloud = not any(h in base.lower() for h in (
            "localhost", "127.0.0.1", "0.0.0.0", "::1",
            "192.168.", "10.", "host.docker.internal",
        ))
        err_kind = type(e).__name__
        err_str = str(e)
        is_auth = "401" in err_str or "403" in err_str or "Unauthorized" in err_str or "Forbidden" in err_str
        if is_cloud and is_auth:
            return (f"Cloud endpoint {base} returned auth error ({err_kind}: {e}). "
                    f"Either the API key is wrong/expired, or this provider doesn't "
                    f"expose /v1/models. To list LOCAL models instead, the user can "
                    f"run `/brain local` (CLI) or switch in Settings → Chat brain (GUI). "
                    f"Don't scan the disk for model files.")
        if is_cloud:
            return (f"Cloud endpoint {base} unreachable ({err_kind}: {e}). "
                    f"Hearth's /v1/models proxy isn't responding. If the user wants "
                    f"local models, suggest `/brain local`. Don't scan the disk.")
        return (f"Couldn't reach the local server at {base}: {err_kind}: {e}. "
                f"Is LM Studio (or your endpoint) running? Don't scan the disk for model files.")


def _learn_environment(p: Dict) -> str:
    from .environment import learn_environment
    endpoint = _RUNTIME_INFO.get("endpoint") or os.getenv("LOCAL_API_BASE")
    return learn_environment(endpoint=endpoint)


# ---------------------------------------------------------------------------
# Media generation (image + video). The returned strings use the same marker
# convention as view_image so the GUI's existing media-render path picks them
# up. CLI parses for the path and pops the OS default viewer.
# ---------------------------------------------------------------------------

def _generate_image(p: Dict) -> str:
    from . import imagine
    r = imagine.generate_image(
        prompt=(p.get("prompt") or "").strip(),
        n=_coerce_int(p.get("n"), 1),
        aspect_ratio=(p.get("aspect_ratio") or "1:1"),
        resolution=(p.get("resolution") or "1k"),
    )
    if not r.get("ok"):
        return f"Error: {r.get('error')}"
    paths = r.get("paths") or []
    if not paths:
        return "Error: provider returned no images."
    # Emit one marker per file so the GUI renders each inline. CLI picks the
    # first one to auto-open in the default viewer.
    lines = [f"Generated {len(paths)} image{'s' if len(paths) != 1 else ''} via {r.get('provider')} "
             f"({r.get('model')}): \"{r.get('prompt')[:80]}\""]
    for path in paths:
        try:
            size = os.path.getsize(path)
        except OSError:
            size = 0
        lines.append(f"__JARVIS_IMAGE__ {path} ({size} bytes, png)")
    return "\n".join(lines)


def _generate_video(p: Dict) -> str:
    from . import imagine
    r = imagine.start_video(
        prompt=(p.get("prompt") or "").strip(),
        duration=_coerce_int(p.get("duration"), 5),
        aspect_ratio=(p.get("aspect_ratio") or "16:9"),
        resolution=(p.get("resolution") or "720p"),
        image_url=(p.get("image_url") or None),
    )
    if not r.get("ok"):
        return f"Error: {r.get('error')}"
    return (
        f"Video generation started (task_id: {r.get('task_id')}).\n"
        f"  model:  {r.get('model')}\n"
        f"  prompt: \"{r.get('prompt')[:80]}\"\n"
        f"{r.get('hint', '')}\n\n"
        f"NEXT STEP: tell the user the video is cooking and they can ask "
        f"'is the video ready?' later. Don't poll in a tight loop."
    )


def _check_video_task(p: Dict) -> str:
    from . import imagine
    task_id = (p.get("task_id") or "").strip()
    if not task_id:
        return "Error: task_id is required."
    r = imagine.check_video_task(task_id)
    if not r.get("ok"):
        return f"Error: {r.get('error')}"
    status = r.get("status", "unknown")
    if status == "done":
        path = r.get("path") or "(no local path)"
        return (
            f"Video DONE.\n"
            f"  path:     {path}\n"
            f"  duration: {r.get('duration', '?')}s\n"
            f"__JARVIS_VIDEO__ {path}"
        )
    if status == "pending":
        return f"Video still cooking (task {task_id}). Check back in 10-20 seconds."
    return f"Video task {task_id}: status={status}. {r.get('error', '')}".strip()


def _list_generations(p: Dict) -> str:
    from . import imagine
    r = imagine.list_recent_tasks(limit=10)
    tasks = r.get("tasks") or []
    if not tasks:
        return "No generation tasks on file yet."
    lines = [f"{len(tasks)} recent generation task{'s' if len(tasks) != 1 else ''}:"]
    for t in tasks:
        ts = time.strftime("%Y-%m-%d %H:%M",
                           time.localtime(t.get("created_at", 0)))
        lines.append(
            f"  [{t.get('status','?'):<8}] {t.get('kind','?'):<5} "
            f"{ts}  {t.get('task_id', '?')[:12]}…  "
            f"\"{(t.get('prompt') or '')[:50]}\""
        )
    return "\n".join(lines)


# ============================================================
# DISPATCH
# ============================================================

_HANDLERS = {
    "read_file": _read_file,
    "write_file": _write_file,
    "edit_file": _edit_file,
    "list_archive": _list_archive,
    "extract_archive_file": _extract_archive_file,
    "summarize_file": _summarize_file,
    "search_chats":   lambda p: __import__("hearth.session_search", fromlist=["search","format_matches"]).format_matches(__import__("hearth.session_search", fromlist=["search"]).search(p.get("query", ""), int(p.get("limit") or 8))),
    "set_reminder":   lambda p: __import__("hearth.reminders", fromlist=["set_reminder"]).set_reminder(
        p.get("when", ""), p.get("what", ""), p.get("recurring", ""),
        action_tool=p.get("action_tool", ""), action_args=p.get("action_args"),
        tag=p.get("tag", "")),
    "snooze_reminder": lambda p: __import__("hearth.reminders", fromlist=["snooze_reminder"]).snooze_reminder(
        p.get("id", ""), int(p.get("minutes", 10))),
    "list_reminders": lambda p: __import__("hearth.reminders", fromlist=["list_reminders"]).list_reminders(bool(p.get("include_fired"))),
    "cancel_reminder": lambda p: {"ok": __import__("hearth.reminders", fromlist=["cancel_reminder"]).cancel_reminder(p.get("id", ""))},
    "spawn_subagent": lambda p: __import__("hearth.subagents", fromlist=["spawn_subagent"]).spawn_subagent(
        p.get("persona", ""), p.get("prompt", ""),
        max_turns=min(int(p.get("max_turns") or 0) or 0, 20) or None,
        mode=(p.get("mode") or "sync"),
        name=p.get("name", "")),
    "list_subagent_personas": lambda p: __import__("hearth.subagents", fromlist=["list_personas"]).list_personas(),
    "get_subagent_result": lambda p: __import__("hearth.subagents", fromlist=["get_subagent_result"]).get_subagent_result(p.get("agent_id", "")),
    "list_directory": _list_directory,
    "create_directory": _create_directory,
    "delete_path": _delete_path,
    "move_path": _move_path,
    "grep_search": _grep_search,
    "glob_files": _glob_files,
    "find_file": _find_file,
    "web_search": _web_search,
    "web_fetch": _web_fetch,
    "run_command": _run_command,
    # Background-job tracking — see hearth/jobs.py. Lazy-import so a missing
    # ~/Jarvis/jobs/ dir or a stale module doesn't break tool registration.
    "start_job":  lambda p: __import__("hearth.jobs", fromlist=["start_job"]).start_job(
        p.get("command", ""), cwd=p.get("cwd"),
        shell=(p.get("shell") or "powershell"),
        description=p.get("description", "")),
    "job_status": lambda p: __import__("hearth.jobs", fromlist=["get_job"]).get_job(p.get("job_id", "")),
    "job_wait":   lambda p: __import__("hearth.jobs", fromlist=["wait_job"]).wait_job(
        p.get("job_id", ""), timeout_s=min(float(p.get("timeout_s") or 30), 300.0)),
    "job_kill":   lambda p: __import__("hearth.jobs", fromlist=["kill_job"]).kill_job(p.get("job_id", "")),
    "job_list":   lambda p: __import__("hearth.jobs", fromlist=["list_jobs"]).list_jobs(active_only=bool(p.get("active_only"))),
    "generate_image":    _generate_image,
    "generate_video":    _generate_video,
    "check_video_task":  _check_video_task,
    "list_generations":  _list_generations,
    "ask_user":   _ask_user,
    "system_info": _system_info,
    "list_processes": _list_processes,
    "network_info": _network_info,
    "get_battery": _get_battery,
    "list_installed_apps": _list_installed_apps,
    "disk_usage": _disk_usage,
    "list_jobs": lambda p: __import__("hearth.jobs", fromlist=["list_jobs"]).list_jobs(bool(p.get("active_only"))),
    "get_job_result": lambda p: __import__("hearth.jobs", fromlist=["get_job_result"]).get_job_result(p.get("job_id", "")),
    "locate_path": _locate_path,
    "open_app": _open_app,
    "open_url": _open_url,
    "open_in_browser": _open_in_browser,
    "list_browsers": _list_browsers,
    "screenshot": _screenshot,
    "view_image": _view_image,
    "clipboard_read": _clipboard_read,
    "clipboard_write": _clipboard_write,
    "memory_save": _memory_save,
    "memory_recall": _memory_recall,
    "memory_list": _memory_list,
    "memory_forget": _memory_forget,
    "edit_soul":   lambda p: __import__("hearth.memory", fromlist=["write_soul"]).write_soul(p.get("content", "")),
    "append_soul": lambda p: __import__("hearth.memory", fromlist=["append_soul"]).append_soul(p.get("line", "")),
    "read_soul":   lambda p: (__import__("hearth.memory", fromlist=["read_soul"]).read_soul() or "(soul.md is empty — write one with edit_soul or call draft_soul for a starter)"),
    "draft_soul":  lambda p: __import__("hearth.memory", fromlist=["draft_soul"]).draft_soul(),
    "set_voice": _set_voice,
    "list_voices": _list_voices,
    "forge_generate": _forge_generate,
    "forge_status": _forge_status,
    "forge_shutdown": _forge_shutdown,
    "end_session": _end_session,
    "get_time": _get_time,
    "whoami": _whoami,
    "list_models": _list_models,
    "learn_environment": _learn_environment,
    "notify": _notify,
}


# ----- PLUGIN SYSTEM ------------------------------------------------------
# Local, private, self-improving: the agent can author its own tools and any
# ~/Jarvis/plugins/*.py is auto-loaded. See hearth/plugins.py. plugins.py takes
# the registry as args (no import of tools) so there's no circular import.
from . import plugins as _plugins


def _create_plugin(p: Dict) -> str:
    name = (p.get("name") or "").strip()
    code = p.get("code") or ""
    if not name or not code:
        return ("Error: create_plugin needs 'name' (lower_snake_case) and 'code' "
                "(the FULL plugin module: a module-level TOOL dict + a run(args) function).")
    return _plugins.save_and_register(name, code, WORKSPACE, TOOL_DEFINITIONS, _HANDLERS)


def _list_plugins(p: Dict) -> str:
    return _plugins.list_plugins(WORKSPACE)


def _delete_plugin(p: Dict) -> str:
    return _plugins.delete_plugin(p.get("name", ""), WORKSPACE, TOOL_DEFINITIONS, _HANDLERS)


_HANDLERS["create_plugin"] = _create_plugin
_HANDLERS["list_plugins"] = _list_plugins
_HANDLERS["delete_plugin"] = _delete_plugin


# ----- INTERACTIVE BROWSER (Playwright-driven Chromium; see hearth/browse.py) -----
def _browse(p: Dict) -> str:
    from . import browse as _b
    return _b.browse(p)


def _browse_click(p: Dict) -> str:
    from . import browse as _b
    return _b.browse_click(p)


def _browse_type(p: Dict) -> str:
    from . import browse as _b
    return _b.browse_type(p)


def _browse_close(p: Dict) -> str:
    from . import browse as _b
    return _b.browse_close(p)


def _browse_scroll(p: Dict) -> str:
    from . import browse as _b
    return _b.browse_scroll(p)


def _browse_key(p: Dict) -> str:
    from . import browse as _b
    return _b.browse_key(p)


_HANDLERS["browse"] = _browse
_HANDLERS["browse_click"] = _browse_click
_HANDLERS["browse_type"] = _browse_type
_HANDLERS["browse_scroll"] = _browse_scroll
_HANDLERS["browse_key"] = _browse_key
_HANDLERS["browse_close"] = _browse_close

# The browser tools register ONLY when Playwright is installed (opt-in via
# `install.ps1 -Browser`). Otherwise the model would see browse + web_fetch +
# web_search and not know which to reach for, and we'd waste prompt tokens on
# tools that can't run. HEARTH_ENABLE_BROWSER=1 forces on, =0 forces off.
import importlib.util as _ilu
_b_env = os.environ.get("HEARTH_ENABLE_BROWSER")
_browser_enabled = _b_env == "1" or (_b_env != "0" and _ilu.find_spec("playwright") is not None)

for _bt in (
    {
        "name": "browse",
        "description": (
            "Drive a REAL web browser (controlled Chromium) — the user SEES the window. "
            "Navigate, get the rendered text + a numbered list of clickable links/buttons, "
            "then browse_click / browse_type. The session stays open across calls. "
            "PREFER browse whenever you might need to KEEP control of what's on screen — "
            "e.g. playing a video the user may then want changed ('next one', 'search "
            "something else'), browsing a site, filling a form. With browse YOU can do "
            "all that next; with open_url you can't (it hands off to their browser and you "
            "lose control). Only use open_url/open_in_browser for a pure one-off handoff. "
            "Leave the browse session OPEN while the user is still using it. It's a SINGLE "
            "tab — each browse() navigation REPLACES the current page, so for 'open these "
            "in separate tabs for me' use open_url (once per link), not browse. For a quick "
            "text grab (you read it), use web_fetch; to search, web_search."
        ),
        "parameters": {
            "type": "object",
            "properties": {"url": {"type": "string", "description": "URL to open (https:// assumed if omitted). Leave empty to re-read the current page."}},
        },
    },
    {
        "name": "browse_click",
        "description": "Click a link or button on the CURRENT browser page by its visible text (use the exact text from the CLICKABLE list that browse returned). Returns the new page.",
        "parameters": {
            "type": "object",
            "properties": {"text": {"type": "string", "description": "Visible text of the link/button to click."}},
            "required": ["text"],
        },
    },
    {
        "name": "browse_scroll",
        "description": "Smooth-scroll the current browser page so you can see more content. Use this when the user asks for 'everything on the page', a long article, or when initial summary cuts off mid-section. direction='down' (default), 'up', 'top', or 'bottom'. Optional pixels= for a specific amount; otherwise scrolls one viewport.",
        "parameters": {
            "type": "object",
            "properties": {
                "direction": {"type": "string", "description": "down (default) / up / top / bottom"},
                "pixels": {"type": "integer", "description": "Optional - how far to scroll. Defaults to one viewport."},
            },
        },
    },
    {
        "name": "browse_type",
        "description": "Type into a field on the current browser page (a search box, login field, etc.). Set submit=true to press Enter after.",
        "parameters": {
            "type": "object",
            "properties": {
                "text": {"type": "string", "description": "What to type."},
                "field": {"type": "string", "description": "Optional: the field's label/placeholder. If omitted, types into the first text input."},
                "submit": {"type": "boolean", "description": "Press Enter after typing (e.g. to run a search)."},
            },
            "required": ["text"],
        },
    },
    {
        "name": "browse_key",
        "description": (
            "Press a keyboard key/shortcut on the current browser page — for media "
            "controls and shortcuts that aren't clickable buttons. YouTube: key='f' "
            "fullscreen, 'k' or ' ' play/pause, 'm' mute, 't' theater, 'j'/'l' seek "
            "10s back/forward. Also 'Escape', 'ArrowUp'/'ArrowDown' volume, combos like "
            "'Control+L'. To fullscreen a video just call browse_key with key='f' (the "
            "video is auto-focused first). Optional focus= text/element to click before "
            "the keypress lands."
        ),
        "parameters": {
            "type": "object",
            "properties": {
                "key": {"type": "string", "description": "Key or combo, e.g. 'f' (fullscreen), 'k' (play/pause), 'm' (mute), 'Control+L'."},
                "focus": {"type": "string", "description": "Optional: text/element to click first so the keypress targets it."},
            },
            "required": ["key"],
        },
    },
    {
        "name": "browse_close",
        "description": "Close the browser session. ONLY when the user is truly DONE with it. If they want to WATCH a video, READ a page, or are still looking at what you opened, LEAVE IT OPEN — closing it yanks away what they're viewing. When in doubt, don't close.",
        "parameters": {"type": "object", "properties": {}},
    },
):
    if _browser_enabled:
        TOOL_DEFINITIONS.append(_bt)

TOOL_DEFINITIONS.append({
    "name": "create_plugin",
    "description": (
        "Write a NEW local tool (plugin) for yourself when no existing tool fits "
        "a capability the user needs — then use it immediately and forever after. "
        "100% local, saved to ~/Jarvis/plugins/<name>.py. The `code` must be a "
        "complete Python module defining EXACTLY two module-level names:\n"
        "  TOOL = {\"name\": \"<same as name>\", \"description\": \"...\", "
        "\"parameters\": {\"type\": \"object\", \"properties\": {...}, \"required\": [...]}}\n"
        "  def run(args: dict) -> str: ...  # returns a string result\n"
        "Use the Python stdlib freely. Don't shadow a built-in tool name. Keep it "
        "SMALL — one focused capability, a few lines (a 100+ line module usually "
        "means you're overcomplicating it). If create_plugin rejects your code, FIX "
        "that error and call create_plugin AGAIN — never fall back to write_file / "
        "run_command. Only create a plugin for a genuinely missing capability."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "name": {"type": "string", "description": "lower_snake_case tool name; also the filename."},
            "code": {"type": "string", "description": "Full plugin module source (TOOL dict + run(args) function)."},
        },
        "required": ["name", "code"],
    },
})

TOOL_DEFINITIONS.append({
    "name": "list_plugins",
    "description": "List the self-authored/installed plugins in ~/Jarvis/plugins/ (name, status, description). Use to see what custom tools exist before creating or deleting one.",
    "parameters": {"type": "object", "properties": {}},
})

TOOL_DEFINITIONS.append({
    "name": "delete_plugin",
    "description": "Delete an installed plugin by name (removes its file + unregisters the tool). Only plugins, never a built-in tool. Use when a plugin is broken/unwanted or the user asks to remove a custom tool.",
    "parameters": {
        "type": "object",
        "properties": {"name": {"type": "string", "description": "The plugin/tool name to delete."}},
        "required": ["name"],
    },
})

# ----- SKILLS (prose + asset bundles; distinct from plugins) -----
# Skills live as folders with SKILL.md + scripts/. The catalog (name + one-line
# description) is in the system prompt; full body only loads when the model
# decides to USE a skill. See hearth/skills_loader.py.
def _list_skills(p: Dict) -> str:
    from . import skills_loader as _sl
    items = _sl.list_skills(include_body=False)
    if not items:
        return "(no skills installed — drop a folder with SKILL.md in ~/Jarvis/skills/)"
    return json.dumps(items, ensure_ascii=False, indent=2)


def _load_skill(p: Dict) -> str:
    from . import skills_loader as _sl
    res = _sl.load_skill(p.get("name", ""))
    return json.dumps(res, ensure_ascii=False, indent=2)


def _create_skill(p: Dict) -> str:
    from . import skills_loader as _sl
    res = _sl.create_skill(
        name=p.get("name", ""),
        description=p.get("description", ""),
        body=p.get("body", ""),
        scripts=p.get("scripts") or None,
    )
    return json.dumps(res, ensure_ascii=False, indent=2)


_HANDLERS["list_skills"] = _list_skills
_HANDLERS["load_skill"] = _load_skill
_HANDLERS["create_skill"] = _create_skill

TOOL_DEFINITIONS.append({
    "name": "list_skills",
    "description": (
        "List available skills (bundled + user-installed). Each entry is "
        "{name, description, version, folder, source}. Skills are prose + asset "
        "bundles (e.g. make-pdf, make-pptx, make-xlsx) the user invokes via "
        "existing tools after you load_skill(<name>) for the full instructions. "
        "Catalog also appears in your system prompt — use this tool when you "
        "want the structured response."
    ),
    "parameters": {"type": "object", "properties": {}},
})

TOOL_DEFINITIONS.append({
    "name": "load_skill",
    "description": (
        "Load the full SKILL.md body + asset manifest for one skill. Call this "
        "BEFORE following its workflow so you have the exact steps + the "
        "bundled script paths (scripts/, references/, assets/). Returns "
        "{ok, name, description, body, scripts, references, assets, folder, ...}. "
        "Use the `folder` field to construct script paths."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "name": {"type": "string", "description": "Skill slug from list_skills (e.g. 'make-pdf')."},
        },
        "required": ["name"],
    },
})

TOOL_DEFINITIONS.append({
    "name": "create_skill",
    "description": (
        "Author a NEW user skill saved to ~/Jarvis/skills/<name>/. Use when "
        "you find yourself running the same multi-step workflow twice — "
        "crystallize it as a skill so future-you (and the user) can invoke "
        "it by name. `name` is lower-kebab-case (e.g. 'summarize-youtube'). "
        "`description` is the one-line when-to-use (the catalog summary). "
        "`body` is the markdown workflow body (steps, defaults, don'ts). "
        "Optional `scripts` is {filename: source} written under scripts/. "
        "Available immediately via list_skills + load_skill."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "name": {"type": "string"},
            "description": {"type": "string"},
            "body": {"type": "string"},
            "scripts": {"type": "object",
                "description": "Optional {filename: source_text} dict of helper scripts."},
        },
        "required": ["name", "description", "body"],
    },
})

def _install_skill(p: Dict) -> str:
    from . import skills_install as _si
    source = (p.get("source") or "").strip()
    if not source:
        return json.dumps({"ok": False, "error":
            "source required: a GitHub repo (owner/repo[/subdir][@ref]) or a local folder path"})
    if not p.get("confirm"):
        # Phase 1: fetch + inspect, disclose what it can do, DO NOT install.
        m = _si.inspect_source(source)
        if not m.get("ok"):
            return json.dumps(m, ensure_ascii=False)
        keep = ("name", "description", "version", "author", "declared_tools",
                "risky_tools", "scripts", "ships_code", "risky",
                "already_installed", "shadows_bundled")
        out = {k: m[k] for k in keep if k in m}
        out["ok"] = True
        out["action_required"] = (
            "Show the user the skill's name + what it does, and tell them it "
            + ("SHIPS SCRIPTS and can run shell commands / modify files (its "
               "scripts execute via run_command, which still prompts you)."
               if m.get("risky") else "uses only safe tools.")
            + " Install ONLY after the user agrees, by calling install_skill "
              "again with confirm=true.")
        return json.dumps(out, ensure_ascii=False, indent=2)
    # Phase 2: user confirmed -> install (re-fetches; placing files is harmless
    # until the skill is actually run, and running its scripts is gated by
    # run_command's own permission prompt).
    res = _si.install_skill(source, consent=lambda _m: True, force=bool(p.get("force")))
    return json.dumps(res, ensure_ascii=False, indent=2)


_HANDLERS["install_skill"] = _install_skill

TOOL_DEFINITIONS.append({
    "name": "install_skill",
    "description": (
        "Install a shareable skill from a GitHub repo or a local folder into "
        "~/Jarvis/skills/. Source forms: 'owner/repo', 'owner/repo@branch', "
        "'owner/repo/subdir', a full github.com URL, or a local path. "
        "TWO-PHASE for safety: call first WITHOUT confirm to fetch + inspect — "
        "you get {name, description, scripts, risky_tools, risky}. Relay that to "
        "the user (especially if `risky` — it ships scripts that run shell "
        "commands), and only after they agree, call again with confirm=true to "
        "install. Use when the user pastes a skill link or asks to add a skill. "
        "After install, call load_skill(name) to use it."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "source": {"type": "string", "description":
                "GitHub repo (owner/repo[/subdir][@ref]), a github.com URL, or a local folder path."},
            "confirm": {"type": "boolean", "description":
                "false/omitted = inspect only (disclose, no install). true = install (after the user agreed)."},
            "force": {"type": "boolean", "description":
                "Overwrite if a skill with the same name is already installed."},
        },
        "required": ["source"],
    },
})

def _read_inbox(p: Dict) -> str:
    from . import email_tools as _e
    res = _e.read_inbox(limit=p.get("limit", 10),
                        unread_only=bool(p.get("unread_only")),
                        folder=p.get("folder") or "INBOX",
                        with_body=bool(p.get("with_body")))
    return json.dumps(res, ensure_ascii=False, indent=2)


def _send_email(p: Dict) -> str:
    from . import email_tools as _e
    res = _e.send_email(to=p.get("to", ""), subject=p.get("subject", ""),
                        body=p.get("body", ""), cc=p.get("cc", ""),
                        reply_to=p.get("reply_to", ""))
    return json.dumps(res, ensure_ascii=False, indent=2)


_HANDLERS["read_inbox"] = _read_inbox
_HANDLERS["send_email"] = _send_email

TOOL_DEFINITIONS.append({
    "name": "read_inbox",
    "description": (
        "Read recent email from the user's configured mailbox (IMAP, read-only). "
        "Returns {from, subject, date, date_iso, body?} per message, newest first. "
        "Set unread_only=true for just unread, with_body=true to include the text "
        "body (truncated). Requires the user to have set up email (app password) — "
        "if not, the result explains how. Don't guess the user's address."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "limit": {"type": "integer", "description": "How many recent messages (1-50, default 10)."},
            "unread_only": {"type": "boolean", "description": "Only unread messages."},
            "folder": {"type": "string", "description": "Mailbox folder (default INBOX)."},
            "with_body": {"type": "boolean", "description": "Include the (truncated) plain-text body."},
        },
    },
})

TOOL_DEFINITIONS.append({
    "name": "send_email",
    "description": (
        "Send a plain-text email from the user's configured address (SMTP). "
        "Confirm the recipient, subject, and body with the user before sending "
        "unless they were explicit. Requires email setup (app password); if not "
        "configured the result explains how. This actually sends mail — treat it "
        "like any outward action."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "to": {"type": "string", "description": "Recipient address(es), comma-separated."},
            "subject": {"type": "string"},
            "body": {"type": "string"},
            "cc": {"type": "string", "description": "Optional CC, comma-separated."},
            "reply_to": {"type": "string", "description": "Optional Reply-To."},
        },
        "required": ["to", "subject", "body"],
    },
})

def _focus_window_macos(name: str) -> str:
    """Raise a macOS app/window by name via AppleScript. Matches against
    running application names (substring, case-insensitive)."""
    list_script = 'tell application "System Events" to get name of (processes where background only is false)'
    try:
        r = subprocess.run(["osascript", "-e", list_script],
                           capture_output=True, text=True, timeout=10)
    except FileNotFoundError:
        return "focus_window needs osascript (it ships with macOS — are you on macOS?)."
    except Exception as e:
        return f"Error listing windows: {e}"
    apps = [a.strip() for a in (r.stdout or "").split(",") if a.strip()]
    match = next((a for a in apps if name.lower() in a.lower()), None)
    if not match:
        sample = ", ".join(sorted(apps)[:12])
        return f"No open app matching '{name}'. Open apps: {sample or '(none)'}"
    try:
        subprocess.run(["osascript", "-e", f'tell application "{match}" to activate'],
                       capture_output=True, text=True, timeout=10)
    except Exception as e:
        return f"Error focusing '{match}': {e}"
    return f"Brought '{match}' to the front."


def _focus_window_linux(name: str) -> str:
    """Raise an X11 window by title substring. Prefers wmctrl (-a does a
    substring activate); falls back to xdotool. Tells the user what to install
    if neither is present (common on minimal/Wayland setups)."""
    wmctrl = shutil.which("wmctrl")
    if wmctrl:
        # List first so we can report matches / give a useful miss message.
        try:
            lst = subprocess.run([wmctrl, "-l"], capture_output=True, text=True, timeout=10)
        except Exception as e:
            return f"Error listing windows: {e}"
        rows = [ln for ln in (lst.stdout or "").splitlines() if ln.strip()]
        titles = [" ".join(ln.split(None, 3)[3:]) for ln in rows if len(ln.split(None, 3)) >= 4]
        if not any(name.lower() in t.lower() for t in titles):
            sample = ", ".join(sorted({t for t in titles if t})[:12])
            return f"No open window matching '{name}'. Open windows: {sample or '(none)'}"
        try:
            subprocess.run([wmctrl, "-a", name], capture_output=True, text=True, timeout=10)
        except Exception as e:
            return f"Error focusing '{name}': {e}"
        return f"Brought a window matching '{name}' to the front."
    xdotool = shutil.which("xdotool")
    if xdotool:
        try:
            r = subprocess.run([xdotool, "search", "--name", name, "windowactivate"],
                               capture_output=True, text=True, timeout=10)
        except Exception as e:
            return f"Error focusing '{name}': {e}"
        if r.returncode != 0:
            return f"No open window matching '{name}' (via xdotool)."
        return f"Brought a window matching '{name}' to the front."
    return ("focus_window needs wmctrl or xdotool on Linux. Install with "
            "'sudo apt install wmctrl' (Debian/Ubuntu/Mint) or your distro's "
            "package manager. (Note: window raising is X11-only — it won't work "
            "under a pure Wayland session.)")


def _focus_window(p: Dict) -> str:
    """Bring an already-open window to the front by (partial) title match.
    Like open_app, but for windows that are ALREADY open — raise/focus them."""
    name = (p.get("name") or "").strip()
    if not name:
        return "Error: window name required (a substring of the window title)."
    if sys.platform == "darwin":
        return _focus_window_macos(name)
    if sys.platform != "win32":
        return _focus_window_linux(name)
    try:
        import win32con
        import win32gui
    except Exception:
        return "focus_window needs pywin32 (pip install pywin32)."
    matches: List = []
    titles: List[str] = []

    def _cb(hwnd, _):
        if win32gui.IsWindowVisible(hwnd):
            t = win32gui.GetWindowText(hwnd)
            if t:
                titles.append(t)
                if name.lower() in t.lower():
                    matches.append((hwnd, t))

    try:
        win32gui.EnumWindows(_cb, None)
    except Exception as e:
        return f"Error enumerating windows: {e}"
    if not matches:
        sample = ", ".join(sorted({t for t in titles if t})[:12])
        return (f"No open window matching '{name}'. Open windows: {sample or '(none)'}")
    hwnd, title = matches[0]
    try:
        win32gui.ShowWindow(hwnd, win32con.SW_RESTORE)  # un-minimize if needed
    except Exception:
        pass
    try:
        win32gui.SetForegroundWindow(hwnd)
    except Exception:
        # Foreground-lock fallback: float it to top briefly, then release.
        try:
            win32gui.SetWindowPos(hwnd, win32con.HWND_TOPMOST, 0, 0, 0, 0,
                                  win32con.SWP_NOMOVE | win32con.SWP_NOSIZE | win32con.SWP_NOACTIVATE)
            win32gui.SetWindowPos(hwnd, win32con.HWND_NOTOPMOST, 0, 0, 0, 0,
                                  win32con.SWP_NOMOVE | win32con.SWP_NOSIZE | win32con.SWP_NOACTIVATE)
            win32gui.BringWindowToTop(hwnd)
        except Exception:
            pass
    extra = f" ({len(matches)} matched; brought the first)" if len(matches) > 1 else ""
    return f"Brought '{title}' to the front.{extra}"


def _send_to_phone(p: Dict) -> str:
    """Push a message FROM this PC TO one of the user's connected channels
    (Discord DM / Telegram / ntfy). Outbound counterpart to the bridges, which
    only reply to incoming messages. Reads the saved bridge config in ~/.hearth."""
    import json as _j
    import urllib.request as _u
    import urllib.parse as _up
    channel = (p.get("channel") or "").strip().lower()
    message = (p.get("message") or "").strip()
    if not message:
        return "Error: message is required."
    base = os.path.join(os.path.expanduser("~"), ".hearth")

    def _cfg(fn):
        try:
            with open(os.path.join(base, fn), encoding="utf-8") as f:
                return _j.load(f) or {}
        except Exception:
            return {}

    if channel == "telegram":
        c = _cfg("phone_bridge.json")
        tok = (c.get("bot_token") or "").strip()
        ids = c.get("allowed_chat_ids") or []
        if not tok or not ids:
            return "Error: Telegram isn't set up (need a bot token + your chat id in Settings - Reach from phone)."
        try:
            data = _up.urlencode({"chat_id": ids[0], "text": message}).encode()
            req = _u.Request(f"https://api.telegram.org/bot{tok}/sendMessage", data=data,
                             headers={"Content-Type": "application/x-www-form-urlencoded", "User-Agent": "Hearth"})
            _u.urlopen(req, timeout=12)
            return "Sent to your Telegram."
        except Exception as e:
            return f"Error sending to Telegram: {type(e).__name__}: {e}"

    if channel == "discord":
        c = _cfg("discord_bridge.json")
        tok = (c.get("bot_token") or "").strip()
        uids = c.get("allowed_user_ids") or []
        if not tok or not uids:
            return "Error: Discord isn't set up (need a bot token + your user id in Settings - Reach from phone)."
        try:
            h = {"Authorization": f"Bot {tok}", "Content-Type": "application/json", "User-Agent": "Hearth (https://github.com/0pen-sourcer/hearth, 0.7)"}
            # Open (or reuse) a DM channel with the owner, then post.
            dm_req = _u.Request("https://discord.com/api/v10/users/@me/channels",
                                data=_j.dumps({"recipient_id": str(uids[0])}).encode(), headers=h)
            chan = _j.loads(_u.urlopen(dm_req, timeout=12).read().decode())
            msg_req = _u.Request(f"https://discord.com/api/v10/channels/{chan['id']}/messages",
                                 data=_j.dumps({"content": message[:1900]}).encode(), headers=h)
            _u.urlopen(msg_req, timeout=12)
            return "Sent to your Discord (DM)."
        except Exception as e:
            return f"Error sending to Discord: {type(e).__name__}: {e}"

    if channel == "ntfy":
        topic = ""
        try:
            from . import reminders as _r  # reuses the configured ntfy topic
            topic = _r._ntfy_topic() if hasattr(_r, "_ntfy_topic") else ""
        except Exception:
            topic = ""
        topic = topic or os.environ.get("HEARTH_NTFY_TOPIC", "") or (_cfg("discord_bridge.json").get("ntfy_topic") or "")
        if not topic:
            return "Error: no ntfy topic set (Settings - Reach from phone - ntfy)."
        try:
            req = _u.Request(f"https://ntfy.sh/{topic}", data=message.encode("utf-8"),
                             headers={"Title": "Hearth", "User-Agent": "Hearth"})
            _u.urlopen(req, timeout=12)
            return "Pushed to your phone via ntfy."
        except Exception as e:
            return f"Error pushing via ntfy: {type(e).__name__}: {e}"

    return "Error: channel must be one of: discord, telegram, ntfy."


_HANDLERS["send_to_phone"] = _send_to_phone

TOOL_DEFINITIONS.append({
    "name": "send_to_phone",
    "description": (
        "Push a message FROM this PC TO the user's phone on a connected channel. "
        "Use when the user says 'send this to my Discord / Telegram', 'text me this', "
        "'ping my phone', or wants a result delivered off the PC. Channels: 'discord' "
        "(DMs the owner via the bot), 'telegram' (messages the owner), 'ntfy' (push "
        "notification). The bridges must be set up in Settings - Reach from phone; if a "
        "channel isn't configured this returns a clear error. This is OUTBOUND - the "
        "bridges already handle INCOMING messages and reply on their own."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "channel": {"type": "string", "description": "discord | telegram | ntfy"},
            "message": {"type": "string", "description": "The text to send."},
        },
        "required": ["channel", "message"],
    },
})

_HANDLERS["focus_window"] = _focus_window

TOOL_DEFINITIONS.append({
    "name": "focus_window",
    "description": (
        "Bring an ALREADY-OPEN window to the front / focus it, by a substring of "
        "its title (e.g. 'Spotify', 'chrome', 'notepad', 'Elden Ring'). Use when "
        "the user says 'switch to / bring up / show me / pull up X' and X is "
        "already running — this raises the existing window instead of launching a "
        "new one (that's open_app). Returns the windows it sees if none match."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "name": {"type": "string", "description": "Substring of the target window's title."},
        },
        "required": ["name"],
    },
})

TOOL_DEFINITIONS.append({
    "name": "read_pdf_large",
    "description": (
        "Map-reduce summarize a VERY large PDF (hundreds of pages) that won't "
        "fit any single context. Splits the PDF into page windows, summarizes "
        "each window with a cheap local LLM call, then reduces the per-chunk "
        "summaries into one structured overview (themes, key points, "
        "per-section bullets). Use this — not read_file — when the user asks to "
        "'summarize this whole book/report'. Extracted text + chunk summaries "
        "are cached so a re-ask is instant. mode='background' returns a job_id "
        "immediately and writes the summary to <pdf>_summary.md (overnight use "
        "case); poll get_job_result(job_id) or read the file when done."
    ),
    "parameters": {
        "type": "object",
        "properties": {
            "path": {"type": "string", "description": "PDF path. Absolute or relative to workspace."},
            "mode": {"type": "string", "enum": ["map_reduce", "background"],
                     "description": "'map_reduce' (default, blocks until done) or 'background' (returns a job_id)."},
            "chunk_pages": {"type": "integer", "description": "Pages per window. Default 12."},
            "focus": {"type": "string", "description": "Optional topic to emphasize (e.g. 'the war years', 'financials')."},
        },
        "required": ["path"],
    },
})


def _read_pdf_large(p: Dict) -> str:
    path = (p.get("path") or "").strip()
    if not path:
        return "Error: read_pdf_large needs a 'path'."
    try:
        path = _resolve_read(path)
    except PermissionError as e:
        return f"Error: {e}"
    chunk_pages = _coerce_int(p.get("chunk_pages"), 12) or 12
    focus = (p.get("focus") or "").strip()
    mode = (p.get("mode") or "map_reduce").strip().lower()
    from . import pdf_mapreduce as _pmr
    if mode == "background":
        return json.dumps(_pmr.run_in_background(path, chunk_pages=chunk_pages,
                                                 focus=focus),
                          ensure_ascii=False, default=str)
    return _pmr.run_map_reduce(path, chunk_pages=chunk_pages, focus=focus)


_HANDLERS["read_pdf_large"] = _read_pdf_large

# Auto-load user/agent plugins. Fully guarded — a broken plugin is skipped and
# can NEVER take down the core tools.
try:
    _loaded_plugins = _plugins.load_plugins(WORKSPACE, TOOL_DEFINITIONS, _HANDLERS)
except Exception:
    _loaded_plugins = []


ACTIVITY_LOG = os.path.join(LOGS_DIR, "activity.jsonl")


def _log_activity(event: str, **fields: Any) -> None:
    rec = {"ts": datetime.now().isoformat(timespec="seconds"), "event": event, **fields}
    try:
        with open(ACTIVITY_LOG, "a", encoding="utf-8") as f:
            f.write(json.dumps(rec, ensure_ascii=False, default=str) + "\n")
    except OSError:
        pass


def execute_tool(name: str, args: Optional[Dict] = None) -> str:
    """Run a tool by name. Returns a string (truncated to per-tool cap).
    MCP-bridged tools (name starts with 'mcp_') route through the live
    MCP client session instead of the static _HANDLERS dict."""
    args = args or {}
    if name.startswith("mcp_"):
        try:
            from . import mcp_client
            for td in mcp_client.list_remote_tools():
                if td["name"] == name:
                    r = mcp_client.call_tool(td["_mcp_server"],
                                             td["_mcp_tool"], args)
                    return r.get("output") if r.get("ok") else \
                           f"Error ({name}): {r.get('error', 'unknown')}"
            return f"Error: MCP tool '{name}' not currently connected"
        except Exception as e:
            return f"Error ({name}): {type(e).__name__}: {e}"
    handler = _HANDLERS.get(name)
    if not handler:
        return f"Error: unknown tool '{name}'. Known: {', '.join(sorted(_HANDLERS))}"
    _log_activity("call", tool=name, args=args)
    t0 = datetime.now()
    try:
        result = handler(args)
    except PermissionError as e:
        msg = f"Error ({name}): {e}"
        _log_activity("error", tool=name, error=str(e))
        return msg
    except FileNotFoundError as e:
        msg = f"Error ({name}): file not found: {e}"
        _log_activity("error", tool=name, error=str(e))
        return msg
    except KeyError as e:
        msg = f"Error ({name}): missing required parameter {e}"
        _log_activity("error", tool=name, error=str(e))
        return msg
    except Exception as e:
        msg = f"Error ({name}): {type(e).__name__}: {e}"
        _log_activity("error", tool=name, error=str(e))
        return msg
    if not isinstance(result, str):
        result = json.dumps(result, ensure_ascii=False, default=str)
    truncated = _trunc(result, RESULT_CAPS.get(name, DEFAULT_CAP))
    dt_ms = int((datetime.now() - t0).total_seconds() * 1000)
    _log_activity("result", tool=name, chars=len(truncated), ms=dt_ms)
    return truncated


# ============================================================
# TOOL DIET — defer niche tool SCHEMAS off the default prompt
# ============================================================
# Every handler still EXISTS and still RUNS if called by name — deferring only
# hides a tool's schema from the per-turn `tools` list to cut prompt overhead
# (~130 tokens/tool). The model rediscovers deferred tools via the always-on
# `load_tools` meta-tool; and because execute_tool dispatches purely by name, a
# model that calls a deferred tool directly still works. So this is safe to
# default on — worst case the model uses a niche tool slightly less often, never
# a hard break. Opt out entirely with HEARTH_ALL_TOOLS=1 (loads every schema).
_DEFERRED_TOOLS = {
    # image / video generation (cloud). forge_generate is intentionally NOT
    # deferred: when a local Forge install is detected it surfaces directly so
    # weaker local models can do local image gen without the load_tools hop
    # (they tend to spiral instead of discovering deferred tools). Its
    # secondary controls stay deferred.
    "generate_image", "generate_video", "check_video_task", "list_generations",
    "forge_status", "forge_shutdown",
    # self-extending
    "create_plugin", "list_plugins", "delete_plugin", "create_skill",
    # soul / persona editing
    "edit_soul", "append_soul", "read_soul", "draft_soul",
    # archive
    "list_archive", "extract_archive_file",
    # large-PDF map-reduce (niche; rediscovered via load_tools when the user
    # asks to summarize a whole book)
    "read_pdf_large",
    # extra system info (system_info covers the common case)
    "network_info", "disk_usage", "list_installed_apps", "learn_environment",
    "list_models",
    # reminders niche (set/list/cancel stay core)
    "snooze_reminder", "study_reminder",
    # voice selection (the UI handles voice; rare as a tool call)
    "set_voice", "list_voices",
    # browser niche (browse/click/type/scroll stay core)
    "browse_key", "browse_close",
    # duplicate job controls (start_job/list_jobs/get_job_result stay core)
    "job_kill", "job_list", "job_status", "job_wait",
    # single-purpose utilities / demos
    "color_hex2rgb", "text_encoder_tool", "entity_graph_extractor",
    "website_status_tool", "tic_tac_toe",
    # email (opt-in, needs an app password) — surfaces via load_tools when the
    # user mentions email, so it never clutters the default tool list
    "read_inbox", "send_email",
    "end_session",
}
_unlocked_tools: "set[str]" = set()
_TOOL_DIET = os.environ.get("HEARTH_ALL_TOOLS", "") not in ("1", "true", "yes")

_LOAD_TOOLS_SCHEMA = {
    "type": "function",
    "function": {
        "name": "load_tools",
        "description": (
            "Reveal extra built-in tools that aren't loaded by default (kept off "
            "to save context). Call this with a short query naming what you need, "
            "then call the tool it returns. Groups available on demand: "
            "'image generation' (forge_*), 'plugins' (write your own tools), "
            "'soul' (edit your persona), 'archive' (list/extract zips), "
            "'system' (network/disk/installed apps/battery), 'voice' (pick a "
            "voice), 'end_session'. Pass 'all' to load everything."),
        "parameters": {
            "type": "object",
            "properties": {
                "query": {"type": "string",
                          "description": "What kind of tool you need — a keyword, a group name, or 'all'."},
            },
            "required": ["query"],
        },
    },
}


def _tool_active(name: str) -> bool:
    if not _TOOL_DIET:
        return True
    return (name not in _DEFERRED_TOOLS) or (name in _unlocked_tools)


def _any_tool_still_deferred() -> bool:
    return any(n in _DEFERRED_TOOLS and n not in _unlocked_tools
               for n in (t["name"] for t in TOOL_DEFINITIONS))


def unlock_tools(query: str = "") -> "List[Dict[str, Any]]":
    """Move deferred tools matching `query` (name / category / description
    substring, or 'all') into the active set so their schemas ship next turn.
    Returns the full tool-defs that matched."""
    q = (query or "").strip().lower()
    toks = [w for w in q.replace("-", " ").replace("/", " ").split() if w]
    matched: "List[Dict[str, Any]]" = []
    for td in TOOL_DEFINITIONS:
        n = td["name"]
        if n not in _DEFERRED_TOOLS:
            continue
        cat = _TOOL_CATEGORY.get(n, "").lower()
        hay = f"{n.lower()} {cat} {td['description'].lower()}"
        # 'all'/empty → everything; else match if any query word (or its
        # singular) appears in the tool's name/category/description.
        hit = (not toks) or q == "all" or any(
            (w in hay) or (w.rstrip("s") in hay) for w in toks)
        if hit:
            _unlocked_tools.add(n)
            matched.append(td)
    return matched


def _h_load_tools(p: "Dict[str, Any]") -> str:
    matched = unlock_tools(p.get("query", ""))
    if not matched:
        return ("No extra tools matched that. Groups you can ask for: image "
                "generation, plugins, soul, archive, system, voice, end_session "
                "(or 'all').")
    lines = ["Loaded these tools — you can call them now:"]
    for td in matched:
        props = ", ".join((td["parameters"].get("properties") or {}).keys())
        lines.append(f"- {td['name']}({props}): {td['description']}")
    return "\n".join(lines)


_HANDLERS["load_tools"] = _h_load_tools


# ============================================================
# WORKSPACE-PATH NORMALIZATION
# ============================================================
# Tool descriptions are authored with the default "~/Jarvis/..." workspace.
# But the user can rename the agent (folder becomes ~/<NewName>/) or relocate
# the workspace entirely (D:\Hearth, etc.). If we shipped the literal "~/Jarvis"
# to the model after a rename, we'd be handing it a path that doesn't exist —
# exactly the kind of stale spec that makes a model hallucinate about its own
# environment. Rewrite every description (and nested param description) to the
# REAL workspace path once, at import, so the model always knows where it lives.

def _workspace_display() -> str:
    """The workspace as the model should see it: ~/<name> when under the home
    dir, else the absolute path. Forward slashes (tilde-style, expanduser-safe)."""
    home = os.path.expanduser("~")
    ws = WORKSPACE
    try:
        if os.path.commonpath([os.path.normcase(ws), os.path.normcase(home)]) == os.path.normcase(home):
            return "~/" + os.path.relpath(ws, home).replace(os.sep, "/")
    except Exception:
        pass
    return ws.replace(os.sep, "/")


def _normalize_tool_workspace_paths() -> None:
    disp = _workspace_display()
    if disp == "~/Jarvis":
        return  # default workspace — descriptions already correct, no-op
    def _fix(s):
        return s.replace("~/Jarvis", disp) if isinstance(s, str) else s
    for td in TOOL_DEFINITIONS:
        if isinstance(td.get("description"), str):
            td["description"] = _fix(td["description"])
        props = (td.get("parameters") or {}).get("properties") or {}
        for p in props.values():
            if isinstance(p, dict) and isinstance(p.get("description"), str):
                p["description"] = _fix(p["description"])


_normalize_tool_workspace_paths()


# ============================================================
# PROVIDER FORMAT CONVERTERS
# ============================================================

def to_openai_tools() -> List[Dict[str, Any]]:
    """All tools the model can call this turn: built-ins + any MCP-bridged
    tools currently connected. MCP tools come in lazily (their sessions
    take a few seconds to spawn at boot) so the list grows as servers come
    up — no model restart needed."""
    out = [
        {
            "type": "function",
            "function": {
                "name": td["name"],
                "description": td["description"],
                "parameters": td["parameters"],
            },
        }
        for td in TOOL_DEFINITIONS
        if _tool_active(td["name"])
    ]
    # Discovery meta-tool: only ship it while something is still deferred, so it
    # vanishes once the model has unlocked everything (or HEARTH_ALL_TOOLS=1).
    if _TOOL_DIET and _any_tool_still_deferred():
        out.append(_LOAD_TOOLS_SCHEMA)
    try:
        from . import mcp_client
        for td in mcp_client.list_remote_tools():
            out.append({
                "type": "function",
                "function": {
                    "name": td["name"],
                    "description": td["description"],
                    "parameters": td["parameters"],
                },
            })
    except Exception:
        pass
    return out


# Tool → category, so listings group sensibly instead of dumping a flat wall
# (e.g. the 4 browser tools read as ONE "Web & browser" capability, not 4
# count-padding entries). Order here = display order. Unknown names (user
# plugins) fall into "Custom / plugins".
_CATEGORY_ORDER = [
    "Files & docs", "Web & browser", "System & apps", "Memory",
    "Reminders & alerts", "Email", "Voice", "Self-extending (plugins)",
    "Image generation", "Session",
]
_TOOL_CATEGORY = {
    # Files & docs
    "read_file": "Files & docs", "write_file": "Files & docs", "edit_file": "Files & docs",
    "list_directory": "Files & docs",
    "create_directory": "Files & docs", "delete_path": "Files & docs", "move_path": "Files & docs",
    "find_file": "Files & docs", "grep_search": "Files & docs", "glob_files": "Files & docs",
    "locate_path": "Files & docs", "list_archive": "Files & docs", "extract_archive_file": "Files & docs",
    "read_pdf_large": "Files & docs",
    # Web & browser
    "web_search": "Web & browser", "web_fetch": "Web & browser",
    "open_url": "Web & browser", "open_in_browser": "Web & browser", "list_browsers": "Web & browser",
    "browse": "Web & browser", "browse_click": "Web & browser", "browse_type": "Web & browser", "browse_scroll": "Web & browser", "browse_key": "Web & browser",
    "browse_close": "Web & browser",
    # System & apps
    "run_command": "System & apps", "system_info": "System & apps", "list_processes": "System & apps",
    "network_info": "System & apps", "get_battery": "System & apps", "list_installed_apps": "System & apps",
    "disk_usage": "System & apps", "open_app": "System & apps", "screenshot": "System & apps",
    "focus_window": "System & apps",
    "list_jobs": "Background jobs", "get_job_result": "Background jobs",
    "view_image": "System & apps", "clipboard_read": "System & apps", "clipboard_write": "System & apps",
    "get_time": "System & apps", "whoami": "System & apps", "list_models": "System & apps",
    "learn_environment": "System & apps",
    # Memory
    "memory_save": "Memory", "memory_recall": "Memory", "memory_list": "Memory",
    "memory_forget": "Memory", "search_chats": "Memory",
    "edit_soul": "Memory", "append_soul": "Memory", "read_soul": "Memory",
    # Reminders & alerts
    "set_reminder": "Reminders & alerts", "list_reminders": "Reminders & alerts",
    "cancel_reminder": "Reminders & alerts", "snooze_reminder": "Reminders & alerts",
    "spawn_subagent": "Sub-agents", "list_subagent_personas": "Sub-agents",
    "get_subagent_result": "Sub-agents",
    "notify": "Reminders & alerts",
    # Voice
    "set_voice": "Voice", "list_voices": "Voice",
    # Self-extending
    "create_plugin": "Self-extending (plugins)", "list_plugins": "Self-extending (plugins)",
    "delete_plugin": "Self-extending (plugins)",
    "install_skill": "Self-extending (plugins)",
    "read_inbox": "Email", "send_email": "Email",
    # Image generation
    "generate_image": "Image generation", "generate_video": "Image generation",
    "check_video_task": "Image generation", "list_generations": "Image generation",
    # Session
    "end_session": "Session",
}


def tools_by_category() -> "List[Tuple[str, List[Dict[str, Any]]]]":
    """Group the live TOOL_DEFINITIONS by category, in display order. Returns
    [(category, [tooldef, ...]), ...]. Unmapped tools (user plugins) → 'Custom / plugins'."""
    groups: Dict[str, List[Dict[str, Any]]] = {}
    for td in TOOL_DEFINITIONS:
        cat = _TOOL_CATEGORY.get(td["name"], "Custom / plugins")
        groups.setdefault(cat, []).append(td)
    order = _CATEGORY_ORDER + [c for c in groups if c not in _CATEGORY_ORDER]
    return [(c, groups[c]) for c in order if c in groups]


def to_claude_tools() -> List[Dict[str, Any]]:
    return [
        {"name": td["name"], "description": td["description"], "input_schema": td["parameters"]}
        for td in TOOL_DEFINITIONS
    ]


def to_gemini_tools():
    """Lazy: only build if google.genai is installed."""
    try:
        from google.genai import types  # type: ignore
    except ImportError:
        raise RuntimeError("Install google-genai to use Gemini: pip install google-genai")

    tmap = {"string": "STRING", "integer": "INTEGER", "number": "NUMBER",
            "boolean": "BOOLEAN", "array": "ARRAY", "object": "OBJECT"}

    def schema(s: Dict) -> Any:
        kw: Dict[str, Any] = {"type": tmap.get(s.get("type", "string"), "STRING")}
        if "description" in s:
            kw["description"] = s["description"]
        if "properties" in s:
            kw["properties"] = {k: schema(v) for k, v in s["properties"].items()}
        if "required" in s and s["required"]:
            kw["required"] = s["required"]
        if "items" in s:
            kw["items"] = schema(s["items"])
        return types.Schema(**kw)

    decls = []
    for td in TOOL_DEFINITIONS:
        kw = {"name": td["name"], "description": td["description"]}
        if td["parameters"].get("properties"):
            kw["parameters"] = schema(td["parameters"])
        decls.append(types.FunctionDeclaration(**kw))
    return [types.Tool(function_declarations=decls)]


# ============================================================
# CONTEXT WINDOW MANAGEMENT
# ----------------------------------------------------------------
# Two strategies, used together:
#
# 1) trim_to_budget()  — Void's pattern. Greedy weight-based truncation.
#    Protects system message + most-recent turns. Aggressively trims old
#    assistant turns BEFORE touching user turns. Char-level cuts, no
#    deletions, so tool_call links remain valid.
#
# 2) compact_history() — Graphify's pattern. When the conversation is
#    genuinely too long for trimming alone, summarize the OLDEST chunk into
#    a single synthetic system message and drop the originals. Caller
#    supplies a summarizer callback (the LLM client) so this stays
#    provider-agnostic.
#
# CHARS_PER_TOKEN of 4 is conservative — better to overestimate token cost
# than to overflow the model.
# ============================================================

CHARS_PER_TOKEN = 4


def _msg_chars(m: Dict[str, Any]) -> int:
    content = m.get("content")
    n = len(content) if isinstance(content, str) else 0
    for tc in m.get("tool_calls") or []:
        try:
            n += len(tc["function"]["name"]) + len(tc["function"].get("arguments") or "")
        except (KeyError, TypeError):
            pass
    return n


def _msg_weight(m: Dict[str, Any], idx: int, total: int) -> float:
    role = m.get("role")
    if role == "system":
        return 0.01  # never trim
    # Protect first 2 and last 3 turns (recency bias)
    if idx < 2 or idx >= total - 3:
        return 0.05
    if role == "tool":
        return 8.0   # tool outputs are huge and stale fastest
    if role == "assistant":
        return 10.0  # narration; prefer to trim
    return 1.0       # user — trim last


def trim_to_budget(messages: List[Dict[str, Any]], context_window: int,
                   reserved_output: int = 0) -> List[Dict[str, Any]]:
    """Trim message contents (not whole messages) until total fits.

    context_window is in TOKENS. We convert to chars via CHARS_PER_TOKEN.
    Returns a NEW list — does not mutate input.
    """
    if reserved_output <= 0:
        reserved_output = max(context_window // 4, 1024)
    budget_chars = max(0, (context_window - reserved_output) * CHARS_PER_TOKEN)

    msgs = [dict(m) for m in messages]
    total = sum(_msg_chars(m) for m in msgs)
    if total <= budget_chars:
        return msgs

    # Greedy phase 1: while over budget, find the heaviest weight*size message
    # and cut its content to half (min 120 chars). Repeat.
    while total > budget_chars:
        best_idx = -1
        best_score = 0.0
        for i, m in enumerate(msgs):
            chars = _msg_chars(m)
            if chars <= 120:
                continue
            score = _msg_weight(m, i, len(msgs)) * chars
            if score > best_score:
                best_score = score
                best_idx = i
        if best_idx < 0:
            break  # nothing left to trim by content
        m = msgs[best_idx]
        if isinstance(m.get("content"), str):
            cur = m["content"]
            cut = max(120, len(cur) // 2)
            if cut < len(cur):
                m["content"] = cur[:cut] + "\n…[trimmed]"
        new_total = sum(_msg_chars(x) for x in msgs)
        if new_total >= total:
            break  # no progress
        total = new_total

    # Final safety net: NO MATTER WHAT the caller hands us, the output must
    # contain at least one user role. LM Studio's chat templates (especially
    # the Harmonic / Hermes finetunes) crash hard with "No user query found
    # in messages" when this invariant is violated, and the trace doesn't
    # tell you which trimming path lost it. We re-check at the end.
    if total > budget_chars and len(msgs) > 2:
        # index of the last user message
        last_user = -1
        for i in range(len(msgs) - 1, -1, -1):
            if msgs[i].get("role") == "user":
                last_user = i
                break
        # Protected tail = everything from last_user onward. Protected head =
        # system message (index 0 if present).
        sys_present = msgs and msgs[0].get("role") == "system"
        head_keep = 1 if sys_present else 0
        # Drop messages between head and last_user, oldest first, until we fit
        # or there's nothing droppable left.
        drop_lo = head_keep
        while total > budget_chars and last_user > drop_lo:
            # Removing index drop_lo; everything shifts down, last_user too.
            removed = msgs.pop(drop_lo)
            last_user -= 1
            total -= _msg_chars(removed)
        # A dropped tool message can orphan a tool_calls assistant turn or
        # vice-versa; clean leading orphan tool messages after the head.
        while len(msgs) > head_keep and msgs[head_keep].get("role") == "tool":
            total -= _msg_chars(msgs[head_keep])
            msgs.pop(head_keep)

    # Final invariant: SOME user turn must exist in the output. If the
    # incoming history had a user message and we somehow dropped it, restore
    # the most-recent one from the original. If the incoming had NO user
    # at all (rare: caller bug), synthesize a continue-marker so the chat
    # template can render rather than crash.
    if not any(m.get("role") == "user" for m in msgs):
        original_user = next(
            (m for m in reversed(messages) if m.get("role") == "user"),
            None,
        )
        msgs.append(original_user or {
            "role": "user",
            "content": "Continue using the results above.",
        })

    return msgs


def _truncate_kept_tool_results(msgs: List[Dict[str, Any]],
                                 max_chars: int = 600) -> List[Dict[str, Any]]:
    """Replace long `tool` role payloads in the kept tail with a short
    truncation marker. The model just needs to know the tool ran and got
    a result; the FULL text from 8 turns ago doesn't help and eats budget.
    Browse results in particular are massive (~3500 chars each)."""
    out: List[Dict[str, Any]] = []
    for m in msgs:
        if m.get("role") == "tool":
            content = m.get("content") or ""
            if isinstance(content, str) and len(content) > max_chars:
                trunc = (f"{content[:max_chars].rstrip()} "
                         f"…[truncated, full result was {len(content)} chars]")
                m = {**m, "content": trunc}
        out.append(m)
    return out


# Marker on a compaction summary message. Doubles as the recovery key: the NEXT
# compaction recognizes a prior summary by this prefix and folds it in instead of
# re-summarizing it (which drifts into summary-of-a-summary). Single source of
# truth — if this text drifts, iterative recovery silently breaks.
_SUMMARY_PREFIX = "Earlier-conversation summary (compacted to save context):"


def dedup_tool_results(messages: List[Dict[str, Any]],
                       min_chars: int = 200) -> List[Dict[str, Any]]:
    """Cheap pre-pass: replace OLDER duplicate tool outputs with a back-reference
    so an identical big result (e.g. the same browse/read_file dump repeated
    across turns) isn't stored — and later summarized — N times. The most-recent
    copy is kept verbatim. This only rewrites the `content` of `tool` messages;
    it never drops a message or touches tool_call_id, so tool_call/tool_result
    pairing stays valid. Returns a NEW list (never mutates the input). Runs
    without an LLM call, so it can shrink context before compaction even fires.
    """
    import hashlib as _hl
    out = [dict(m) for m in messages]
    seen: set = set()
    for i in range(len(out) - 1, -1, -1):   # newest first -> the newest copy wins
        m = out[i]
        if m.get("role") != "tool":
            continue
        c = m.get("content")
        if not isinstance(c, str) or len(c) < min_chars:
            continue
        h = _hl.md5(c.encode("utf-8", "replace")).hexdigest()[:12]
        if h in seen:
            out[i] = {**m, "content":
                      "[Duplicate tool output — identical to a more recent call; omitted to save context.]"}
        else:
            seen.add(h)
    return out


def compact_history(messages: List[Dict[str, Any]],
                    summarize: Any,
                    keep_recent: int = 8,
                    target_chars: int = 0) -> List[Dict[str, Any]]:
    """Replace older turns with a single summary system message.

    `summarize(text) -> str` is a caller-supplied callback (sync) — usually
    a small LLM call. We hand it the concatenated old turns and expect a
    short digest back. Recent turns are preserved verbatim.

    `target_chars` (optional): if set, after compaction the result is
    re-measured; if it's STILL above target_chars, the kept tail's tool
    results are truncated and the keep_recent window is halved + re-tried.
    This catches the case where the last 8 messages include 4 browse
    results at 3500 chars each — compaction "succeeded" but kept 14K of
    bulky tool output and the next chat call still wedges.

    Returns: [system, summary_msg, *recent_turns]. If there's nothing to
    compact, returns the input unchanged.
    """
    if len(messages) <= keep_recent + 2:
        return list(messages)

    # Cheap pre-pass: collapse repeated identical tool dumps before we spend an
    # LLM call summarizing them. Often the single biggest token sink.
    messages = dedup_tool_results(messages)

    # split: leading system + everything-up-to-keep_recent + recent tail
    head: List[Dict[str, Any]] = []
    if messages and messages[0].get("role") == "system":
        head = [messages[0]]
        body = messages[1:]
    else:
        body = list(messages)

    if len(body) <= keep_recent:
        return list(messages)

    # Walk the tail back until it contains at least one `user` message.
    # LM Studio's Jinja templates require a user role somewhere — without
    # this, compaction in the middle of a long tool chain would produce
    # [system, summary_system, assistant, tool, assistant, tool, ...] with
    # zero user messages, and the next API call would crash with
    # "No user query found in messages.".
    expand_to = keep_recent
    while expand_to < len(body):
        if any(m.get("role") == "user" for m in body[-expand_to:]):
            break
        expand_to += 1
    # If even the entire body has no user message, bail out: don't compact.
    if not any(m.get("role") == "user" for m in body[-expand_to:]):
        return list(messages)

    to_compact = body[:-expand_to]
    recent = body[-expand_to:]
    # Also: never start the kept tail with a `tool` message — that role
    # only makes sense as a reply to an immediately-prior tool_calls turn.
    # Drop leading orphan tool messages.
    while recent and recent[0].get("role") == "tool":
        to_compact.append(recent.pop(0))

    # Iterative summary: if a PRIOR compaction summary is among the turns we're
    # about to fold away, pull it out and feed it to the summarizer as
    # established context rather than re-summarizing a summary (which drifts).
    prev_summary = ""
    _kept_to_compact: List[Dict[str, Any]] = []
    for m in to_compact:
        c = m.get("content")
        if (m.get("role") == "system" and isinstance(c, str)
                and c.startswith(_SUMMARY_PREFIX)):
            prev_summary = c[len(_SUMMARY_PREFIX):].lstrip(" :\n")
        else:
            _kept_to_compact.append(m)
    to_compact = _kept_to_compact

    transcript_parts = []
    for m in to_compact:
        role = m.get("role", "?")
        content = m.get("content") or ""
        if isinstance(content, str) and content:
            transcript_parts.append(f"[{role}] {content}")
        for tc in m.get("tool_calls") or []:
            try:
                transcript_parts.append(
                    f"[{role}→tool] {tc['function']['name']}({tc['function'].get('arguments','')})"
                )
            except (KeyError, TypeError):
                pass

    transcript = "\n".join(transcript_parts)
    if prev_summary:
        transcript = ("PREVIOUS SUMMARY:\n" + prev_summary
                      + "\n\n---\nNEWER TURNS TO FOLD IN:\n" + transcript)
    try:
        summary = summarize(transcript)
    except Exception as e:
        summary = f"[summary failed: {e}]"

    summary_msg = {
        "role": "system",
        "content": _SUMMARY_PREFIX + "\n" + (summary or "[empty]"),
    }
    result = head + [summary_msg] + recent

    # Post-compact tightening: if a target was given AND we're still over,
    # the kept tail is the problem (long browse / read_file results in
    # recent turns). Truncate kept tool payloads first; if still over,
    # halve the kept window and re-run (recursive but bounded by len).
    if target_chars and target_chars > 0:
        current = sum(_msg_chars(m) for m in result)
        if current > target_chars:
            result = head + [summary_msg] + _truncate_kept_tool_results(recent)
            current = sum(_msg_chars(m) for m in result)
        if current > target_chars and keep_recent > 2:
            # Halve and retry — guaranteed to converge since each pass
            # either fits or shrinks the kept window.
            new_keep = max(2, keep_recent // 2)
            return compact_history(messages, summarize,
                                   keep_recent=new_keep,
                                   target_chars=target_chars)
    return result


def estimate_tokens(messages: List[Dict[str, Any]]) -> int:
    """Cheap token estimate for the whole message list."""
    return sum(_msg_chars(m) for m in messages) // CHARS_PER_TOKEN
